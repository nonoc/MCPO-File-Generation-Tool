import re
import os
import ast
import csv
import json
import uuid
import emoji
import math
import time
import base64
import shutil
import datetime
import tarfile
import zipfile
import py7zr
import logging
import requests
from requests import get, post
from requests.auth import HTTPBasicAuth
import threading
import markdown2
import tempfile
from pathlib import Path
from lxml import etree
from PIL import Image
from docx import Document
from docx.shared import Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.shared import qn
from docx.oxml.ns import nsdecls
from docx.oxml import parse_xml
from docx.shared import Pt as DocxPt
from bs4 import BeautifulSoup, NavigableString
from mcp.server.fastmcp import FastMCP, Context
from mcp.server.session import ServerSession
from openpyxl import Workbook, load_workbook
from openpyxl.comments import Comment
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt as PptPt
from pptx.enum.shapes import PP_PLACEHOLDER 
from pptx.parts.image import Image
from pptx.enum.text import MSO_AUTO_SIZE
from io import BytesIO
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem, Image as ReportLabImage
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.enums import TA_LEFT
from reportlab.lib.units import mm

SCRIPT_VERSION = "0.8.1"

URL = os.getenv('OWUI_URL')
TOKEN = os.getenv('JWT_SECRET') ## will be deleted in 1.0.0

PERSISTENT_FILES = os.getenv("PERSISTENT_FILES", "false")
FILES_DELAY = int(os.getenv("FILES_DELAY", 60)) 

EXPORT_DIR_ENV = os.getenv("FILE_EXPORT_DIR")
EXPORT_DIR = (EXPORT_DIR_ENV or r"/output").rstrip("/")
os.makedirs(EXPORT_DIR, exist_ok=True)

BASE_URL_ENV = os.getenv("FILE_EXPORT_BASE_URL")
BASE_URL = (BASE_URL_ENV or "http://localhost:9003/files").rstrip("/")

LOG_LEVEL_ENV = os.getenv("LOG_LEVEL")
LOG_FORMAT_ENV = os.getenv(
    "LOG_FORMAT", "%(asctime)s %(levelname)s %(name)s - %(message)s"
)

DOCS_TEMPLATE_PATH = os.getenv("DOCS_TEMPLATE_DIR", "/rootPath/templates")
PPTX_TEMPLATE = None
DOCX_TEMPLATE = None
XLSX_TEMPLATE = None
PPTX_TEMPLATE_PATH = None
DOCX_TEMPLATE_PATH = None
XLSX_TEMPLATE_PATH = None

if DOCS_TEMPLATE_PATH and os.path.exists(DOCS_TEMPLATE_PATH):
    logging.debug(f"Template Folder: {DOCS_TEMPLATE_PATH}")
    for root, dirs, files in os.walk(DOCS_TEMPLATE_PATH):
        for file in files:
            fpath = os.path.join(root, file)
            if file.lower().endswith(".pptx") and PPTX_TEMPLATE_PATH is None:
                PPTX_TEMPLATE_PATH = fpath
                logging.debug(f"PPTX template: {PPTX_TEMPLATE_PATH}")
            elif file.lower().endswith(".docx") and DOCX_TEMPLATE_PATH is None:
                DOCX_TEMPLATE_PATH = fpath
            elif file.lower().endswith(".xlsx") and XLSX_TEMPLATE_PATH is None:
                XLSX_TEMPLATE_PATH = fpath
    if PPTX_TEMPLATE_PATH:
        try:
            PPTX_TEMPLATE = Presentation(PPTX_TEMPLATE_PATH)
            logging.debug(f"Using PPTX template: {PPTX_TEMPLATE_PATH}")
        except Exception as e:
            logging.warning(f"PPTX template failed to load : {e}")
            PPTX_TEMPLATE = None
    else:
        logging.debug("No PPTX template found. Creation of a blank document.")
        PPTX_TEMPLATE = None

    if DOCX_TEMPLATE_PATH and os.path.exists(DOCS_TEMPLATE_PATH):
        try:
            DOCX_TEMPLATE = Document(DOCX_TEMPLATE_PATH)
            logging.debug(f"Using DOCX template: {DOCX_TEMPLATE_PATH}")
        except Exception as e:
            logging.warning(f"DOCX template failed to load : {e}")
            DOCX_TEMPLATE = None
    else:
        logging.debug("No DOCX template found. Creation of a blank document.")
        DOCX_TEMPLATE = None
    
    XLSX_TEMPLATE_PATH = os.path.join("/rootPath/templates","Default_Template.xlsx")

    if XLSX_TEMPLATE_PATH:
        try:
            XLSX_TEMPLATE = load_workbook(XLSX_TEMPLATE_PATH)
            logging.debug(f"Using XLSX template: {XLSX_TEMPLATE_PATH}")
        except Exception as e:
            logging.warning(f"Failed to load XLSX template: {e}")
            XLSX_TEMPLATE = None
    else:
        logging.debug("No XLSX template found. Creation of a blank document.")
        XLSX_TEMPLATE = None


def search_image(query):
    log.debug(f"Searching for image with query: '{query}'")
    image_source = os.getenv("IMAGE_SOURCE", "unsplash")

    if image_source == "unsplash":
        return search_unsplash(query)
    elif image_source == "local_sd":
        return search_local_sd(query)
    elif image_source == "pexels":
        return search_pexels(query)
    else:
        log.warning(f"Image source unknown : {image_source}")
        return None

def search_local_sd(query: str):
    log.debug(f"Searching for local SD image with query: '{query}'")
    SD_URL = os.getenv("LOCAL_SD_URL")
    SD_USERNAME = os.getenv("LOCAL_SD_USERNAME")
    SD_PASSWORD = os.getenv("LOCAL_SD_PASSWORD")
    DEFAULT_MODEL = os.getenv("LOCAL_SD_DEFAULT_MODEL", "sd_xl_base_1.0.safetensors")
    DEFAULT_STEPS = int(os.getenv("LOCAL_SD_STEPS", 20))
    DEFAULT_WIDTH = int(os.getenv("LOCAL_SD_WIDTH", 512))
    DEFAULT_HEIGHT = int(os.getenv("LOCAL_SD_HEIGHT", 512))
    DEFAULT_CFG_SCALE = float(os.getenv("LOCAL_SD_CFG_SCALE", 1.5))
    DEFAULT_SCHEDULER = os.getenv("LOCAL_SD_SCHEDULER", "Karras")
    DEFAULT_SAMPLE = os.getenv("LOCAL_SD_SAMPLE", "Euler a")

    if not SD_URL:
        log.warning("LOCAL_SD_URL is not defined.")
        return None

    payload = {
        "prompt": query.strip(),
        "steps": DEFAULT_STEPS,
        "width": DEFAULT_WIDTH,
        "height": DEFAULT_HEIGHT,
        "cfg_scale": DEFAULT_CFG_SCALE,
        "sampler_name": DEFAULT_SAMPLE,
        "scheduler": DEFAULT_SCHEDULER,
        "enable_hr": False,
        "hr_upscaler": "Latent",
        "seed": -1,
        "override_settings": {
            "sd_model_checkpoint": DEFAULT_MODEL
        }
    }

    try:
        url = f"{SD_URL}/sdapi/v1/txt2img"
        log.debug(f"Sending request to local SD API at {url}")
        response = requests.post(
            url,
            json=payload,
            headers={"Content-Type": "application/json"},
            auth=HTTPBasicAuth(SD_USERNAME, SD_PASSWORD),
            timeout=30
        )
        response.raise_for_status()
        data = response.json()

        images = data.get("images", [])
        if not images:
            log.warning(f"No image generated for the request : '{query}'")
            return None

        image_b64 = images[0]
        image_data = base64.b64decode(image_b64)

        folder_path = _generate_unique_folder()
        filename = f"{query.replace(' ', '_')}.png"
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)

        with open(filepath, "wb") as f:
            f.write(image_data)

        return _public_url(folder_path, filename)

    except requests.exceptions.Timeout:
        log.error(f"Timeout during generation for : '{query}'")
    except requests.exceptions.RequestException as e:
        log.error(f"Network error : {e}")
    except Exception as e:
        log.error(f"Unexpected error : {e}")

    return None

def search_unsplash(query):
    log.debug(f"Searching Unsplash for query: '{query}'")
    api_key = os.getenv("UNSPLASH_ACCESS_KEY")
    if not api_key:
        log.warning("UNSPLASH_ACCESS_KEY is not set. Cannot search for images.")
        return None
    url = "https://api.unsplash.com/search/photos"
    params = {
        "query": query,
        "per_page": 1,
        "orientation": "landscape"
    }
    headers = {"Authorization": f"Client-ID {api_key}"}
    log.debug(f"Sending request to Unsplash API")
    try:
        response = requests.get(url, params=params, headers=headers)
        log.debug(f"Unsplash API response status: {response.status_code}")
        response.raise_for_status() 
        data = response.json()
        if data.get("results"):
            image_url = data["results"][0]["urls"]["regular"]
            log.debug(f"Found image URL for '{query}': {image_url}")
            return image_url
        else:
            log.debug(f"No results found on Unsplash for query: '{query}'")
    except requests.exceptions.RequestException as e:
        log.error(f"Network error while searching image for '{query}': {e}")
    except json.JSONDecodeError as e:
        log.error(f"Error decoding JSON from Unsplash for '{query}': {e}")
    except Exception as e:
        log.error(f"Unexpected error searching image for '{query}': {e}")
    return None 

def search_pexels(query):
    log.debug(f"Searching Pexels for query: '{query}'")
    api_key = os.getenv("PEXELS_ACCESS_KEY")
    if not api_key:
        log.warning("PEXELS_ACCESS_KEY is not set. Cannot search for images.")
        return None
    url = "https://api.pexels.com/v1/search"
    params = {
        "query": query,
        "per_page": 1,
        "orientation": "landscape"
    }
    headers = {"Authorization": f"{api_key}"}
    log.debug(f"Sending request to Pexels API")
    try:
        response = requests.get(url, params=params, headers=headers)
        log.debug(f"Pexels API response status: {response.status_code}")
        response.raise_for_status() 
        data = response.json()
        if data.get("photos"):
            image_url = data["photos"][0]["src"]["large"]
            log.debug(f"Found image URL for '{query}': {image_url}")
            return image_url
        else:
            log.debug(f"No results found on Pexels for query: '{query}'")
    except requests.exceptions.RequestException as e:
        log.error(f"Network error while searching image for '{query}': {e}")
    except json.JSONDecodeError as e:
        log.error(f"Error decoding JSON from Pexels for '{query}': {e}")
    except Exception as e:
        log.error(f"Unexpected error searching image for '{query}': {e}")
    return None

def _resolve_log_level(val: str | None) -> int:
    if not val:
        return logging.INFO
    v = val.strip()
    if v.isdigit():
        try:
            return int(v)
        except ValueError:
            return logging.INFO
    return getattr(logging, v.upper(), logging.INFO)

logging.basicConfig(
    level=_resolve_log_level(LOG_LEVEL_ENV),
    format=LOG_FORMAT_ENV,
)
log = logging.getLogger("file_export_mcp")
log.setLevel(_resolve_log_level(LOG_LEVEL_ENV))
log.info("Effective LOG_LEVEL -> %s", logging.getLevelName(log.level))

mcp = FastMCP("file_export")

def dynamic_font_size(content_list, max_chars=400, base_size=28, min_size=12):
    total_chars = sum(len(line) for line in content_list)
    ratio = total_chars / max_chars if max_chars > 0 else 1
    if ratio <= 1:
        return PptPt(base_size)
    else:
        new_size = int(base_size / ratio)
        return PptPt(max(min_size, new_size))

def _public_url(folder_path: str, filename: str) -> str:
    """Build a stable public URL for a generated file."""
    folder = os.path.basename(folder_path).lstrip("/")
    name = filename.lstrip("/")
    return f"{BASE_URL}/{folder}/{name}"

def _generate_unique_folder() -> str:
    folder_name = f"export_{uuid.uuid4().hex[:10]}_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}"
    folder_path = os.path.join(EXPORT_DIR, folder_name)
    os.makedirs(folder_path, exist_ok=True)
    return folder_path

def _generate_filename(folder_path: str, ext: str, filename: str = None) -> tuple[str, str]:
    if not filename:
        filename = f"export_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.{ext}"
    base, ext = os.path.splitext(filename)
    filepath = os.path.join(folder_path, filename)
    counter = 1
    while os.path.exists(filepath):
        filename = f"{base}_{counter}{ext}"
        filepath = os.path.join(folder_path, filename)
        counter += 1
    return filepath, filename

styles = getSampleStyleSheet()
styles.add(ParagraphStyle(
    name="CustomHeading1",
    parent=styles["Heading1"],
    textColor=colors.HexColor("#0A1F44"),
    fontSize=18,
    spaceAfter=16,
    spaceBefore=12,
    alignment=TA_LEFT
))
styles.add(ParagraphStyle(
    name="CustomHeading2",
    parent=styles["Heading2"],
    textColor=colors.HexColor("#1C3F77"),
    fontSize=14,
    spaceAfter=12,
    spaceBefore=10,
    alignment=TA_LEFT
))
styles.add(ParagraphStyle(
    name="CustomHeading3",
    parent=styles["Heading3"],
    textColor=colors.HexColor("#3A6FB0"), 
    fontSize=12,
    spaceAfter=10,
    spaceBefore=8,
    alignment=TA_LEFT
))
styles.add(ParagraphStyle(
    name="CustomNormal",
    parent=styles["Normal"],
    fontSize=11,
    leading=14,
    alignment=TA_LEFT
))
styles.add(ParagraphStyle(
    name="CustomListItem",
    parent=styles["Normal"],
    fontSize=11,
    leading=14,
    alignment=TA_LEFT
))
styles.add(ParagraphStyle(
    name="CustomCode",
    parent=styles["Code"],
    fontSize=10,
    leading=12,
    fontName="Courier",
    backColor=colors.HexColor("#F5F5F5"),
    borderColor=colors.HexColor("#CCCCCC"),
    borderWidth=1,
    leftIndent=10,
    rightIndent=10,
    topPadding=5,
    bottomPadding=5
))

def render_text_with_emojis(text: str) -> str:
    if not text:
        return ""
    try:
        converted = emoji.emojize(text, language="alias")
        return converted
    except Exception as e:
        log.error(f"Error in emoji conversion: {e}")
        return text

def process_list_items(ul_or_ol_element, is_ordered=False):
    items = []
    bullet_type = '1' if is_ordered else 'bullet'
    for li in ul_or_ol_element.find_all('li', recursive=False):
        li_text_parts = []
        for content in li.contents:
            if isinstance(content, NavigableString):
                li_text_parts.append(str(content))
            elif content.name not in ['ul', 'ol']:
                 li_text_parts.append(content.get_text())
        li_text = ''.join(li_text_parts).strip()
        list_item_paragraph = None
        if li_text:
            rendered_text = render_text_with_emojis(li_text)
            list_item_paragraph = Paragraph(rendered_text, styles["CustomListItem"])
        sub_lists = li.find_all(['ul', 'ol'], recursive=False)
        sub_flowables = []
        if list_item_paragraph:
             sub_flowables.append(list_item_paragraph)
        for sub_list in sub_lists:
            is_sub_ordered = sub_list.name == 'ol'
            nested_items = process_list_items(sub_list, is_sub_ordered)
            if nested_items:
                nested_list_flowable = ListFlowable(
                    nested_items,
                    bulletType='1' if is_sub_ordered else 'bullet',
                    leftIndent=10 * mm,
                    bulletIndent=5 * mm,
                    spaceBefore=2,
                    spaceAfter=2
                )
                sub_flowables.append(nested_list_flowable)
        if sub_flowables:
            items.append(ListItem(sub_flowables))
    return items

def render_html_elements(soup):
    log.debug("Starting render_html_elements...")
    story = []
    element_count = 0
    for elem in soup.children:
        element_count += 1
        log.debug(f"Processing element #{element_count}: {type(elem)}, name={getattr(elem, 'name', 'NavigableString')}")
        if isinstance(elem, NavigableString):
            text = str(elem).strip()
            if text:
                log.debug(f"Adding Paragraph from NavigableString: {text[:50]}...")
                story.append(Paragraph(render_text_with_emojis(text), styles["CustomNormal"]))
                story.append(Spacer(1, 6))
        elif hasattr(elem, 'name'):
            tag_name = elem.name
            log.debug(f"Handling tag: <{tag_name}>")
            if tag_name == "h1":
                text = render_text_with_emojis(elem.get_text().strip())
                log.debug(f"Adding H1: {text[:50]}...")
                story.append(Paragraph(text, styles["CustomHeading1"]))
                story.append(Spacer(1, 10))
            elif tag_name == "h2":
                text = render_text_with_emojis(elem.get_text().strip())
                log.debug(f"Adding H2: {text[:50]}...")
                story.append(Paragraph(text, styles["CustomHeading2"]))
                story.append(Spacer(1, 8))
            elif tag_name == "h3":
                text = render_text_with_emojis(elem.get_text().strip())
                log.debug(f"Adding H3: {text[:50]}...")
                story.append(Paragraph(text, styles["CustomHeading3"]))
                story.append(Spacer(1, 6))
            elif tag_name == "p":
                imgs = elem.find_all("img")
                if imgs:
                    for img_tag in imgs:
                        src = img_tag.get("src")
                        alt = img_tag.get("alt", "[Image]")
                        try:
                            if src and src.startswith("http"):
                                log.debug(f"Downloading image from URL: {src}")
                                response = requests.get(src)
                                response.raise_for_status()
                                img_data = BytesIO(response.content)
                                img = ReportLabImage(img_data, width=200, height=150)  # ✅ CORRIGÉ
                            else:
                                log.debug(f"Loading local image: {src}")
                                img = ReportLabImage(src, width=200, height=150)  # ✅ CORRIGÉ
                            story.append(img)
                            story.append(Spacer(1, 10))
                        except Exception as e:
                            log.error(f"Error loading image {src}: {e}")
                            story.append(Paragraph(f"[Image: {alt}]", styles["CustomNormal"]))
                            story.append(Spacer(1, 6))
                else:
                    text = render_text_with_emojis(elem.get_text().strip())
                    if text:
                        log.debug(f"Adding Paragraph: {text[:50]}...")
                        story.append(Paragraph(text, styles["CustomNormal"]))
                        story.append(Spacer(1, 6))
            elif tag_name in ["ul", "ol"]:
                is_ordered = tag_name == "ol"
                log.debug(f"Processing list (ordered={is_ordered})...")
                items = process_list_items(elem, is_ordered)
                if items:
                    log.debug(f"Adding ListFlowable with {len(items)} items")
                    story.append(ListFlowable(items,
                        bulletType='1' if is_ordered else 'bullet',
                        leftIndent=10 * mm,
                        bulletIndent=5 * mm,
                        spaceBefore=6,
                        spaceAfter=10
                    ))
            elif tag_name == "blockquote":
                text = render_text_with_emojis(elem.get_text().strip())
                if text:
                    log.debug(f"Adding Blockquote: {text[:50]}...")
                    story.append(Paragraph(f"{text}", styles["CustomNormal"]))
                    story.append(Spacer(1, 8))
            elif tag_name in ["code", "pre"]:
                text = elem.get_text().strip()
                if text:
                    log.debug(f"Adding Code/Pre block: {text[:50]}...")
                    story.append(Paragraph(text, styles["CustomCode"]))
                    story.append(Spacer(1, 6 if tag_name == "code" else 8))
            elif tag_name == "img":
                src = elem.get("src")
                alt = elem.get("alt", "[Image]")
                log.debug(f"Found <img> tag. src='{src}', alt='{alt}'")
                if src is not None: 
                    try:
                        if src.startswith("image_query:"):

                            query = src.replace("image_query:", "").strip()
                            log.debug(f"Handling image_query: '{query}'")
                            image_url = search_image(query)
                            if image_url:
                                log.debug(f"Downloading image from Unsplash URL: {image_url}")
                                response = requests.get(image_url)
                                log.debug(f"Image download response status: {response.status_code}")
                                response.raise_for_status()
                                img_data = BytesIO(response.content)
                                img = ReportLabImage(img_data, width=200, height=150)
                                log.debug("Adding ReportLab Image object to story (Unsplash)")
                                story.append(img)
                                story.append(Spacer(1, 10))
                            else:
                                log.warning(f"No image found for query: {query}")
                                story.append(Paragraph(f"[Image non trouvee pour: {query}]", styles["CustomNormal"]))
                                story.append(Spacer(1, 6))
                        elif src.startswith("http"):
                            log.debug(f"Downloading image from direct URL: {src}")
                            response = requests.get(src)
                            log.debug(f"Image download response status: {response.status_code}")
                            response.raise_for_status()
                            img_data = BytesIO(response.content)
                            img = ReportLabImage(img_data, width=200, height=150)
                            log.debug("Adding ReportLab Image object to story (Direct URL)")
                            story.append(img)
                            story.append(Spacer(1, 10))
                        else:
                            log.debug(f"Loading local image: {src}")
                            if os.path.exists(src):
                                img = ReportLabImage(src, width=200, height=150)
                                log.debug("Adding ReportLab Image object to story (Local)")
                                story.append(img)
                                story.append(Spacer(1, 10))
                            else:
                               log.error(f"Local image file not found: {src}")
                               story.append(Paragraph(f"[Image locale non trouvee: {src}]", styles["CustomNormal"]))
                               story.append(Spacer(1, 6))
                    except requests.exceptions.RequestException as e:
                        log.error(f"Network error loading image {src}: {e}")
                        story.append(Paragraph(f"[Image (erreur reseau): {alt}]", styles["CustomNormal"]))
                        story.append(Spacer(1, 6))
                    except Exception as e:
                        log.error(f"Error processing image {src}: {e}", exc_info=True) 
                        story.append(Paragraph(f"[Image: {alt}]", styles["CustomNormal"]))
                        story.append(Spacer(1, 6))
                else:
                    log.warning("Image tag found with no 'src' attribute.")
                    story.append(Paragraph(f"[Image: {alt} (source manquante)]", styles["CustomNormal"]))
                    story.append(Spacer(1, 6))
            elif tag_name == "br":
                log.debug("Adding Spacer for <br>")
                story.append(Spacer(1, 6))
            else:
                text = elem.get_text().strip()
                if text:
                    log.debug(f"Adding Paragraph for unknown tag <{tag_name}>: {text[:50]}...")
                    story.append(Paragraph(render_text_with_emojis(text), styles["CustomNormal"]))
                    story.append(Spacer(1, 6))
    log.debug(f"Finished render_html_elements. Story contains {len(story)} elements.")
    return story

def _cleanup_files(folder_path: str, delay_minutes: int):
    def delete_files():
        time.sleep(delay_minutes * 60)
        try:
            import shutil
            shutil.rmtree(folder_path) 
            log.debug(f"Folder {folder_path} deleted.")
        except Exception as e:
            logging.error(f"Error deleting files : {e}")
    thread = threading.Thread(target=delete_files)
    thread.start()

def _convert_markdown_to_structured(markdown_content):
    """
    Converts Markdown content into a structured format for Word
    
    Args:
        markdown_content (str): Markdown content
        
    Returns:
        list: List of objects with 'text' and 'type'
    """
    if not markdown_content or not isinstance(markdown_content, str):
        return []
    
    lines = markdown_content.split('\n')
    structured = []
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        if line.startswith('# '):
            structured.append({"text": line[2:].strip(), "type": "title"})
        elif line.startswith('## '):
            structured.append({"text": line[3:].strip(), "type": "heading"})
        elif line.startswith('### '):
            structured.append({"text": line[4:].strip(), "type": "subheading"})
        elif line.startswith('#### '):
            structured.append({"text": line[5:].strip(), "type": "subheading"})
        elif line.startswith('- '):
            structured.append({"text": line[2:].strip(), "type": "bullet"})
        elif line.startswith('* '):
            structured.append({"text": line[2:].strip(), "type": "bullet"})
        elif line.startswith('**') and line.endswith('**'):
            structured.append({"text": line[2:-2].strip(), "type": "bold"})
        else:
            structured.append({"text": line, "type": "paragraph"})
    
    return structured

def _create_excel(data: list[list[str]], filename: str, folder_path: str | None = None, title: str | None = None) -> dict:
    log.debug("Creating Excel file with optional template")
    if folder_path is None:
        folder_path = _generate_unique_folder()
    
    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "xlsx")

    # Detect structured content (list of dicts with 'type' field)
    is_structured = (
        isinstance(data, list)
        and data
        and isinstance(data[0], dict)
        and "type" in data[0]
    )
    
    # Detect raw 2D matrix (list of lists with primitive values)
    is_raw_matrix = (
        isinstance(data, list)
        and data
        and isinstance(data[0], list)
        and all(
            isinstance(cell, (str, int, float, bool, type(None)))
            for row in data[:5]
            for cell in row[:10]
        )
    )
    
    # For raw 2D matrices and structured content: create clean workbook
    if is_raw_matrix or is_structured:
        log.debug("Raw 2D matrix or structured content detected, creating clean workbook")
        wb = Workbook()
    elif XLSX_TEMPLATE:
        try:
            log.debug("Loading XLSX template...")
            wb = load_workbook(XLSX_TEMPLATE_PATH) 
            log.debug(f"Template loaded with {len(wb.sheetnames)} sheet(s)")
        except Exception as e:
            log.warning(f"Failed to load XLSX template: {e}")
            wb = Workbook()
    else:
        log.debug("No XLSX template available, creating new workbook")
        wb = Workbook()

    ws = wb.active

    from openpyxl.utils import get_column_letter 

    
    if title:
        ws.title = "".join(c for c in title if c.isalnum() or c in (' ', '-', '_')).strip()[:31]
        title_cell_found = False
        for row in ws.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str) and "title" in cell.value.lower():
                    cell.value = title 
                    from openpyxl.styles import Font
                    log.debug(f"Title '{title}' replaced in the cell {get_column_letter(cell.column)}{cell.row} containing  'title'")
                    title_cell_found = True
                    break
            if title_cell_found:
                break
    
    start_row, start_col = 1, 1
    if ws.auto_filter and ws.auto_filter.ref:
        try:
            from openpyxl.utils import range_boundaries
            start_col, start_row, _, _ = range_boundaries(ws.auto_filter.ref)
        except: pass

    if not data:
        wb.save(filepath)
        return {"success": True, "filepath": filepath, "filename": filename}

    template_border = ws.cell(start_row, start_col).border
    has_borders = template_border and any([template_border.top.style, template_border.bottom.style, 
                                          template_border.left.style, template_border.right.style])
    
    for r in range(max(len(data) + 10, 50)):
        for c in range(max(len(data[0]) + 5, 20)):
            cell = ws.cell(row=start_row + r, column=start_col + c)
            
            if r < len(data) and c < len(data[0]):
                cell.value = data[r][c]
                if r == 0 and data[r][c]:  
                    from openpyxl.styles import Font
                    cell.font = Font(bold=True)
                if has_borders:  
                    from openpyxl.styles import Border
                    cell.border = Border(top=template_border.top, bottom=template_border.bottom,
                                       left=template_border.left, right=template_border.right)
            else:
                cell.value = None
                if cell.has_style:
                    from openpyxl.styles import Font, PatternFill, Border, Alignment
                    cell.font, cell.fill, cell.border, cell.alignment = Font(), PatternFill(), Border(), Alignment()

    if ws.auto_filter:
        ws.auto_filter.ref = f"{get_column_letter(start_col)}{start_row}:{get_column_letter(start_col + len(data[0]) - 1)}{start_row + len(data) - 1}"
    
    for c in range(len(data[0])):
        max_len = max(len(str(data[r][c])) for r in range(len(data)))
        ws.column_dimensions[get_column_letter(start_col + c)].width = min(max_len + 2, 150)

    wb.save(filepath)

    return {"url": _public_url(folder_path, fname), "path": filepath}
def _create_csv(data: list[list[str]], filename: str, folder_path: str | None = None) -> dict:
    log.debug("Creating CSV file")
    if folder_path is None:
        folder_path = _generate_unique_folder()

    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "csv")

    with open(filepath, "w", newline="", encoding="utf-8") as f:
        if isinstance(data, list):
            csv.writer(f).writerows(data)
        else:
            csv.writer(f).writerow([data])

    return {"url": _public_url(folder_path, fname), "path": filepath}

def _create_pdf(text: str | list[str], filename: str, folder_path: str | None = None) -> dict:    
    log.debug("Creating PDF file")
    if folder_path is None:
        folder_path = _generate_unique_folder()
    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "pdf")

    md_parts = []
    if isinstance(text, list):
        for item in text:
            if isinstance(item, str):
                md_parts.append(item)
            elif isinstance(item, dict):
                t = item.get("type")
                if t == "title":
                    md_parts.append(f"# {item.get('text','')}")
                elif t == "subtitle":
                    md_parts.append(f"## {item.get('text','')}")
                elif t == "paragraph":
                    md_parts.append(item.get("text",""))
                elif t == "list":
                    md_parts.append("\n".join([f"- {x}" for x in item.get("items",[])]))
                elif t in ("image","image_query"):
                    query = item.get("query","")
                    if query:
                        md_parts.append(f"![Image](image_query: {query})")
    else:
        md_parts = [str(text or "")]
        
    md_text = "\n\n".join(md_parts)    
   
    def replace_image_query(match):
        query = match.group(1).strip()
        image_url = search_image(query)
        return f'\n\n<img src="{image_url}" alt="Image: {query}" />\n\n' if image_url else ""

    md_text = re.sub(r'!\[[^\]]*\]\(\s*image_query:\s*([^)]+)\)', replace_image_query, md_text)
    html = markdown2.markdown(md_text, extras=['fenced-code-blocks','tables','break-on-newline','cuddled-lists'])
    soup = BeautifulSoup(html, "html.parser")
    story = render_html_elements(soup) or [Paragraph("Empty Content", styles["CustomNormal"])]

    doc = SimpleDocTemplate(filepath, topMargin=72, bottomMargin=72, leftMargin=72, rightMargin=72)
    try:
        doc.build(story)
    except Exception as e:
        log.error(f"Error building PDF {fname}: {e}", exc_info=True)
        doc2 = SimpleDocTemplate(filepath)
        doc2.build([Paragraph("Error in PDF generation", styles["CustomNormal"])])

    return {"url": _public_url(folder_path, fname), "path": filepath}

def _create_presentation(slides_data: list[dict], filename: str, folder_path: str | None = None, title: str | None = None) -> dict:
    if folder_path is None:
        folder_path = _generate_unique_folder()
    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "pptx")
      
    use_template = False
    prs = None
    title_layout = None
    content_layout = None

    if PPTX_TEMPLATE:
        try:
            log.debug("Attempting to load template...")
            src = PPTX_TEMPLATE
            if hasattr(PPTX_TEMPLATE, "slides") and hasattr(PPTX_TEMPLATE, "save"):
                log.debug("Template is a Presentation object, converting to BytesIO")
                buf = BytesIO()
                PPTX_TEMPLATE.save(buf); buf.seek(0)
                src = buf

            tmp = Presentation(src)
            log.debug(f"Template loaded with {len(tmp.slides)} slides")
            if len(tmp.slides) >= 1:
                prs = tmp
                use_template = True

                title_layout = prs.slides[0].slide_layout
                content_layout = prs.slides[1].slide_layout if len(prs.slides) >= 2 else prs.slides[0].slide_layout
                log.debug("Using template layouts")

                for i in range(len(prs.slides) - 1, 0, -1):
                    rId = prs.slides._sldIdLst[i].rId 
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[i]
        except Exception:
            log.error(f"Error loading template: {e}")
            use_template = False
            prs = None

    if not use_template:
        log.debug("No valid template, creating new presentation with default layouts")
        prs = Presentation()
        title_layout = prs.slide_layouts[0]
        content_layout = prs.slide_layouts[1]

    if use_template:
        log.debug("Using template title slide")
        tslide = prs.slides[0]
        if tslide.shapes.title:
            tslide.shapes.title.text = title or ""
            for p in tslide.shapes.title.text_frame.paragraphs:
                for r in p.runs:
                    title_info = next(({'size': PptPt(int(child.attrib.get('sz', 2800))/100), 'bold': child.attrib.get('b', '0') == '1'} for child in title_layout.element.iter() if 'defRPr' in child.tag.split('}')[-1] and 'sz' in child.attrib), {'size': PptPt(28), 'bold': True})

                    r.font.size = title_info['size'] 
                    r.font.bold = title_info['bold']
    else:
        log.debug("Creating new title slide")
        tslide = prs.slides.add_slide(title_layout)
        if tslide.shapes.title:
            tslide.shapes.title.text = title or ""
            for p in tslide.shapes.title.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = PptPt(28); r.font.bold = True

    EMU_PER_IN = 914400
    slide_w_in = prs.slide_width / EMU_PER_IN
    slide_h_in = prs.slide_height / EMU_PER_IN
    log.debug(f"Slide dimensions: {slide_w_in} x {slide_h_in} inches")

    page_margin = 0.5
    gutter = 0.3

    for i, slide_data in enumerate(slides_data):
        log.debug(f"Processing slide {i+1}: {slide_data.get('title', 'Untitled')}")
        if not isinstance(slide_data, dict):
            log.warning(f"Slide data is not a dict, skipping slide {i+1}")
            continue

        slide_title = slide_data.get("title", "Untitled")
        content_list = slide_data.get("content", [])
        if not isinstance(content_list, list):
            content_list = [content_list]
        log.debug(f"Adding slide with title: '{slide_title}'")
        slide = prs.slides.add_slide(content_layout)

        if slide.shapes.title:
            slide.shapes.title.text = slide_title
            for p in slide.shapes.title.text_frame.paragraphs:
                for r in p.runs:
                    title_info = next(({'size': PptPt(int(child.attrib.get('sz', 2800))/100), 'bold': child.attrib.get('b', '0') == '1'} for child in content_layout.element.iter() if 'defRPr' in child.tag.split('}')[-1] and 'sz' in child.attrib), {'size': PptPt(28), 'bold': True})

                    r.font.size = title_info['size'] 
                    r.font.bold = title_info['bold']

        content_shape = None
        try:
            for ph in slide.placeholders:
                try:
                    if ph.placeholder_format.idx == 1:
                        content_shape = ph; break
                except Exception:
                    pass
            if content_shape is None:
                for ph in slide.placeholders:
                    try:
                        if ph.placeholder_format.idx != 0:
                            content_shape = ph; break
                    except Exception:
                        pass
        except Exception:
            log.error(f"Error finding content placeholder: {e}")
            pass

        title_bottom_in = 1.0 
        if slide.shapes.title:
            try:
                title_bottom_emu = slide.shapes.title.top + slide.shapes.title.height
                title_bottom_in = max(title_bottom_emu / EMU_PER_IN, 1.0)
                title_bottom_in += 0.2
            except Exception:
                title_bottom_in = 1.2 

        if content_shape is None:

            content_shape = slide.shapes.add_textbox(Inches(page_margin), Inches(title_bottom_in), Inches(slide_w_in - 2*page_margin), Inches(slide_h_in - title_bottom_in - page_margin))
            log.debug("Creating new textbox for content")
        tf = content_shape.text_frame
        try:
            tf.clear()
        except Exception:
            log.error(f"Error clearing text frame: {e}")
            pass
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        try:
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.05)
            tf.margin_bottom = Inches(0.05)
        except Exception:
            pass

        content_left_in, content_top_in = page_margin, title_bottom_in
        content_width_in = slide_w_in - 2*page_margin
        content_height_in = slide_h_in - (title_bottom_in + page_margin)

        image_query = slide_data.get("image_query")
        if image_query:
            image_url = search_image(image_query)
            if image_url:
                log.debug(f"Searching for image query: '{image_query}'")
                try:
                    log.debug(f"Downloading image from URL: {image_url}")
                    response = requests.get(image_url, timeout=30)
                    response.raise_for_status()
                    image_data = response.content
                    image_stream = BytesIO(image_data)
                    pos = slide_data.get("image_position", "right")
                    size = slide_data.get("image_size", "medium")
                    if size == "small":
                        img_w_in, img_h_in = 2.0, 1.5
                    elif size == "large":
                        img_w_in, img_h_in = 4.0, 3.0
                    else:
                        img_w_in, img_h_in = 3.0, 2.0
                    log.debug(f"Image dimensions: {img_w_in} x {img_h_in} inches")

                    if pos == "left":
                        img_left_in = page_margin
                        img_top_in = title_bottom_in
                        content_left_in = img_left_in + img_w_in + gutter
                        content_top_in = title_bottom_in
                        content_width_in = max(slide_w_in - page_margin - content_left_in, 2.5)
                        content_height_in = slide_h_in - (title_bottom_in + page_margin)
                    elif pos == "right":
                        img_left_in = max(slide_w_in - page_margin - img_w_in, page_margin)
                        img_top_in = title_bottom_in
                        content_left_in = page_margin
                        content_top_in = title_bottom_in
                        content_width_in = max(img_left_in - gutter - content_left_in, 2.5)
                        content_height_in = slide_h_in - (title_bottom_in + page_margin)
                    elif pos == "top":
                        img_left_in = slide_w_in - page_margin - img_w_in
                        img_top_in = title_bottom_in
                        content_left_in = page_margin
                        content_top_in = img_top_in + img_h_in + gutter
                        content_width_in = slide_w_in - 2*page_margin
                        content_height_in = max(slide_h_in - page_margin - content_top_in, 2.0)
                    elif pos == "bottom":
                        img_left_in = slide_w_in - page_margin - img_w_in
                        img_top_in = max(slide_h_in - page_margin - img_h_in, page_margin)
                        content_left_in = page_margin
                        content_top_in = title_bottom_in
                        content_width_in = slide_w_in - 2*page_margin
                        content_height_in = max(img_top_in - gutter - content_top_in, 2.0)
                    else:
                        img_left_in = max(slide_w_in - page_margin - img_w_in, page_margin)
                        img_top_in = title_bottom_in
                        content_left_in = page_margin
                        content_top_in = title_bottom_in
                        content_width_in = max(img_left_in - gutter - content_left_in, 2.5)
                        content_height_in = slide_h_in - (title_bottom_in + page_margin)

                    slide.shapes.add_picture(image_stream, Inches(img_left_in), Inches(img_top_in), Inches(img_w_in), Inches(img_h_in))
                    log.debug(f"Image added at position: left={img_left_in}, top={img_top_in}")
                except Exception:
                    pass

        try:
            content_shape.left = Inches(content_left_in)
            content_shape.top = Inches(content_top_in)
            content_shape.width = Inches(content_width_in)
            content_shape.height = Inches(content_height_in)
        except Exception:
            pass

        approx_chars_per_in = 9.5
        approx_lines_per_in = 1.6
        safe_width = max(content_width_in, 0.1)
        safe_height = max(content_height_in, 0.1)
        est_capacity = int(safe_width * approx_chars_per_in * safe_height * approx_lines_per_in)
        font_size = dynamic_font_size(content_list, max_chars=max(est_capacity, 120), base_size=24, min_size=12)

        try:
            tf = content_shape.text_frame
        except Exception:
            try:
                tf = content_shape.text_frame
            except Exception:
                log.warning("Could not access text frame for content shape")
                continue

        if not tf.paragraphs:
            tf.add_paragraph()
        for idx, line in enumerate(content_list):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = str(line) if line is not None else ""
            run.font.size = font_size
            p.space_after = PptPt(6)

    prs.save(filepath)
    return {"url": _public_url(folder_path, fname), "path": filepath}

def _create_word(content: list[dict] | str, filename: str, folder_path: str | None = None, title: str | None = None) -> dict:
    log.debug("Creating Word document")

    if isinstance(content, str):
        content = _convert_markdown_to_structured(content)
    elif not isinstance(content, list):
        content = []

    if folder_path is None:
        folder_path = _generate_unique_folder()
    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "docx")

    use_template = False
    doc = None

    if DOCX_TEMPLATE:
        try:
            src = DOCX_TEMPLATE
            if hasattr(DOCX_TEMPLATE, "paragraphs") and hasattr(DOCX_TEMPLATE, "save"):
                buf = BytesIO()
                DOCX_TEMPLATE.save(buf)
                buf.seek(0)
                src = buf

            doc = Document(src)
            use_template = True
            log.debug("Using DOCX template")

            for element in doc.element.body:
                if element.tag.endswith('}p') or element.tag.endswith('}tbl'):
                    doc.element.body.remove(element)

        except Exception as e:
            log.warning(f"Failed to load DOCX template: {e}")
            use_template = False
            doc = None

    if not use_template:
        doc = Document()
        log.debug("Creating new Word document without template")

    if title:
        title_paragraph = doc.add_paragraph(title)
        try:
            title_paragraph.style = doc.styles['Title']
        except KeyError:
            try:
                title_paragraph.style = doc.styles['Heading 1']
            except KeyError:
                run = title_paragraph.runs[0] if title_paragraph.runs else title_paragraph.add_run()
                run.font.size = DocxPt(20)
                run.font.bold = True
        title_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        log.debug("Document title added")

    for item in content or []:
        if isinstance(item, str):
            doc.add_paragraph(item)
        elif isinstance(item, dict):
            if item.get("type") == "image_query":
                new_item = {
                    "type": "image",
                    "query": item.get("query")
                }
                image_query = new_item.get("query")
                if image_query:
                    log.debug(f"Image search for the query : {image_query}")
                    image_url = search_image(image_query)
                    if image_url:
                        response = requests.get(image_url)
                        image_data = BytesIO(response.content)
                        doc.add_picture(image_data, width=Inches(6))
                        log.debug("Image successfully added")
                    else:
                        log.warning(f"Image search for : '{image_query}'")
            elif "type" in item:
                item_type = item.get("type")
                if item_type == "title":
                    paragraph = doc.add_paragraph(item.get("text", ""))
                    try:
                        paragraph.style = doc.styles['Heading 1']
                    except KeyError:
                        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                        run.font.size = DocxPt(18)
                        run.font.bold = True
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    log.debug("Title added")
                elif item_type == "subtitle":
                    paragraph = doc.add_paragraph(item.get("text", ""))
                    try:
                        paragraph.style = doc.styles['Heading 2']
                    except KeyError:
                        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                        run.font.size = DocxPt(16)
                        run.font.bold = True
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    log.debug("Subtitle added")
                elif item_type == "paragraph":
                    doc.add_paragraph(item.get("text", ""))
                    log.debug("Paragraph added")
                elif item_type == "list":
                    items = item.get("items", [])
                    for i, item_text in enumerate(items):
                        paragraph = doc.add_paragraph(item_text)
                        try:
                            paragraph.style = doc.styles['List Bullet']
                        except KeyError:
                            paragraph.style = doc.styles['Normal']
                    log.debug("List added")
                elif item_type == "image":
                    image_query = item.get("query")
                    if image_query:
                        log.debug(f"Image search for the query : {image_query}")
                        image_url = search_image(image_query)
                        if image_url:
                            response = requests.get(image_url)
                            image_data = BytesIO(response.content)
                            doc.add_picture(image_data, width=Inches(6))
                            log.debug("Image successfully added")
                        else:
                            log.warning(f"Image search for : '{image_query}'")
                elif item_type == "table":
                    data = item.get("data", [])
                    if data:
                        template_table_style = None
                        if use_template and DOCX_TEMPLATE:
                            try:
                                for table in DOCX_TEMPLATE.tables:
                                    if table.style:
                                        template_table_style = table.style
                                        break
                            except Exception:
                                pass
                        
                        table = doc.add_table(rows=len(data), cols=len(data[0]) if data else 0)
                        
                        if template_table_style:
                            try:
                                table.style = template_table_style
                                log.debug(f"Applied template table style: {template_table_style.name}")
                            except Exception as e:
                                log.debug(f"Could not apply template table style: {e}")
                        else:
                            try:
                                for style_name in ['Table Grid', 'Light Grid Accent 1', 'Medium Grid 1 Accent 1', 'Light List Accent 1']:
                                    try:
                                        table.style = doc.styles[style_name]
                                        log.debug(f"Applied built-in table style: {style_name}")
                                        break
                                    except KeyError:
                                        continue
                            except Exception as e:
                                log.debug(f"Could not apply any table style: {e}")
                        
                        for i, row in enumerate(data):
                            for j, cell in enumerate(row):
                                cell_obj = table.cell(i, j)
                                cell_obj.text = str(cell)
                                
                                if i == 0:
                                    for paragraph in cell_obj.paragraphs:
                                        for run in paragraph.runs:
                                            run.font.bold = True
                                        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        
                        if not template_table_style:
                            try:
                                tbl = table._tbl
                                tblPr = tbl.tblPr
                                tblBorders = parse_xml(r'<w:tblBorders {}><w:top w:val="single" w:sz="4" w:space="0" w:color="000000"/><w:left w:val="single" w:sz="4" w:space="0" w:color="000000"/><w:bottom w:val="single" w:sz="4" w:space="0" w:color="000000"/><w:right w:val="single" w:sz="4" w:space="0" w:color="000000"/><w:insideH w:val="single" w:sz="4" w:space="0" w:color="000000"/><w:insideV w:val="single" w:sz="4" w:space="0" w:color="000000"/></w:tblBorders>'.format(nsdecls('w')))
                                tblPr.append(tblBorders)
                            except Exception as e:
                                log.debug(f"Could not add table borders: {e}")
                        
                        log.debug("Table added with improved styling")
            elif "text" in item:
                doc.add_paragraph(item["text"])
                log.debug("Paragraph added")
    
    doc.save(filepath)
    return {"url": _public_url(folder_path, fname), "path": filepath}

def _create_raw_file(content: str, filename: str | None, folder_path: str | None = None) -> dict:
    log.debug("Creating raw file")
    if folder_path is None:
        folder_path = _generate_unique_folder()

    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "txt")

    if fname.lower().endswith(".xml") and isinstance(content, str) and not content.strip().startswith("<?xml"):
        content = f'<?xml version="1.0" encoding="UTF-8"?>\n{content}'

    with open(filepath, "w", encoding="utf-8") as f:
        f.write(content or "")

    return {"url": _public_url(folder_path, fname), "path": filepath}

def upload_file(file_path: str, filename: str, file_type: str, token: str) -> dict:
    """
    Upload a file to OpenWebUI server.
    """
    url = f"{URL}/api/v1/files/"
    headers = {
        'Authorization': token,
        'Accept': 'application/json'
    }
    
    with open(file_path, 'rb') as f:
        files = {'file': f}
        response = post(url, headers=headers, files=files)

    if response.status_code != 200:
        return {"error": {"message": f'Error uploading file: {response.status_code}'}}
    else:
        return {
            "file_path_download": f"[Download {filename}.{file_type}](/api/v1/files/{response.json()['id']}/content)"
        }

def download_file(file_id: str, token: str) -> BytesIO:
    """
    Download a file from OpenWebUI server.
    """
   
    url = f"{URL}/api/v1/files/{file_id}/content"
    headers = {
        'Authorization': token,
        'Accept': 'application/json'
    }
    
    response = get(url, headers=headers)
    
    if response.status_code != 200:
        return {"error": {"message": f'Error downloading the file: {response.status_code}'}}
    else:
        return BytesIO(response._content)

def _extract_paragraph_style_info(para):
    """Extract detailed style information from a paragraph"""
    if not para.runs:
        return {}
    
    first_run = para.runs[0]
    return {
        "font_name": first_run.font.name,
        "font_size": first_run.font.size,
        "bold": first_run.font.bold,
        "italic": first_run.font.italic,
        "underline": first_run.font.underline,
        "color": first_run.font.color.rgb if first_run.font.color else None
    }

def _extract_cell_style_info(cell):
    """Extract style information from a cell"""
    return {
        "style": cell.style.name if hasattr(cell, 'style') else None,
        "text_alignment": cell.paragraphs[0].alignment if cell.paragraphs else None
    }

def _apply_text_to_paragraph(para, new_text):
    """
    Apply new text to a paragraph while preserving formatting.
    """
    original_style = para.style
    original_alignment = para.alignment
    
    original_run_format = None
    if para.runs:
        first_run = para.runs[0]
        original_run_format = {
            "font_name": first_run.font.name,
            "font_size": first_run.font.size,
            "bold": first_run.font.bold,
            "italic": first_run.font.italic,
            "underline": first_run.font.underline,
            "color": first_run.font.color.rgb if first_run.font.color and first_run.font.color.rgb else None
        }
    
    for _ in range(len(para.runs)):
        para._element.remove(para.runs[0]._element)
    
    if isinstance(new_text, list):
        for i, text_item in enumerate(new_text):
            if i > 0:
                para.add_run("\n")
            run = para.add_run(str(text_item))
            if original_run_format:
                _apply_run_formatting(run, original_run_format)
    else:
        run = para.add_run(str(new_text))
        if original_run_format:
            _apply_run_formatting(run, original_run_format)
    
    if original_style:
        try:
            para.style = original_style
        except Exception:
            pass
    if original_alignment is not None:
        try:
            para.alignment = original_alignment
        except Exception:
            pass


def _apply_run_formatting(run, format_dict):
    """
    Apply formatting from a dict to a run.
    """
    try:
        if format_dict.get("font_name"):
            run.font.name = format_dict["font_name"]
    except Exception:
        pass
    
    try:
        if format_dict.get("font_size"):
            run.font.size = format_dict["font_size"]
    except Exception:
        pass
    
    try:
        if format_dict.get("bold") is not None:
            run.font.bold = format_dict["bold"]
    except Exception:
        pass
    
    try:
        if format_dict.get("italic") is not None:
            run.font.italic = format_dict["italic"]
    except Exception:
        pass
    
    try:
        if format_dict.get("underline") is not None:
            run.font.underline = format_dict["underline"]
    except Exception:
        pass
    
    try:
        if format_dict.get("color"):
            from docx.shared import RGBColor
            run.font.color.rgb = format_dict["color"]
    except Exception:
        pass

@mcp.tool(
    name="full_context_document",
    title="Return the structure of a document (docx, xlsx, pptx)",
    description="Return the structure, content, and metadata of a document based on its type (docx, xlsx, pptx). Unified output format with index, type, style, and text."
)

async def full_context_document(
    file_id: str,
    file_name: str,
    headers: dict = None,
    ctx: Context[ServerSession, None] = None
) -> dict:
    """
    Return the structure of a document (docx, xlsx, pptx) based on its file extension.
    The function detects the file type and processes it accordingly.
    Returns:
        dict: A JSON object with the structure of the document.
    """
    user_token = TOKEN
    if headers:
        auth_header = headers.get("authorization")
        if auth_header:
            user_token = auth_header
            logging.info("Using authorization from MCPO forwarded headers")
        else:
            logging.warning("Forwarded headers present but no authorization found")
    else:
        logging.info("ℹNo forwarded headers, using admin TOKEN fallback")
    try:
        user_file = download_file(file_id=file_id,token=user_token) 

        if isinstance(user_file, dict) and "error" in user_file:
            return json.dumps(user_file, indent=4, ensure_ascii=False)

        file_extension = os.path.splitext(file_name)[1].lower()
        file_type = file_extension.lstrip('.')

        structure = {
            "file_name": file_name,
            "file_id": file_id,
            "type": file_type,
            "slide_id_order": [],
            "body": [],
        }
        index_counter = 1

        if file_type == "docx":
            doc = Document(user_file)
            
            para_id_counter = 1
            
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                
                style = para.style.name
                style_info = _extract_paragraph_style_info(para)
                element_type = "heading" if style.startswith("Heading") else "paragraph"
                
                para_xml_id = para_id_counter
                
                structure["body"].append({
                    "index": para_id_counter,
                    "para_xml_id": para_xml_id,
                    "id_key": f"pid:{para_xml_id}",
                    "type": element_type,
                    "style": style,
                    "style_info": style_info,
                    "text": text
                })
                para_id_counter += 1
            
            for table_idx, table in enumerate(doc.tables):
                table_xml_id = id(table._element)
                table_info = {
                    "index": para_id_counter,
                    "table_xml_id": table_xml_id,
                    "id_key": f"tid:{table_xml_id}",
                    "type": "table",
                    "style": "Table",
                    "table_id": table_idx,
                    "rows": len(table.rows),
                    "columns": len(table.rows[0].cells) if table.rows else 0,
                    "cells": []
                }
                
                for row_idx, row in enumerate(table.rows):
                    row_data = []
                    for col_idx, cell in enumerate(row.cells):
                        cell_xml_id = id(cell._element)
                        cell_text = cell.text.strip()
                        cell_data = {
                            "row": row_idx,
                            "column": col_idx,
                            "cell_xml_id": cell_xml_id,
                            "id_key": f"tid:{table_xml_id}/cid:{cell_xml_id}",
                            "text": cell_text,
                            "style": cell.style.name if hasattr(cell, 'style') else None
                        }
                        row_data.append(cell_data)
                    table_info["cells"].append(row_data)
                
                structure["body"].append(table_info)
                para_id_counter += 1

        elif file_type == "xlsx":
            wb = load_workbook(user_file, read_only=True, data_only=True)

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                for row_idx, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                    for col_idx, cell in enumerate(row, start=1):
                        if cell is None or str(cell).strip() == "":
                            continue
                        col_letter = sheet.cell(row=row_idx, column=col_idx).column_letter
                        cell_ref = f"{col_letter}{row_idx}"
                        structure["body"].append({
                            "index": cell_ref,
                            "type": "cell",
                            "text": str(cell)
                        })
                        index_counter += 1

        elif file_type == "pptx":
            prs = Presentation(user_file)
            structure["slide_id_order"] = [int(s.slide_id) for s in prs.slides]
            for slide_idx, slide in enumerate(prs.slides):
                title_shape = slide.shapes.title if hasattr(slide.shapes, "title") else None
                title_text = title_shape.text.strip() if (title_shape and getattr(title_shape, "text", "").strip()) else ""

                slide_obj = {
                    "index": slide_idx,
                    "slide_id": int(slide.slide_id),
                    "id_key": f"sid:{int(slide.slide_id)}",
                    "title": title_text,
                    "shapes": []
                }
                

                for shape_idx, shape in enumerate(slide.shapes):
                    key = f"s{slide_idx}/sh{shape_idx}"
                    if hasattr(shape, "image"):
                        shape_id_val = getattr(shape, "shape_id", None) or shape._element.cNvPr.id  
                        slide_obj["shapes"].append({
                            "shape_idx": shape_idx,
                            "shape_id": shape_id_val,
                            "idx_key": key, 
                            "id_key": f"sid:{int(slide.slide_id)}/shid:{int(shape_id_val)}",
                            "kind": "image"
                        })
                        continue

                    if hasattr(shape, "text_frame") and shape.text_frame:
                        kind = "title" if (title_shape is not None and shape is title_shape) else "textbox"

                        paragraphs = []
                        for p in shape.text_frame.paragraphs:
                            text = "".join(run.text for run in p.runs) if p.runs else p.text
                            text = (text or "").strip()
                            if text != "":
                                paragraphs.append(text)

                        shape_id_val = getattr(shape, "shape_id", None) or shape._element.cNvPr.id 
                        slide_obj["shapes"].append({
                            "shape_idx": shape_idx,
                            "shape_id": shape_id_val,
                            "idx_key": key, 
                            "id_key": f"sid:{int(slide.slide_id)}/shid:{int(shape_id_val)}",
                            "kind": kind,
                            "paragraphs": paragraphs
                        })
                        continue

                    # --- NEW: table extraction ---
                    if getattr(shape, "has_table", False):
                        tbl = shape.table
                        rows = []
                        for r in tbl.rows:
                            row_cells = []
                            for c in r.cells:
                                # collect full text
                                if hasattr(c, "text_frame") and c.text_frame:
                                    paras = []
                                    for p in c.text_frame.paragraphs:
                                        t = "".join(run.text for run in p.runs) if p.runs else p.text
                                        t = (t or "").strip()
                                        if t:
                                            paras.append(t)
                                    cell_text = "\n".join(paras)
                                else:
                                    cell_text = (getattr(c, "text", "") or "").strip()
                                row_cells.append(cell_text)
                            rows.append(row_cells)

                        shape_id_val = getattr(shape, "shape_id", None) or shape._element.cNvPr.id
                        slide_obj["shapes"].append({
                            "shape_idx": shape_idx,
                            "shape_id": shape_id_val,
                            "idx_key": key,
                            "id_key": f"sid:{int(slide.slide_id)}/shid:{int(shape_id_val)}",
                            "kind": "table",
                            "rows": rows  # list of lists: each inner list = one row's cell texts
                        })
                        continue

                structure["body"].append(slide_obj)

        else:
            return json.dumps({
                "error": {"message": f"Unsupported file type: {file_type}. Only docx, xlsx, and pptx are supported."}
            }, indent=4, ensure_ascii=False)

        return json.dumps(structure, indent=4, ensure_ascii=False)

    except Exception as e:
        return json.dumps({"error": {"message": str(e)}}, indent=4, ensure_ascii=False)

def add_auto_sized_review_comment(cell, text, author="AI Reviewer"):
    """
    Adds a note to an Excel cell, adjusting the width and height
    so that all the text is visible.
    """
    if not text:
        return

    avg_char_width = 7
    px_per_line = 15
    base_width = 200
    max_width = 500
    min_height = 40

    width = min(max_width, base_width + len(text) * 2)
    chars_per_line = max(1, width // avg_char_width)
    lines = 0
    for paragraph in text.split('\n'):
        lines += math.ceil(len(paragraph) / chars_per_line)
    height = max(min_height, lines * px_per_line)

    comment = Comment(text, author)
    comment.width = width
    comment.height = height
    cell.comment = comment


def _snapshot_runs(p):
    """Return a list of {'text': str, 'font': {...}} for each run in a paragraph."""
    runs = []
    for r in p.runs:
        f = r.font
        font_spec = {
            "name": f.name,
            "size": f.size,
            "bold": f.bold,
            "italic": f.italic,
            "underline": f.underline,
            "color_rgb": getattr(getattr(f.color, "rgb", None), "rgb", None) or getattr(f.color, "rgb", None)
        }
        runs.append({"text": r.text or "", "font": font_spec})
    return runs

def _apply_font(run, font_spec):
    """Apply font specifications to a run."""
    if not font_spec:
        return
    f = run.font
    try:
        if font_spec.get("name") is not None:
            f.name = font_spec["name"]
        if font_spec.get("size") is not None:
            f.size = font_spec["size"]
        if font_spec.get("bold") is not None:
            f.bold = font_spec["bold"]
        if font_spec.get("italic") is not None:
            f.italic = font_spec["italic"]
        if font_spec.get("underline") is not None:
            f.underline = font_spec["underline"]
        rgb = font_spec.get("color_rgb")
        if rgb is not None:
            try:
                f.color.rgb = rgb
            except Exception:
                pass
    except Exception:
        pass

def _set_text_with_runs(shape, new_content):
    """
    Set the text of a shape while preserving the original run-level formatting.
    """

    if not (hasattr(shape, "text_frame") and shape.text_frame):
        return
    tf = shape.text_frame

    if isinstance(new_content, list):
        lines = [str(item) for item in new_content]
    else:
        lines = [str(new_content or "")]

    original_para_styles = []
    original_para_runs = []     

    for p in tf.paragraphs:
        level = int(getattr(p, "level", 0) or 0)
        alignment = getattr(p, "alignment", None)
        original_para_styles.append({"level": level, "alignment": alignment})
        original_para_runs.append(_snapshot_runs(p))

    tf.clear()

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if (i == 0 and tf.paragraphs) else tf.add_paragraph()

        if original_para_styles:
            style = original_para_styles[i] if i < len(original_para_styles) else original_para_styles[-1]
            p.level = style.get("level", 0)
            if style.get("alignment") is not None:
                p.alignment = style["alignment"]

        runs_spec = (
            original_para_runs[i] if i < len(original_para_runs)
            else (original_para_runs[-1] if original_para_runs else [])
        )

        if not runs_spec:
            r = p.add_run()
            r.text = line
            continue

        n = len(runs_spec)
        total = len(line)

        if total == 0:
            for spec in runs_spec:
                r = p.add_run()
                r.text = ""
                _apply_font(r, spec["font"])
        else:
            base, rem = divmod(total, n)
            sizes = [base + (1 if k < rem else 0) for k in range(n)]
            pos = 0
            for k, spec in enumerate(runs_spec):
                seg = line[pos:pos + sizes[k]]
                pos += sizes[k]
                r = p.add_run()
                r.text = seg
                _apply_font(r, spec["font"])

def shape_by_id(slide, shape_id):
    sid = int(shape_id)
    for sh in slide.shapes:
        val = getattr(sh, "shape_id", None) or getattr(getattr(sh, "_element", None), "cNvPr", None)
        val = int(getattr(val, "id", val)) if val is not None else None
        if val == sid:
            return sh
    return None


def ensure_slot_textbox(slide, slot):
    slot = (slot or "").lower()

    def _get(ph_name):
        return getattr(PP_PLACEHOLDER, ph_name, None)

    TITLE = _get("TITLE")
    CENTER_TITLE = _get("CENTER_TITLE")
    SUBTITLE = _get("SUBTITLE")
    BODY = _get("BODY")
    CONTENT = _get("CONTENT")
    OBJECT = _get("OBJECT")

    title_types = {t for t in (TITLE, CENTER_TITLE, SUBTITLE) if t is not None}
    body_types  = {t for t in (BODY, CONTENT, OBJECT) if t is not None}

    def find_placeholder(accepted_types):
        for sh in slide.shapes:
            if not getattr(sh, "is_placeholder", False):
                continue
            pf = getattr(sh, "placeholder_format", None)
            if not pf:
                continue
            try:
                if pf.type in accepted_types:
                    return sh
            except Exception:
                pass
        return None

    if slot == "title":
        ph = find_placeholder(title_types)
        if ph:
            return ph

    if slot == "body":
        ph = find_placeholder(body_types)
        if ph:
            return ph

    from pptx.util import Inches
    if slot == "title":
        return slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    if slot == "body":
        return slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    return slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.5))

def _layout_has(layout, want_title=False, want_body=False):
    has_title = has_body = False
    for ph in getattr(layout, "placeholders", []):
        pf = getattr(ph, "placeholder_format", None)
        t = getattr(pf, "type", None) if pf else None
        if t in (getattr(PP_PLACEHOLDER, "TITLE", None),
                 getattr(PP_PLACEHOLDER, "CENTER_TITLE", None),
                 getattr(PP_PLACEHOLDER, "SUBTITLE", None)):
            has_title = True
        if t in (getattr(PP_PLACEHOLDER, "BODY", None),
                 getattr(PP_PLACEHOLDER, "CONTENT", None),
                 getattr(PP_PLACEHOLDER, "OBJECT", None)):
            has_body = True
    return (not want_title or has_title) and (not want_body or has_body)

def _pick_layout_for_slots(prs, anchor_slide, needs_title, needs_body):
    if anchor_slide and _layout_has(anchor_slide.slide_layout, needs_title, needs_body):
        return anchor_slide.slide_layout
    for layout in prs.slide_layouts:
        if _layout_has(layout, needs_title, needs_body):
            return layout
    return anchor_slide.slide_layout if anchor_slide else prs.slide_layouts[-1]

def _collect_needs(edit_items):
    needs = {}
    for tgt, _ in edit_items:
        if not isinstance(tgt, str):
            continue
        m = re.match(r"^(n\d+):slot:(title|body)$", tgt.strip(), flags=re.I)
        if m:
            ref, slot = m.group(1), m.group(2).lower()
            needs.setdefault(ref, {"title": False, "body": False})
            needs[ref][slot] = True
    return needs

def _body_placeholder_bounds(slide):
    """Return (left, top, width, height) for the body/content area if possible, else None."""
    try:
        for shp in slide.shapes:
            phf = getattr(shp, "placeholder_format", None)
            if phf is not None:
                # BODY placeholder is the content region on most layouts
                if str(getattr(phf, "type", "")).endswith("BODY"):
                    return shp.left, shp.top, shp.width, shp.height
    except Exception:
        pass
    return None

def _add_table_from_matrix(slide, matrix):
    """
    Create a table on the slide sized to the matrix (rows x cols) and fill it.
    The table is placed over the body placeholder bounds if available,
    else within 1-inch margins.
    Returns the created table shape.
    """
    if not isinstance(matrix, (list, tuple)) or not matrix or not isinstance(matrix[0], (list, tuple)):
        return None

    rows = len(matrix)
    cols = len(matrix[0])

    # determine placement rectangle
    rect = _body_placeholder_bounds(slide)
    if rect:
        left, top, width, height = rect
    else:
        # safe default margins
        left = Inches(1)
        top = Inches(1.2)
        # try to use slide size when available
        try:
            prs = slide.part.presentation
            width = prs.slide_width - Inches(2)
            height = prs.slide_height - Inches(2.2)
        except Exception:
            width = Inches(8)
            height = Inches(4.5)

    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = tbl_shape.table

    # fill cells
    for r in range(rows):
        for c in range(cols):
            try:
                table.cell(r, c).text = "" if matrix[r][c] is None else str(matrix[r][c])
            except Exception:
                pass

    return tbl_shape


def _resolve_donor_simple(order, slides_by_id, anchor_id, kind):
    """
    kind: 'insert_after' or 'insert_before'
    Rules:
      insert_after(anchor):
        - if anchor is first -> donor = next slide if exists else anchor
        - else               -> donor = anchor
      insert_before(anchor):
        - if anchor is last  -> donor = previous slide if exists else anchor
        - else               -> donor = anchor
    """
    if not order:
        return None
    if anchor_id not in order:
        # anchor not found
        return slides_by_id.get(order[1]) if len(order) > 1 else slides_by_id.get(order[0])

    pos = order.index(anchor_id)
    last_idx = len(order) - 1

    if kind == "insert_after":
        if pos == 0:
            # after first slide
            return slides_by_id.get(order[pos + 1]) if pos + 1 <= last_idx else slides_by_id.get(anchor_id)
        else:
            return slides_by_id.get(anchor_id)

    # insert_before
    if pos == last_idx:
        # before last slide
        return slides_by_id.get(order[pos - 1]) if pos - 1 >= 0 else slides_by_id.get(anchor_id)
    else:
        return slides_by_id.get(anchor_id)

def _set_table_from_matrix(shape, data):
    # data = list[list[Any]]; trims to current table size
    tbl = shape.table
    max_r = len(tbl.rows)
    max_c = len(tbl.columns)
    for r, row_vals in enumerate(data):
        if r >= max_r:
            break
        for c, val in enumerate(row_vals):
            if c >= max_c:
                break
            tbl.cell(r, c).text = ""  # clear
            tbl.cell(r, c).text = "" if val is None else str(val)

@mcp.tool()
async def edit_document(
    file_id: str,
    file_name: str,
    edits: dict,
    headers: dict = None,
    ctx: Context[ServerSession, None] = None
) -> dict:
    """
    Edits a document (docx, xlsx, pptx) using structured operations.

    Args:
        file_id: Unique identifier for the document.
        file_name: Name of the document file.
        edits: Dictionary with:
            - "ops": List of structural changes.
            - "content_edits": List of content updates.

    ## Supported Formats

    ### PPTX (PowerPoint)
    - ops: 
        - ["insert_after", <slide_id>, "nK"]
        - ["insert_after", <slide_id>, "nK", {"layout_like_sid": <slide_id>}]
        - ["insert_before", <slide_id>, "nK"]
        - ["insert_before", <slide_id>, "nK", {"layout_like_sid": <slide_id>}]
        - ["delete_slide", slide_id]
    - content_edits:
        - Edit a text shape
            ["sid:<slide_id>/shid:<shape_id>", text_or_list]
        - Edit a table
            ["sid:<slide_id>/shid:<shape_id>", [[row1_col1, row1_col2], [row2_col1, row2_col2], ...]]
        - Edit title or body or table of a newly inserted slide
            ["nK:slot:title", text_or_list]
            ["nK:slot:body", text_or_list]
            ["nK:slot:table", [[row1_col1, row1_col2], [row2_col1, row2_col2], ...]]

    ### DOCX (Word)
    - ops:
        - ["insert_after", para_xml_id, "nK"]
        - ["insert_before", para_xml_id, "nK"]
        - ["delete_paragraph", para_xml_id]
    - content_edits:
        - ["pid:<para_xml_id>", text_or_list]
        - ["tid:<table_xml_id>/cid:<cell_xml_id>", text]
        - ["nK", text_or_list]

    ### XLSX (Excel)
    - ops:
        - ["insert_row", "sheet_name", row_idx]
        - ["delete_row", "sheet_name", row_idx]
        - ["insert_column", "sheet_name", col_idx]
        - ["delete_column", "sheet_name", col_idx]
    - content_edits:
        - ["<ref>", value]

    ## Notes
    - Always call `full_context_document()` first to get IDs.
    - Use cell refs like "A1", "B5".
    - Formatting is preserved.
    - Returns a download link to the edited file.
    """
    temp_folder = f"/app/temp/{uuid.uuid4()}"
    os.makedirs(temp_folder, exist_ok=True)
    user_token = TOKEN
    if headers:
        auth_header = headers.get("authorization")
        if auth_header:
            user_token = auth_header
            logging.info("✅ Using authorization from MCPO forwarded headers")
        else:
            logging.warning("⚠️ Forwarded headers present but no authorization found")
    else:
        logging.info("ℹ️ No forwarded headers, using admin TOKEN fallback")
    try:
        user_file = download_file(file_id, token=user_token)
        if isinstance(user_file, dict) and "error" in user_file:
            return json.dumps(user_file, indent=4, ensure_ascii=False)

        file_extension = os.path.splitext(file_name)[1].lower()
        file_type = file_extension.lstrip('.')

        edited_path = None
        response = None

        if file_type == "docx":
            try:
                doc = Document(user_file)
                
                para_by_xml_id = {}
                table_by_xml_id = {}
                cell_by_xml_id = {}
                
                para_id_counter = 1
                for para in doc.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    para_by_xml_id[para_id_counter] = para
                    para_id_counter += 1
          
                for table in doc.tables:
                    table_xml_id = id(table._element)
                    table_by_xml_id[table_xml_id] = table
                    for row in table.rows:
                        for cell in row.cells:
                            cell_xml_id = id(cell._element)
                            cell_by_xml_id[cell_xml_id] = cell
          
                if isinstance(edits, dict):
                    ops = edits.get("ops", []) or []
                    edit_items = edits.get("content_edits", []) or []
                else:
                    ops = []
                    edit_items = edits
          
                new_refs = {}
                
                for op in ops:
                    if not isinstance(op, (list, tuple)) or not op:
                        continue
                    kind = op[0]
                    
                    if kind == "insert_after" and len(op) >= 3:
                        anchor_xml_id = int(op[1])
                        new_ref = op[2]
                        
                        anchor_para = para_by_xml_id.get(anchor_xml_id)
                        if anchor_para:
                            para_index = doc.paragraphs.index(anchor_para)
          	       	      
                            new_para = doc.add_paragraph()
          	       	      
                            anchor_element = anchor_para._element
                            parent = anchor_element.getparent()
                            parent.insert(parent.index(anchor_element) + 1, new_para._element)
          	       	      
                            new_para.style = anchor_para.style
      	      				
                            new_xml_id = id(new_para._element)
                            new_refs[new_ref] = new_xml_id
                            para_by_xml_id[new_xml_id] = new_para
      	      		
                    elif kind == "insert_before" and len(op) >= 3:
                        anchor_xml_id = int(op[1])
                        new_ref = op[2]
                        
                        anchor_para = para_by_xml_id.get(anchor_xml_id)
                        if anchor_para:
                            new_para = doc.add_paragraph()
          	       	      
                            anchor_element = anchor_para._element
                            parent = anchor_element.getparent()
                            parent.insert(parent.index(anchor_element), new_para._element)
          	       	      
                            new_para.style = anchor_para.style
      	      				
                            new_xml_id = id(new_para._element)
                            new_refs[new_ref] = new_xml_id
                            para_by_xml_id[new_xml_id] = new_para
      	      		
                    elif kind == "delete_paragraph" and len(op) >= 2:
                        para_xml_id = int(op[1])
                        para = para_by_xml_id.get(para_xml_id)
                        if para:
                            parent = para._element.getparent()
                            parent.remove(para._element)
                            para_by_xml_id.pop(para_xml_id, None)
                
                for target, new_text in edit_items:
                    if not isinstance(target, str):
                        continue
      	      		
                    t = target.strip()
                    
                    m = re.match(r"^pid:(\d+)$", t, flags=re.I)
                    if m:
                        para_xml_id = int(m.group(1))
                        para = para_by_xml_id.get(para_xml_id)
                        if para:
                            _apply_text_to_paragraph(para, new_text)
                        continue
      	      		
                    m = re.match(r"^tid:(\d+)/cid:(\d+)$", t, flags=re.I)
                    if m:
                        table_xml_id = int(m.group(1))
                        cell_xml_id = int(m.group(2))
                        cell = cell_by_xml_id.get(cell_xml_id)
                        if cell:
                            for para in cell.paragraphs:
                                for _ in range(len(para.runs)):
                                    para._element.remove(para.runs[0]._element)
          	       	      
                            if cell.paragraphs:
                                first_para = cell.paragraphs[0]
                                first_para.add_run(str(new_text))
                        continue
      	      		
                    m = re.match(r"^n(\d+)$", t, flags=re.I)
                    if m:
                        new_ref = t
                        para_xml_id = new_refs.get(new_ref)
                        if para_xml_id:
                            para = para_by_xml_id.get(para_xml_id)
                            if para:
                                _apply_text_to_paragraph(para, new_text)
                        continue
          
                edited_path = os.path.join(
                    temp_folder, f"{os.path.splitext(file_name)[0]}_edited.docx"
                )
                doc.save(edited_path)
                response = upload_file(
                    file_path=edited_path,
                    filename=f"{os.path.splitext(file_name)[0]}_edited",
                    file_type="docx", 
                    token=user_token
                )
            except Exception as e:
                raise Exception(f"Error during DOCX editing: {e}")

        elif file_type == "xlsx":
            try:
                wb = load_workbook(user_file)
                ws = wb.active

                edit_items = edits.get("content_edits", []) if isinstance(edits, dict) and "content_edits" in edits else edits
          
                for index, new_text in edit_items:
                    try:
                        if isinstance(index, str) and re.match(r"^[A-Z]+[0-9]+$", index.strip().upper()):
                            cell_ref = index.strip().upper()
                        elif isinstance(index, int):
                            cell_ref = f"A{index+1}"
                        else:
                            cell_ref = "A1"

                        cell = ws[cell_ref]
                        cell.value = new_text
                    except Exception:
                        fallback_cell = ws["A1"]
                        fallback_cell.value = new_text

                edited_path = os.path.join(
                    temp_folder, f"{os.path.splitext(file_name)[0]}_edited.xlsx"
                )
                wb.save(edited_path)
                response = upload_file(
                    file_path=edited_path,
                    filename=f"{os.path.splitext(file_name)[0]}_edited",
                    file_type="xlsx", 
                    token=user_token
                )
            except Exception as e:
                raise Exception(f"Error during XLSX editing: {e}")

        elif file_type == "pptx":
            try:
                prs = Presentation(user_file)
                if isinstance(edits, dict):
                    ops = edits.get("ops", []) or []
                    edit_items = edits.get("content_edits", []) or []
                else:
                    ops = []
                    edit_items = edits
                new_ref_needs = _collect_needs(edit_items)
                order = [int(s.slide_id) for s in prs.slides]
                slides_by_id = {int(s.slide_id): s for s in prs.slides}
                new_refs = {}
                
                for op in ops:
                    if not isinstance(op, (list, tuple)) or not op:
                        continue
                    kind = op[0]

                    if kind == "insert_after" and len(op) >= 3:
                        anchor_id = int(op[1])
                        new_ref = op[2]
                        if anchor_id in order:
                            like_sid = None
                            if len(op) >= 4 and isinstance(op[3], dict):
                                like_sid = op[3].get("layout_like_sid")
                            needs = new_ref_needs.get(new_ref, {"title": False, "body": False})
                            # style_donor selection
                            if like_sid and like_sid in slides_by_id:
                                style_donor = slides_by_id[like_sid]
                            else:
                                style_donor = _resolve_donor_simple(order, slides_by_id, anchor_id, kind)
                            layout = _pick_layout_for_slots(prs, style_donor, needs["title"], needs["body"]) if style_donor else prs.slide_layouts[0]
                            new_slide = prs.slides.add_slide(layout)
                            new_sid = int(new_slide.slide_id)

                            sldIdLst = prs.slides._sldIdLst
                            new_sldId = sldIdLst[-1]
                            try:
                                anchor_pos = order.index(anchor_id)
                                sldIdLst.remove(new_sldId)
                                sldIdLst.insert(anchor_pos + 1, new_sldId)
                            except Exception:
                                pass 
                            order.insert(order.index(anchor_id) + 1, new_sid)
                            slides_by_id[new_sid] = new_slide
                            new_refs[new_ref] = new_sid


                    elif kind == "insert_before" and len(op) >= 3:
                        anchor_id = int(op[1])
                        new_ref = op[2]
                        if anchor_id in order:
                            like_sid = None
                            if len(op) >= 4 and isinstance(op[3], dict):
                                like_sid = op[3].get("layout_like_sid")

                            needs = new_ref_needs.get(new_ref, {"title": False, "body": False})

                            # style_donor selection
                            if like_sid and like_sid in slides_by_id:
                                style_donor = slides_by_id[like_sid]
                            else:
                                style_donor = _resolve_donor_simple(order, slides_by_id, anchor_id, kind)

                            layout = _pick_layout_for_slots(prs, style_donor, needs["title"], needs["body"]) if style_donor else prs.slide_layouts[0]
                            new_slide = prs.slides.add_slide(layout)
                            new_sid = int(new_slide.slide_id)

                            sldIdLst = prs.slides._sldIdLst
                            new_sldId = sldIdLst[-1]
                            try:
                                anchor_pos = order.index(anchor_id)
                                sldIdLst.remove(new_sldId)
                                sldIdLst.insert(anchor_pos, new_sldId)
                            except Exception:
                                pass

                            order.insert(order.index(anchor_id), new_sid)
                            slides_by_id[new_sid] = new_slide
                            new_refs[new_ref] = new_sid

                    elif kind == "delete_slide" and len(op) >= 2:
                        sid = int(op[1])
                        if sid in order:
                            i = order.index(sid)
                            sldIdLst = prs.slides._sldIdLst
                            rId = sldIdLst[i].rId
                            prs.part.drop_rel(rId)
                            del sldIdLst[i]
                            order.pop(i)
                            slides_by_id.pop(sid, None)   		 
      		
                for target, new_text in edit_items:
                    if not isinstance(target, str):
                        continue
                    t = target.strip()

                    # >>> ADD: table edit 
                    # target format: sid:<sid>/shid:<shid>  with new_text like [[row1...],[row2...],...]
                    m = re.match(r"^sid:(\d+)/shid:(\d+)$", t, flags=re.I)
                    if m:
                        slide_id = int(m.group(1))
                        shape_id = int(m.group(2))
                        slide = slides_by_id.get(slide_id)
                        if slide:
                            shape = shape_by_id(slide, shape_id)
                            if shape and getattr(shape, "has_table", False) and isinstance(new_text, (list, tuple)) and new_text and isinstance(new_text[0], (list, tuple)):
                                _set_table_from_matrix(shape, new_text)
                                continue
                            
                    # <<< END ADD

                    m = re.match(r"^sid:(\d+)/shid:(\d+)$", t, flags=re.I)
                    if m:
                        slide_id = int(m.group(1))
                        shape_id = int(m.group(2))
                        slide = slides_by_id.get(slide_id)
                        if not slide:
                            continue
                        shape = shape_by_id(slide, shape_id)
                        if not shape:
                            continue
                        _set_text_with_runs(shape, new_text)
                        continue
                    m = re.match(r"^(n\d+):slot:(title|body)$", t, flags=re.I)
                    if m:
                        ref = m.group(1)
                        slot = m.group(2).lower()
                        sid = new_refs.get(ref)
                        if not sid:
                            continue
                        slide = slides_by_id.get(sid)
                        if not slide:
                            continue
                        shape = ensure_slot_textbox(slide, slot)
                        tf = getattr(shape, "text_frame", None)
                        if tf is None:
                            continue
                        if isinstance(new_text, list):
                            tf.clear()
                            tf.text = str(new_text[0]) if new_text else ""
                            for line in new_text[1:]:
                                p = tf.add_paragraph()
                                p.text = str(line)
                                try:
                                    p.level = getattr(tf.paragraphs[0], "level", 0)
                                except Exception:
                                    pass
                        else:
                            tf.text = str(new_text)
                        continue

                    # nK:table  (create a new table on a newly inserted slide nK)
                    m = re.match(r"^(n\d+):slot:table$", t, flags=re.I)
                    if m:
                        ref = m.group(1)
                        sid = new_refs.get(ref)
                        if not sid:
                            continue
                        slide = slides_by_id.get(sid)
                        if not slide:
                            continue
                        # new_text must be a 2D list
                        if isinstance(new_text, (list, tuple)) and new_text and isinstance(new_text[0], (list, tuple)):
                            _add_table_from_matrix(slide, new_text)
                        continue
 

                edited_path = os.path.join(
                    temp_folder, f"{os.path.splitext(file_name)[0]}_edited.pptx"
                )
                prs.save(edited_path)
                response = upload_file(
                    file_path=edited_path,
                    filename=f"{os.path.splitext(file_name)[0]}_edited",
                    file_type="pptx", 
                    token=user_token
                )
            except Exception as e:
                raise Exception(f"Error during PPTX editing: {e}")

        else:
            raise Exception(f"File type not supported: {file_type}")

        shutil.rmtree(temp_folder, ignore_errors=True)

        return response

    except Exception as e:
        shutil.rmtree(temp_folder, ignore_errors=True)
        return json.dumps(
            {"error": {"message": str(e)}},
            indent=4,
            ensure_ascii=False
        )
    
def _get_pptx_namespaces():
    """Returns XML namespaces for PowerPoint"""
    return {
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p15': 'http://schemas.microsoft.com/office/powerpoint/2012/main',
        'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main'
    }

def _add_native_pptx_comment_zip(pptx_path, slide_num, comment_text, author_id, x=100, y=100):
    """
    Add a native PowerPoint comment by directly manipulating the ZIP file.
        Args:
        pptx_path: Path to the PPTX file
        slide_num: Slide number (1-based)
        comment_text: Comment text
        author_id: Author ID
        x: X position in EMU (not pixels!)
        y: Y position in EMU (not pixels!)
    """
    namespaces = _get_pptx_namespaces()
    
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            zf.extractall(temp_path)
        
        authors_file = temp_path / 'ppt' / 'commentAuthors.xml'
        if authors_file.exists():
            root = etree.parse(str(authors_file)).getroot()
            found = False
            for author in root.findall('.//p:cmAuthor', namespaces):
                if author.get('name') == 'AI Reviewer':
                    author_id = int(author.get('id'))
                    found = True
                    break
            
            if not found:
                existing_ids = [int(a.get('id')) for a in root.findall('.//p:cmAuthor', namespaces)]
                author_id = max(existing_ids) + 1 if existing_ids else 0
                author = etree.SubElement(root, f'{{{namespaces["p"]}}}cmAuthor')
                author.set('id', str(author_id))
                author.set('name', 'AI Reviewer')
                author.set('initials', 'AI')
                author.set('lastIdx', '1')
                author.set('clrIdx', str(author_id % 8))
        else:
            authors_file.parent.mkdir(parents=True, exist_ok=True)
            root = etree.Element(
                f'{{{namespaces["p"]}}}cmAuthorLst',
                nsmap={k: v for k, v in namespaces.items() if k in ['p']}
            )
            author = etree.SubElement(root, f'{{{namespaces["p"]}}}cmAuthor')
            author.set('id', str(author_id))
            author.set('name', 'AI Reviewer')
            author.set('initials', 'AI')
            author.set('lastIdx', '1')
            author.set('clrIdx', '0')
            
            rels_file = temp_path / 'ppt' / '_rels' / 'presentation.xml.rels'
            if rels_file.exists():
                rels_root = etree.parse(str(rels_file)).getroot()
                existing_ids = [int(rel.get('Id')[3:]) for rel in rels_root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship')]
                next_rid = max(existing_ids) + 1 if existing_ids else 1
                
                rel = etree.SubElement(rels_root, '{http://schemas.openxmlformats.org/package/2006/relationships}Relationship')
                rel.set('Id', f'rId{next_rid}')
                rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/commentAuthors')
                rel.set('Target', 'commentAuthors.xml')
                
                with open(rels_file, 'wb') as f:
                    f.write(etree.tostring(rels_root, xml_declaration=True, encoding='UTF-8', pretty_print=True))
        
        with open(authors_file, 'wb') as f:
            f.write(etree.tostring(root, xml_declaration=True, encoding='UTF-8', pretty_print=True))
        
        comments_dir = temp_path / 'ppt' / 'comments'
        comments_dir.mkdir(parents=True, exist_ok=True)
        comment_file = comments_dir / f'comment{slide_num}.xml'
        
        if comment_file.exists():
            comments_root = etree.parse(str(comment_file)).getroot()
        else:
            comments_root = etree.Element(
                f'{{{namespaces["p"]}}}cmLst',
                nsmap={k: v for k, v in namespaces.items() if k in ['p']}
            )
            
            slide_rels_file = temp_path / 'ppt' / 'slides' / '_rels' / f'slide{slide_num}.xml.rels'
            if slide_rels_file.exists():
                slide_rels_root = etree.parse(str(slide_rels_file)).getroot()
            else:
                slide_rels_file.parent.mkdir(parents=True, exist_ok=True)
                slide_rels_root = etree.Element(
                    '{http://schemas.openxmlformats.org/package/2006/relationships}Relationships'
                )
            
            existing_ids = [int(rel.get('Id')[3:]) for rel in slide_rels_root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship')]
            next_rid = max(existing_ids) + 1 if existing_ids else 1
            
            rel = etree.SubElement(slide_rels_root, '{http://schemas.openxmlformats.org/package/2006/relationships}Relationship')
            rel.set('Id', f'rId{next_rid}')
            rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments')
            rel.set('Target', f'../comments/comment{slide_num}.xml')
            
            with open(slide_rels_file, 'wb') as f:
                f.write(etree.tostring(slide_rels_root, xml_declaration=True, encoding='UTF-8', pretty_print=True))
        
        existing_ids = [int(c.get('idx')) for c in comments_root.findall('.//p:cm', namespaces)]
        next_id = max(existing_ids) + 1 if existing_ids else 1
        
        comment = etree.SubElement(comments_root, f'{{{namespaces["p"]}}}cm')
        comment.set('authorId', str(author_id))
        comment.set('dt', datetime.datetime.now().isoformat())
        comment.set('idx', str(next_id))
        
        pos = etree.SubElement(comment, f'{{{namespaces["p"]}}}pos')
        pos.set('x', str(int(x)))
        pos.set('y', str(int(y)))
        
        text_elem = etree.SubElement(comment, f'{{{namespaces["p"]}}}text')
        text_elem.text = comment_text
        
        with open(comment_file, 'wb') as f:
            f.write(etree.tostring(comments_root, xml_declaration=True, encoding='UTF-8', pretty_print=True))
        
        content_types_file = temp_path / '[Content_Types].xml'
        if content_types_file.exists():
            ct_root = etree.parse(str(content_types_file)).getroot()
            ns = {'ct': 'http://schemas.openxmlformats.org/package/2006/content-types'}
            
            has_authors = False
            has_comments = False
            
            for override in ct_root.findall('.//ct:Override', ns):
                if override.get('PartName') == '/ppt/commentAuthors.xml':
                    has_authors = True
                if override.get('PartName') == f'/ppt/comments/comment{slide_num}.xml':
                    has_comments = True
            
            if not has_authors:
                override = etree.SubElement(ct_root, '{http://schemas.openxmlformats.org/package/2006/content-types}Override')
                override.set('PartName', '/ppt/commentAuthors.xml')
                override.set('ContentType', 'application/vnd.openxmlformats-officedocument.presentationml.commentAuthors+xml')
            
            if not has_comments:
                override = etree.SubElement(ct_root, '{http://schemas.openxmlformats.org/package/2006/content-types}Override')
                override.set('PartName', f'/ppt/comments/comment{slide_num}.xml')
                override.set('ContentType', 'application/vnd.openxmlformats-officedocument.presentationml.comments+xml')
            
            with open(content_types_file, 'wb') as f:
                f.write(etree.tostring(ct_root, xml_declaration=True, encoding='UTF-8', pretty_print=True))
        
        with zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for file_path in temp_path.rglob('*'):
                if file_path.is_file():
                    arcname = str(file_path.relative_to(temp_path))
                    zf.write(file_path, arcname)
        
        log.debug(f"Native comment added to slide {slide_num} with idx={next_id}")

@mcp.tool(
    name="review_document",
    title="Review and comment on various document types",
    description="Review an existing document of various types (docx, xlsx, pptx), perform corrections and add comments."
)
async def review_document(
    file_id: str,
    file_name: str,
    review_comments: list[tuple[int | str, str]],
    headers: dict = None,
    ctx: Context[ServerSession, None] = None
) -> dict:
    """
    Generic document review function that works with different document types.
    File type is automatically detected from the file extension.
    Returns a markdown hyperlink for downloading the reviewed document.
    
    For Excel files (.xlsx):
    - The index must be a cell reference (e.g. "A1", "B3", "C10")
    - These correspond to the "index" key returned by the full_context_document() function
    - Never use integer values for Excel cells
    
    For Word files (.docx):
    - The index should be a paragraph ID in the format "pid:<para_xml_id>"
    - These correspond to the "id_key" field returned by the full_context_document() function
    
    For PowerPoint files (.pptx):
    - The index should be a slide ID in the format "sid:<slide_id>"
    - These correspond to the "id_key" field returned by the full_context_document() function
    """
    temp_folder = f"/app/temp/{uuid.uuid4()}"
    os.makedirs(temp_folder, exist_ok=True)
    user_token = TOKEN
    if headers:
        auth_header = headers.get("authorization")
        if auth_header:
            user_token = auth_header
            logging.info("✅ Using authorization from MCPO forwarded headers")
        else:
            logging.warning("⚠️ Forwarded headers present but no authorization found")
    else:
        logging.info("ℹ️ No forwarded headers, using admin TOKEN fallback")
    try:
        user_file = download_file(file_id, token=user_token)
        if isinstance(user_file, dict) and "error" in user_file:
            return json.dumps(user_file, indent=4, ensure_ascii=False)

        file_extension = os.path.splitext(file_name)[1].lower()
        file_type = file_extension.lstrip('.')

        reviewed_path = None
        response = None

        if file_type == "docx":
            try:
                doc = Document(user_file)
                paragraphs = list(doc.paragraphs)
                para_by_xml_id = {}
                para_id_counter = 1
                
                for para in doc.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    para_by_xml_id[para_id_counter] = para
                    para_id_counter += 1

                for index, comment_text in review_comments:
                    if isinstance(index, int) and 0 <= index < len(paragraphs):
                        para = paragraphs[index]
                        if para.runs:
                            try:
                                doc.add_comment(
                                    runs=[para.runs[0]],
                                    text=comment_text,
                                    author="AI Reviewer",
                                    initials="AI"
                                )
                            except Exception:
                                para.add_run(f"  [AI Comment: {comment_text}]")
                    elif isinstance(index, str) and index.startswith("pid:"):
                        try:
                            para_xml_id = int(index.split(":")[1])
                            para = para_by_xml_id.get(para_xml_id)
                            if para and para.runs:
                                try:
                                    doc.add_comment(
                                        runs=[para.runs[0]],
                                        text=comment_text,
                                        author="AI Reviewer",
                                        initials="AI"
                                    )
                                except Exception:
                                    para.add_run(f"  [AI Comment: {comment_text}]")
                        except Exception:
                            if isinstance(index, int) and 0 <= index < len(paragraphs):
                                para = paragraphs[index]
                                if para.runs:
                                    try:
                                        doc.add_comment(
                                            runs=[para.runs[0]],
                                            text=comment_text,
                                            author="AI Reviewer",
                                            initials="AI"
                                        )
                                    except Exception:
                                        para.add_run(f"  [AI Comment: {comment_text}]")
                reviewed_path = os.path.join(
                    temp_folder, f"{os.path.splitext(file_name)[0]}_reviewed.docx"
                )
                doc.save(reviewed_path)
                response = upload_file(
                    file_path=reviewed_path,
                    filename=f"{os.path.splitext(file_name)[0]}_reviewed",
                    file_type="docx", 
                    token=user_token
                )
            except Exception as e:
                raise Exception(f"Error during DOCX revision: {e}")

        elif file_type == "xlsx":
            try:
                wb = load_workbook(user_file)
                ws = wb.active

                for index, comment_text in review_comments:
                    try:
                        if isinstance(index, str) and re.match(r"^[A-Z]+[0-9]+$", index.strip().upper()):
                            cell_ref = index.strip().upper()
                        elif isinstance(index, int):
                            cell_ref = f"A{index+1}"
                        else:
                            cell_ref = "A1"

                        cell = ws[cell_ref]
                        add_auto_sized_review_comment(cell, comment_text, author="AI Reviewer")

                    except Exception:
                        fallback_cell = ws["A1"]
                        add_auto_sized_review_comment(fallback_cell, comment_text, author="AI Reviewer")

                reviewed_path = os.path.join(
                    temp_folder, f"{os.path.splitext(file_name)[0]}_reviewed.xlsx"
                )
                wb.save(reviewed_path)
                response = upload_file(
                    file_path=reviewed_path,
                    filename=f"{os.path.splitext(file_name)[0]}_reviewed",
                    file_type="xlsx", 
                    token=user_token
                )
            except Exception as e:
                raise Exception(f"Error: {e}")

        elif file_type == "pptx":
            try:
                temp_pptx = os.path.join(temp_folder, "temp_input.pptx")
                with open(temp_pptx, 'wb') as f:
                    f.write(user_file.read())
                
                prs = Presentation(temp_pptx)
                slides_by_id = {int(s.slide_id): s for s in prs.slides}
                
                comments_by_slide = {}
                
                for index, comment_text in review_comments:
                    slide_num = None
                    slide_id = None
                    
                    if isinstance(index, int) and 0 <= index < len(prs.slides):
                        slide_num = index + 1
                        slide_id = list(slides_by_id.keys())[index]
                    elif isinstance(index, str):
                        if index.startswith("sid:") and "/shid:" in index:
                            try:
                                slide_id = int(index.split("/")[0].replace("sid:", ""))
                                if slide_id in slides_by_id:
                                    slide_num = list(slides_by_id.keys()).index(slide_id) + 1
                            except Exception as e:
                                log.warning(f"Failed to parse shape ID: {e}")
                        elif index.startswith("sid:"):
                            try:
                                slide_id = int(index.split(":")[1])
                                if slide_id in slides_by_id:
                                    slide_num = list(slides_by_id.keys()).index(slide_id) + 1
                            except Exception as e:
                                log.warning(f"Failed to parse slide ID: {e}")
                    
                    if slide_num and slide_id:
                        if slide_num not in comments_by_slide:
                            comments_by_slide[slide_num] = []
                        
                        shape_info = ""
                        if "/shid:" in str(index):
                            try:
                                shape_id = int(str(index).split("/shid:")[1])
                                shape_info = f"[Shape {shape_id}] "
                            except:
                                pass
                        
                        comments_by_slide[slide_num].append(f"{shape_info}{comment_text}")
                comment_offset = 0              
                for slide_num, comments in comments_by_slide.items():
                    comment_start_x = 5000
                    comment_start_y = 1000
                    comment_spacing_y = 1500
                    
                    for idx, comment in enumerate(comments):
                        try:
                            y_position = comment_start_y + (idx * comment_spacing_y)
                            
                            _add_native_pptx_comment_zip(
                                pptx_path=temp_pptx,
                                slide_num=slide_num,
                                comment_text=f"• {comment}",
                                author_id=0,
                                x=comment_start_x,
                                y=y_position
                            )
                            log.debug(f"Native PowerPoint comment added to slide {slide_num} at position x={comment_start_x}, y={y_position}")
                        except Exception as e:
                            log.warning(f"Failed to add native comment to slide {slide_num}: {e}", exc_info=True)
                            prs_fallback = Presentation(temp_pptx)
                            slide = prs_fallback.slides[slide_num - 1]
                            left = top = Inches(0.2)
                            width = Inches(4)
                            height = Inches(1)
                            textbox = slide.shapes.add_textbox(left, top, width, height)
                            text_frame = textbox.text_frame
                            p = text_frame.add_paragraph()
                            p.text = f"AI Reviewer: {comment}"
                            p.font.size = PptPt(10)
                            prs_fallback.save(temp_pptx)

                reviewed_path = os.path.join(
                    temp_folder, f"{os.path.splitext(file_name)[0]}_reviewed.pptx"
                )
                shutil.copy(temp_pptx, reviewed_path)
                
                response = upload_file(
                    file_path=reviewed_path,
                    filename=f"{os.path.splitext(file_name)[0]}_reviewed",
                    file_type="pptx", 
                    token=user_token
                )
            except Exception as e:
                raise Exception(f"Error when revising PPTX: {e}")

        else:
            raise Exception(f"File type not supported : {file_type}")

        shutil.rmtree(temp_folder, ignore_errors=True)

        return response

    except Exception as e:
        shutil.rmtree(temp_folder, ignore_errors=True)
        return json.dumps(
            {"error": {"message": str(e)}},
            indent=4,
            ensure_ascii=False
        )

@mcp.tool()
async def create_file(data: dict, persistent: bool = PERSISTENT_FILES) -> dict:
    """ "{"data": {"format":"pdf","filename":"report.pdf","content":[{"type":"title","text":"..."},{"type":"paragraph","text":"..."}],"title":"..."}}
"{"data": {"format":"docx","filename":"doc.docx","content":[{"type":"title","text":"..."},{"type":"list","items":[...]}],"title":"..."}}"
"{"data": {"format":"pptx","filename":"slides.pptx","slides_data":[{"title":"...","content":[...],"image_query":"...","image_position":"left|right|top|bottom","image_size":"small|medium|large"}],"title":"..."}}"
"{"data": {"format":"xlsx","filename":"data.xlsx","content":[["Header1","Header2"],["Val1","Val2"]],"title":"..."}}"
"{"data": {"format":"csv","filename":"data.csv","content":[[...]]}}"
"{"data": {"format":"txt|xml|py|etc","filename":"file.ext","content":"string"}}" """
    log.debug("Creating file via tool")
    folder_path = _generate_unique_folder()
    format_type = (data.get("format") or "").lower()
    filename = data.get("filename")
    content = data.get("content")
    title = data.get("title")

    if format_type == "pdf":
        result = _create_pdf(content if isinstance(content, list) else [str(content or "")], filename, folder_path=folder_path)
    elif format_type == "pptx":
        result = _create_presentation(data.get("slides_data", []), filename, folder_path=folder_path, title=title)
    elif format_type == "docx":
        result = _create_word(content if content is not None else [], filename, folder_path=folder_path, title=title)
    elif format_type == "xlsx":
        result = _create_excel(content if content is not None else [], filename, folder_path=folder_path, title=title)
    elif format_type == "csv":
        result = _create_csv(content if content is not None else [], filename, folder_path=folder_path)
    else:
        use_filename = filename or f"export.{format_type or 'txt'}"
        result = _create_raw_file(content if content is not None else "", use_filename, folder_path=folder_path)

    if not persistent:
        _cleanup_files(folder_path, FILES_DELAY)

    return {"url": result["url"]}

@mcp.tool()
async def generate_and_archive(files_data: list[dict], archive_format: str = "zip", archive_name: str = None, persistent: bool = PERSISTENT_FILES) -> dict:
    """files_data=[{"format":"pdf","filename":"report.pdf","content":[{"type":"title","text":"..."},{"type":"paragraph","text":"..."}],"title":"..."},{"format":"docx","filename":"doc.docx","content":[{"type":"title","text":"..."},{"type":"list","items":[...]}],"title":"..."},{"format":"pptx","filename":"slides.pptx","slides_data":[{"title":"...","content":[...],"image_query":"...","image_position":"left|right|top|bottom","image_size":"small|medium|large"}],"title":"..."},{"format":"xlsx","filename":"data.xlsx","content":[["Header1","Header2"],["Val1","Val2"]],"title":"..."},{"format":"csv","filename":"data.csv","content":[[...]]},{"format":"txt|xml|py|etc","filename":"file.ext","content":"string"}]"""
    log.debug("Generating archive via tool")
    folder_path = _generate_unique_folder()
    generated_paths: list[str] = []

    for file_info in files_data or []:
        fmt = (file_info.get("format") or "").lower()
        fname = file_info.get("filename")
        content = file_info.get("content")
        title = file_info.get("title")

        try:
            if fmt == "pdf":
                res = _create_pdf(content if isinstance(content, list) else [str(content or "")], fname, folder_path=folder_path)
            elif fmt == "pptx":
                res = _create_presentation(file_info.get("slides_data", []), fname, folder_path=folder_path, title=title)
            elif fmt == "docx":
                res = _create_word(content if content is not None else [], fname, folder_path=folder_path, title=title)
            elif fmt == "xlsx":
                res = _create_excel(content if content is not None else [], fname, folder_path=folder_path, title=title)
            elif fmt == "csv":
                res = _create_csv(content if content is not None else [], fname, folder_path=folder_path)
            else:
                use_fname = fname or f"export.{fmt or 'txt'}"
                res = _create_raw_file(content if content is not None else "", use_fname, folder_path=folder_path)
        except Exception as e:
            log.error(f"Error generating file {fname or '<no name>'}: {e}", exc_info=True)
            raise

        generated_paths.append(res["path"])

    timestamp = datetime.datetime.now().strftime('%Y%m%d_%H%M%S')
    archive_basename = f"{archive_name or 'archive'}_{timestamp}"
    archive_filename = f"{archive_basename}.zip" if archive_format.lower() not in ("7z", "tar.gz") else f"{archive_basename}.{archive_format}"
    archive_path = os.path.join(folder_path, archive_filename)

    if archive_format.lower() == "7z":
        with py7zr.SevenZipFile(archive_path, mode='w') as archive:
            for p in generated_paths:
                archive.write(p, os.path.relpath(p, folder_path))
    elif archive_format.lower() == "tar.gz":
        with tarfile.open(archive_path, "w:gz") as tar:
            for p in generated_paths:
                tar.add(p, arcname=os.path.relpath(p, folder_path))
    else:
        with zipfile.ZipFile(archive_path, 'w') as zipf:
            for p in generated_paths:
                zipf.write(p, os.path.relpath(p, folder_path))

    if not persistent:
        _cleanup_files(folder_path, FILES_DELAY)

    return {"url": _public_url(folder_path, archive_filename)}

if __name__ == "__main__":
    log.info(f"Starting  File Export builtin MCPO v{SCRIPT_VERSION}")
    mcp.run()
