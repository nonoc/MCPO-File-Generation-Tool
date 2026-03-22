"""
Export dispatcher - main entry point for file generation.
"""

import os
import logging
import json
import uuid
import datetime
import time
import base64
import shutil
import tarfile
import zipfile
import py7zr
import requests
from requests import get, post
from requests.auth import HTTPBasicAuth
import threading
from io import BytesIO
from pptx.util import Inches
from pptx.enum.shapes import PP_PLACEHOLDER
from lxml import etree

# Import from new modules (try relative first, then absolute)
try:
    from .shared.constants import (
        EXPORT_DIR,
        BASE_URL,
        FILES_DELAY,
        PERSISTENT_FILES,
        TABLE_SEPARATOR_RE,
        IMAGE_SOURCE_UNSPLASH,
        IMAGE_SOURCE_PEXELS,
        IMAGE_SOURCE_LOCAL_SD,
        DOCS_TEMPLATE_PATH,
        PPTX_TEMPLATE,
        DOCX_TEMPLATE,
        XLSX_TEMPLATE,
        PPTX_TEMPLATE_PATH,
        DOCX_TEMPLATE_PATH,
        XLSX_TEMPLATE_PATH,
    )
    from .shared.utils import (
        _public_url,
        _generate_unique_folder,
        _generate_filename,
        _cleanup_files,
        render_text_with_emojis,
        _extract_paragraph_style_info,
        _extract_cell_style_info,
        dynamic_font_size,
        _resolve_log_level,
    )
    from .shared.images import search_image
    from .processing.block_utils import _convert_markdown_to_structured, flatten_structured_blocks
    from .processing.normalization import normalize_content_for_export, _StructuredContentRenderer
    from .exporters import (
        _create_word,
        _create_pdf,
        _create_excel,
        _create_presentation,
        add_auto_sized_review_comment,
    )
except ImportError:
    from shared.constants import (
        EXPORT_DIR,
        BASE_URL,
        FILES_DELAY,
        PERSISTENT_FILES,
        TABLE_SEPARATOR_RE,
        IMAGE_SOURCE_UNSPLASH,
        IMAGE_SOURCE_PEXELS,
        IMAGE_SOURCE_LOCAL_SD,
        DOCS_TEMPLATE_PATH,
        PPTX_TEMPLATE,
        DOCX_TEMPLATE,
        XLSX_TEMPLATE,
        PPTX_TEMPLATE_PATH,
        DOCX_TEMPLATE_PATH,
        XLSX_TEMPLATE_PATH,
    )
    from shared.utils import (
        _public_url,
        _generate_unique_folder,
        _generate_filename,
        _cleanup_files,
        render_text_with_emojis,
        _extract_paragraph_style_info,
        _extract_cell_style_info,
        dynamic_font_size,
        _resolve_log_level,
    )
    from shared.images import search_image
    from processing.block_utils import _convert_markdown_to_structured, flatten_structured_blocks
    from processing.normalization import normalize_content_for_export, _StructuredContentRenderer
    from exporters import (
        _create_word,
        _create_pdf,
        _create_excel,
        _create_presentation,
        add_auto_sized_review_comment,
    )

log = logging.getLogger(__name__)

# File format types
FORMAT_PDF = "pdf"
FORMAT_DOCX = "docx"
FORMAT_XLSX = "xlsx"
FORMAT_PPTX = "pptx"
FORMAT_CSV = "csv"
FORMAT_TXT = "txt"
FORMAT_XML = "xml"


def _create_csv(data, filename, folder_path=None):
    """Create CSV file."""
    log.debug("Creating CSV file")
    if folder_path is None:
        folder_path = _generate_unique_folder()

    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "csv")

    with open(filepath, "w", newline="", encoding="utf-8") as f:
        if isinstance(data, list):
            import csv
            csv.writer(f).writerows(data)
        else:
            import csv
            csv.writer(f).writerow([data])

    return {"url": _public_url(folder_path, fname), "path": filepath}


def _create_raw_file(content, filename, folder_path=None):
    """Create raw text file."""
    log.debug("Creating raw file")
    if folder_path is None:
        folder_path = _generate_unique_folder()

    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "txt")

    if fname.lower().endswith(".xml") and isinstance(content, str) and not content.strip().startswith("<?xml"):
        content = f'<?xml version="1.0" encoding="UTF-8"?>\n{content}'

    with open(filepath, "w", encoding="utf-8") as f:
        f.write(content or "")

    return {"url": _public_url(folder_path, fname), "path": filepath}


def _extract_document_title(content, root_title):
    """Extract document title using block-first strategy."""
    content_list = []
    if isinstance(content, dict):
        content_list = [content]
    elif isinstance(content, list):
        content_list = content
    
    has_title_block = False
    title_text_from_block = None
    
    for item in content_list:
        if isinstance(item, dict):
            item_type = (item.get("type") or "").strip().lower()
            if item_type == "title":
                has_title_block = True
                text = item.get("text") or item.get("title") or ""
                if text:
                    title_text_from_block = str(text).strip()
                    break
        elif isinstance(item, str):
            if item.strip().startswith("# "):
                has_title_block = True
                title_text_from_block = item.strip()[2:].strip()
                break
    
    if has_title_block:
        log.debug(f"Title block found, using block title: '{title_text_from_block}'")
        return title_text_from_block
    
    if root_title:
        log.debug(f"No title block found, using root title: '{root_title}'")
        return root_title.strip()
    
    return None


def _filter_title_blocks_from_content(content, has_root_title=False):
    """
    Filter out title blocks from content when a root-level title is provided.
    
    Args:
        content: Content to filter (string, dict, or list)
        has_root_title: If True, filter out title blocks; otherwise return content unchanged
        
    Returns:
        Filtered content with title blocks removed
    """
    if not has_root_title:
        return content
    
    if isinstance(content, str):
        # For markdown strings, remove the first # heading if it exists
        lines = content.splitlines()
        if lines and lines[0].strip().startswith("# "):
            return "\n".join(lines[1:]).lstrip()
        return content
    
    if isinstance(content, dict):
        content_type = (content.get("type") or "").strip().lower()
        if content_type == "title":
            return []
        # If it's a dict that's not a title, wrap in list
        return [content]
    
    if isinstance(content, list):
        # Filter out title blocks while preserving other content
        filtered = []
        for item in content:
            if isinstance(item, dict):
                item_type = (item.get("type") or "").strip().lower()
                if item_type != "title":
                    filtered.append(item)
            elif isinstance(item, str):
                if not item.strip().startswith("# "):
                    filtered.append(item)
            else:
                filtered.append(item)
        return filtered if filtered else []
    
    return content


def add_auto_sized_review_comment(cell, text, author="AI Reviewer"):
    """
    Adds a note to an Excel cell, adjusting the width and height
    so that all the text is visible.
    """
    if not text:
        return

    import math
    avg_char_width = 7
    px_per_line = 15
    base_width = 200
    max_width = 500
    min_height = 40

    width = min(max_width, base_width + len(text) * 2)
    chars_per_line = max(1, width // avg_char_width)
    lines = 0
    for paragraph in text.split('\n'):
        lines += math.ceil(len(paragraph) / chars_per_line)
    height = max(min_height, lines * px_per_line)

    from openpyxl.comments import Comment
    comment = Comment(text, author)
    comment.width = width
    comment.height = height
    cell.comment = comment


def _snapshot_runs(p):
    """Return a list of {'text': str, 'font': {...}} for each run in a paragraph."""
    runs = []
    for r in p.runs:
        f = r.font
        font_spec = {
            "name": f.name,
            "size": f.size,
            "bold": f.bold,
            "italic": f.italic,
            "underline": f.underline,
            "color_rgb": getattr(getattr(f.color, "rgb", None), "rgb", None) or getattr(f.color, "rgb", None)
        }
        runs.append({"text": r.text or "", "font": font_spec})
    return runs


def _apply_font(run, font_spec):
    """Apply font specifications to a run."""
    if not font_spec:
        return
    f = run.font
    try:
        if font_spec.get("name") is not None:
            f.name = font_spec["name"]
        if font_spec.get("size") is not None:
            f.size = font_spec["size"]
        if font_spec.get("bold") is not None:
            f.bold = font_spec["bold"]
        if font_spec.get("italic") is not None:
            f.italic = font_spec["italic"]
        if font_spec.get("underline") is not None:
            f.underline = font_spec["underline"]
        rgb = font_spec.get("color_rgb")
        if rgb is not None:
            try:
                f.color.rgb = rgb
            except Exception:
                pass
    except Exception:
        pass


def _set_text_with_runs(shape, new_content):
    """
    Set the text of a shape while preserving the original run-level formatting.
    """
    if not (hasattr(shape, "text_frame") and shape.text_frame):
        return
    tf = shape.text_frame

    if isinstance(new_content, list):
        lines = [str(item) for item in new_content]
    else:
        lines = [str(new_content or "")]

    original_para_styles = []
    original_para_runs = []     

    for p in tf.paragraphs:
        level = int(getattr(p, "level", 0) or 0)
        alignment = getattr(p, "alignment", None)
        original_para_styles.append({"level": level, "alignment": alignment})
        original_para_runs.append(_snapshot_runs(p))

    tf.clear()

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if (i == 0 and tf.paragraphs) else tf.add_paragraph()

        if original_para_styles:
            style = original_para_styles[i] if i < len(original_para_styles) else original_para_styles[-1]
            p.level = style.get("level", 0)
            if style.get("alignment") is not None:
                p.alignment = style["alignment"]

        runs_spec = (
            original_para_runs[i] if i < len(original_para_runs)
            else (original_para_runs[-1] if original_para_runs else [])
        )

        if not runs_spec:
            r = p.add_run()
            r.text = line
            continue

        n = len(runs_spec)
        total = len(line)

        if total == 0:
            for spec in runs_spec:
                r = p.add_run()
                r.text = ""
                _apply_font(r, spec["font"])
        else:
            base, rem = divmod(total, n)
            sizes = [base + (1 if k < rem else 0) for k in range(n)]
            pos = 0
            for k, spec in enumerate(runs_spec):
                seg = line[pos:pos + sizes[k]]
                pos += sizes[k]
                r = p.add_run()
                r.text = seg
                _apply_font(r, spec["font"])


def shape_by_id(slide, shape_id):
    sid = int(shape_id)
    for sh in slide.shapes:
        val = getattr(sh, "shape_id", None) or getattr(getattr(sh, "_element", None), "cNvPr", None)
        val = int(getattr(val, "id", val)) if val is not None else None
        if val == sid:
            return sh
    return None


def ensure_slot_textbox(slide, slot):
    slot = (slot or "").lower()

    def _get(ph_name):
        return getattr(PP_PLACEHOLDER, ph_name, None)

    TITLE = _get("TITLE")
    CENTER_TITLE = _get("CENTER_TITLE")
    SUBTITLE = _get("SUBTITLE")
    BODY = _get("BODY")
    CONTENT = _get("CONTENT")
    OBJECT = _get("OBJECT")

    title_types = {t for t in (TITLE, CENTER_TITLE, SUBTITLE) if t is not None}
    body_types  = {t for t in (BODY, CONTENT, OBJECT) if t is not None}

    def find_placeholder(accepted_types):
        for sh in slide.shapes:
            if not getattr(sh, "is_placeholder", False):
                continue
            pf = getattr(sh, "placeholder_format", None)
            if not pf:
                continue
            try:
                if pf.type in accepted_types:
                    return sh
            except Exception:
                pass
        return None

    if slot == "title":
        ph = find_placeholder(title_types)
        if ph:
            return ph

    if slot == "body":
        ph = find_placeholder(body_types)
        if ph:
            return ph

    if slot == "title":
        return slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    if slot == "body":
        return slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    return slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.5))


def _layout_has(layout, want_title=False, want_body=False):
    has_title = has_body = False
    for ph in getattr(layout, "placeholders", []):
        pf = getattr(ph, "placeholder_format", None)
        t = getattr(pf, "type", None) if pf else None
        if t in (getattr(PP_PLACEHOLDER, "TITLE", None),
                 getattr(PP_PLACEHOLDER, "CENTER_TITLE", None),
                 getattr(PP_PLACEHOLDER, "SUBTITLE", None)):
            has_title = True
        if t in (getattr(PP_PLACEHOLDER, "BODY", None),
                 getattr(PP_PLACEHOLDER, "CONTENT", None),
                 getattr(PP_PLACEHOLDER, "OBJECT", None)):
            has_body = True
    return (not want_title or has_title) and (not want_body or has_body)


def _pick_layout_for_slots(prs, anchor_slide, needs_title, needs_body):
    if anchor_slide and _layout_has(anchor_slide.slide_layout, needs_title, needs_body):
        return anchor_slide.slide_layout
    for layout in prs.slide_layouts:
        if _layout_has(layout, needs_title, needs_body):
            return layout
    return anchor_slide.slide_layout if anchor_slide else prs.slide_layouts[-1]


def _collect_needs(edit_items):
    needs = {}
    for tgt, _ in edit_items:
        if not isinstance(tgt, str):
            continue
        import re
        m = re.match(r"^(n\d+):slot:(title|body)$", tgt.strip(), flags=re.I)
        if m:
            ref, slot = m.group(1), m.group(2).lower()
            needs.setdefault(ref, {"title": False, "body": False})
            needs[ref][slot] = True
    return needs


def _get_pptx_namespaces():
    """Returns XML namespaces for PowerPoint"""
    return {
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p15': 'http://schemas.microsoft.com/office/powerpoint/2012/main',
        'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main'
    }


def _add_native_pptx_comment_zip(pptx_path, slide_num, comment_text, author_id, x=100, y=100):
    """
    Add a native PowerPoint comment by directly manipulating the ZIP file.
    """
    namespaces = _get_pptx_namespaces()
    
    import tempfile
    import os
    from lxml import etree
    
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = os.path.join(temp_dir, "temp_pptx")
        os.makedirs(temp_path)
        
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            zf.extractall(temp_path)
        
        authors_file = os.path.join(temp_path, 'ppt', 'commentAuthors.xml')
        if os.path.exists(authors_file):
            root = etree.parse(authors_file).getroot()
            found = False
            for author in root.findall('.//p:cmAuthor', namespaces):
                if author.get('name') == 'AI Reviewer':
                    author_id = int(author.get('id'))
                    found = True
                    break
            
            if not found:
                existing_ids = [int(a.get('id')) for a in root.findall('.//p:cmAuthor', namespaces)]
                author_id = max(existing_ids) + 1 if existing_ids else 0
                author = etree.SubElement(root, f'{{{namespaces["p"]}}}cmAuthor')
                author.set('id', str(author_id))
                author.set('name', 'AI Reviewer')
                author.set('initials', 'AI')
                author.set('lastIdx', '1')
                author.set('clrIdx', str(author_id % 8))
        else:
            os.makedirs(os.path.join(temp_path, 'ppt'), exist_ok=True)
            root = etree.Element(
                f'{{{namespaces["p"]}}}cmAuthorLst',
                nsmap={k: v for k, v in namespaces.items() if k in ['p']}
            )
            author = etree.SubElement(root, f'{{{namespaces["p"]}}}cmAuthor')
            author.set('id', str(author_id))
            author.set('name', 'AI Reviewer')
            author.set('initials', 'AI')
            author.set('lastIdx', '1')
            author.set('clrIdx', '0')
            
            rels_file = os.path.join(temp_path, 'ppt', '_rels', 'presentation.xml.rels')
            if os.path.exists(rels_file):
                rels_root = etree.parse(rels_file).getroot()
                existing_ids = [int(rel.get('Id')[3:]) for rel in rels_root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship')]
                next_rid = max(existing_ids) + 1 if existing_ids else 1
                
                rel = etree.SubElement(rels_root, '{http://schemas.openxmlformats.org/package/2006/relationships}Relationship')
                rel.set('Id', f'rId{next_rid}')
                rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/commentAuthors')
                rel.set('Target', 'commentAuthors.xml')
                
                with open(rels_file, 'wb') as f:
                    f.write(etree.tostring(rels_root, xml_declaration=True, encoding='UTF-8', pretty_print=True))
        
        with open(authors_file, 'wb') as f:
            f.write(etree.tostring(root, xml_declaration=True, encoding='UTF-8', pretty_print=True))
        
        comments_dir = os.path.join(temp_path, 'ppt', 'comments')
        os.makedirs(comments_dir, exist_ok=True)
        comment_file = os.path.join(comments_dir, f'comment{slide_num}.xml')
        
        if os.path.exists(comment_file):
            comments_root = etree.parse(comment_file).getroot()
        else:
            comments_root = etree.Element(
                f'{{{namespaces["p"]}}}cmLst',
                nsmap={k: v for k, v in namespaces.items() if k in ['p']}
            )
            
            slide_rels_file = os.path.join(temp_path, 'ppt', 'slides', '_rels', f'slide{slide_num}.xml.rels')
            if not os.path.exists(os.path.dirname(slide_rels_file)):
                os.makedirs(os.path.dirname(slide_rels_file), exist_ok=True)
            
            if os.path.exists(slide_rels_file):
                slide_rels_root = etree.parse(slide_rels_file).getroot()
            else:
                slide_rels_root = etree.Element(
                    '{http://schemas.openxmlformats.org/package/2006/relationships}Relationships'
                )
            
            existing_ids = [int(rel.get('Id')[3:]) for rel in slide_rels_root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship')]
            next_rid = max(existing_ids) + 1 if existing_ids else 1
            
            rel = etree.SubElement(slide_rels_root, '{http://schemas.openxmlformats.org/package/2006/relationships}Relationship')
            rel.set('Id', f'rId{next_rid}')
            rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments')
            rel.set('Target', f'../comments/comment{slide_num}.xml')
            
            with open(slide_rels_file, 'wb') as f:
                f.write(etree.tostring(slide_rels_root, xml_declaration=True, encoding='UTF-8', pretty_print=True))
        
        existing_ids = [int(c.get('idx')) for c in comments_root.findall('.//p:cm', namespaces)]
        next_id = max(existing_ids) + 1 if existing_ids else 1
        
        comment = etree.SubElement(comments_root, f'{{{namespaces["p"]}}}cm')
        comment.set('authorId', str(author_id))
        comment.set('dt', datetime.datetime.now().isoformat())
        comment.set('idx', str(next_id))
        
        pos = etree.SubElement(comment, f'{{{namespaces["p"]}}}pos')
        pos.set('x', str(int(x)))
        pos.set('y', str(int(y)))
        
        text_elem = etree.SubElement(comment, f'{{{namespaces["p"]}}}text')
        text_elem.text = comment_text
        
        with open(comment_file, 'wb') as f:
            f.write(etree.tostring(comments_root, xml_declaration=True, encoding='UTF-8', pretty_print=True))
        
        content_types_file = os.path.join(temp_path, '[Content_Types].xml')
        if os.path.exists(content_types_file):
            ct_root = etree.parse(content_types_file).getroot()
            ns = {'ct': 'http://schemas.openxmlformats.org/package/2006/content-types'}
            
            has_authors = False
            has_comments = False
            
            for override in ct_root.findall('.//ct:Override', ns):
                if override.get('PartName') == '/ppt/commentAuthors.xml':
                    has_authors = True
                if override.get('PartName') == f'/ppt/comments/comment{slide_num}.xml':
                    has_comments = True
            
            if not has_authors:
                override = etree.SubElement(ct_root, '{http://schemas.openxmlformats.org/package/2006/content-types}Override')
                override.set('PartName', '/ppt/commentAuthors.xml')
                override.set('ContentType', 'application/vnd.openxmlformats-officedocument.presentationml.commentAuthors+xml')
            
            if not has_comments:
                override = etree.SubElement(ct_root, '{http://schemas.openxmlformats.org/package/2006/content-types}Override')
                override.set('PartName', f'/ppt/comments/comment{slide_num}.xml')
                override.set('ContentType', 'application/vnd.openxmlformats-officedocument.presentationml.comments+xml')
            
            with open(content_types_file, 'wb') as f:
                f.write(etree.tostring(ct_root, xml_declaration=True, encoding='UTF-8', pretty_print=True))
        
        with zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for root_dir, dirs, files in os.walk(temp_path):
                for file in files:
                    file_path = os.path.join(root_dir, file)
                    arcname = os.path.relpath(file_path, temp_path)
                    zf.write(file_path, arcname)
        
        log.debug(f"Native comment added to slide {slide_num} with idx={next_id}")