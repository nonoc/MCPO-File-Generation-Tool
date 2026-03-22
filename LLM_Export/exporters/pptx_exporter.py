"""
PowerPoint (PPTX) exporter module using python-pptx.
"""

import os
import logging
from typing import Any, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.enum.shapes import PP_PLACEHOLDER 
from pptx.enum.text import MSO_AUTO_SIZE
from io import BytesIO

# Import from new modules (try relative first)
try:
    from ..shared.utils import _generate_unique_folder, _public_url, render_text_with_emojis
    from ..shared.constants import PPTX_TEMPLATE, PPTX_TEMPLATE_PATH
except ImportError:
    from shared.utils import _generate_unique_folder, _public_url, render_text_with_emojis
    from shared.constants import PPTX_TEMPLATE, PPTX_TEMPLATE_PATH

log = logging.getLogger(__name__)


def _normalize_pptx_content(content: Any) -> List[str]:
    """
    Normalize content for PowerPoint rendering.
    
    Converts raw block dictionaries to plain text strings for PowerPoint slides.
    Handles title blocks, paragraphs, lists, and other block types.
    
    Args:
        content: Content to normalize (string, dict, list, or raw block)
        
    Returns:
        List of plain text strings for slide content
    """
    if content is None:
        return []
    
    if isinstance(content, str):
        return [content]
    
    if isinstance(content, dict):
        content_type = (content.get("type") or "").strip().lower()
        text = content.get("text") or content.get("title") or content.get("content") or ""
        
        # Handle title blocks
        if content_type == "title":
            return [str(text).strip()] if text else []
        
        # Handle section blocks
        if content_type == "section":
            title = content.get("title") or ""
            children = content.get("children") or []
            result = [str(title).strip()] if title else []
            for child in children:
                result.extend(_normalize_pptx_content(child))
            return result
        
        # Handle paragraph-like blocks
        if content_type in {"paragraph", "text", "body", "description", "summary"}:
            return [str(text).strip()] if text else []
        
        # Handle label paragraph
        if content_type == "label_paragraph":
            label = content.get("label") or ""
            body = content.get("text") or ""
            if label and body:
                return [f"{label}: {body}"]
            elif label:
                return [str(label).strip()]
            elif body:
                return [str(body).strip()]
            return []
        
        # Handle list blocks
        if content_type in {"list", "bullet_list"} or "list" in content_type:
            items = content.get("items") or []
            result = []
            for item in items:
                if isinstance(item, dict):
                    item_text = item.get("text") or ""
                    result.append(str(item_text).strip())
                elif isinstance(item, str):
                    result.append(item.strip())
            return result if result else []
        
        # Handle bullet blocks
        if content_type in {"bullet", "list_item"}:
            text = content.get("text") or ""
            return [str(text).strip()] if text else []
        
        # Fallback: try to extract any text
        if text:
            return [str(text).strip()]
        
        return []
    
    if isinstance(content, list):
        result = []
        for item in content:
            result.extend(_normalize_pptx_content(item))
        return result
    
    return [str(content)]


def dynamic_font_size(content_list: List[str], max_chars: int = 400, base_size: int = 28, min_size: int = 12) -> PptPt:
    """Calculate dynamic font size based on content length."""
    total_chars = sum(len(line) for line in content_list)
    ratio = total_chars / max_chars if max_chars > 0 else 1
    if ratio <= 1:
        return PptPt(base_size)
    else:
        new_size = int(base_size / ratio)
        return PptPt(max(min_size, new_size))


def _create_presentation(slides_data: List[dict], filename: str, folder_path: Optional[str] = None, title: Optional[str] = None) -> dict:
    """
    Create a PowerPoint presentation from slide data.
    
    Args:
        slides_data: List of slide data dictionaries
        filename: Output filename
        folder_path: Optional folder path for output
        title: Optional presentation title
        
    Returns:
        Dict with url and path to generated file
    """
    if folder_path is None:
        folder_path = _generate_unique_folder()
        
    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        from ..shared.utils import _generate_filename
        filepath, fname = _generate_filename(folder_path, "pptx")

    # Handle template
    use_template = False
    prs = None
    title_layout = None
    content_layout = None

    if PPTX_TEMPLATE:
        try:
            log.debug("Attempting to load template...")
            src = PPTX_TEMPLATE
            if hasattr(PPTX_TEMPLATE, "slides") and hasattr(PPTX_TEMPLATE, "save"):
                log.debug("Template is a Presentation object, converting to BytesIO")
                buf = BytesIO()
                PPTX_TEMPLATE.save(buf)
                buf.seek(0)
                src = buf

            tmp = Presentation(src)
            log.debug(f"Template loaded with {len(tmp.slides)} slides")
            if len(tmp.slides) >= 1:
                prs = tmp
                use_template = True

                title_layout = prs.slides[0].slide_layout
                content_layout = prs.slides[1].slide_layout if len(prs.slides) >= 2 else prs.slides[0].slide_layout
                log.debug("Using template layouts")

                for i in range(len(prs.slides) - 1, 0, -1):
                    rId = prs.slides._sldIdLst[i].rId 
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[i]
        except Exception as e:
            log.error(f"Error loading template: {e}")
            use_template = False
            prs = None

    if not use_template:
        log.debug("No valid template, creating new presentation with default layouts")
        prs = Presentation()
        title_layout = prs.slide_layouts[0]
        content_layout = prs.slide_layouts[1]

    # Create title slide
    if use_template:
        log.debug("Using template title slide")
        tslide = prs.slides[0]
        if tslide.shapes.title:
            tslide.shapes.title.text = title or ""
            for p in tslide.shapes.title.text_frame.paragraphs:
                for r in p.runs:
                    title_info = next(({'size': PptPt(int(child.attrib.get('sz', 2800))/100), 'bold': child.attrib.get('b', '0') == '1'} for child in title_layout.element.iter() if 'defRPr' in child.tag.split('}')[-1] and 'sz' in child.attrib), {'size': PptPt(28), 'bold': True})
                    r.font.size = title_info['size'] 
                    r.font.bold = title_info['bold']
    else:
        log.debug("Creating new title slide")
        tslide = prs.slides.add_slide(title_layout)
        if tslide.shapes.title:
            tslide.shapes.title.text = title or ""
            for p in tslide.shapes.title.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = PptPt(28)
                    r.font.bold = True

    # Process slides
    EMU_PER_IN = 914400
    slide_w_in = prs.slide_width / EMU_PER_IN
    slide_h_in = prs.slide_height / EMU_PER_IN
    log.debug(f"Slide dimensions: {slide_w_in} x {slide_h_in} inches")

    page_margin = 0.5
    gutter = 0.3

    for i, slide_data in enumerate(slides_data):
        log.debug(f"Processing slide {i+1}: {slide_data.get('title', 'Untitled')}")
        if not isinstance(slide_data, dict):
            log.warning(f"Slide data is not a dict, skipping slide {i+1}")
            continue

        slide_title = slide_data.get("title", "Untitled")
        content_list = slide_data.get("content", [])
        if not isinstance(content_list, list):
            content_list = [content_list]
        
        # Normalize content to handle raw block dictionaries
        normalized_content = []
        for item in content_list:
            normalized_content.extend(_normalize_pptx_content(item))
        content_list = normalized_content
        
        log.debug(f"Adding slide with title: '{slide_title}'")
        slide = prs.slides.add_slide(content_layout)

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = slide_title
            for p in slide.shapes.title.text_frame.paragraphs:
                for r in p.runs:
                    title_info = next(({'size': PptPt(int(child.attrib.get('sz', 2800))/100), 'bold': child.attrib.get('b', '0') == '1'} for child in content_layout.element.iter() if 'defRPr' in child.tag.split('}')[-1] and 'sz' in child.attrib), {'size': PptPt(28), 'bold': True})
                    r.font.size = title_info['size'] 
                    r.font.bold = title_info['bold']

        # Find content placeholder
        content_shape = _find_content_placeholder(slide)
        
        if content_shape is None:
            content_shape = slide.shapes.add_textbox(
                Inches(page_margin), 
                Inches(1.2), 
                Inches(slide_w_in - 2*page_margin), 
                Inches(slide_h_in - 2.2)
            )
            log.debug("Creating new textbox for content")

        # Set title bottom position
        title_bottom_in = 1.0 
        if slide.shapes.title:
            try:
                title_bottom_emu = slide.shapes.title.top + slide.shapes.title.height
                title_bottom_in = max(title_bottom_emu / EMU_PER_IN, 1.0)
                title_bottom_in += 0.2
            except Exception:
                title_bottom_in = 1.2 

        # Position content shape
        content_left_in = page_margin
        content_top_in = title_bottom_in
        content_width_in = slide_w_in - 2*page_margin
        content_height_in = slide_h_in - (title_bottom_in + page_margin)

        # Handle image query
        image_query = slide_data.get("image_query")
        if image_query:
            image_url = search_image(image_query)
            if image_url:
                try:
                    response = requests.get(image_url, timeout=30)
                    response.raise_for_status()
                    image_data = response.content
                    image_stream = BytesIO(image_data)
                    pos = slide_data.get("image_position", "right")
                    size = slide_data.get("image_size", "medium")
                    
                    if size == "small":
                        img_w_in, img_h_in = 2.0, 1.5
                    elif size == "large":
                        img_w_in, img_h_in = 4.0, 3.0
                    else:
                        img_w_in, img_h_in = 3.0, 2.0

                    if pos == "left":
                        img_left_in = page_margin
                        img_top_in = title_bottom_in
                        content_left_in = img_left_in + img_w_in + gutter
                    elif pos == "right":
                        img_left_in = max(slide_w_in - page_margin - img_w_in, page_margin)
                        img_top_in = title_bottom_in
                        content_left_in = page_margin
                        content_width_in = max(img_left_in - gutter - content_left_in, 2.5)
                    elif pos == "top":
                        img_left_in = slide_w_in - page_margin - img_w_in
                        img_top_in = title_bottom_in
                        content_left_in = page_margin
                        content_top_in = img_top_in + img_h_in + gutter
                        content_height_in = max(slide_h_in - page_margin - content_top_in, 2.0)
                    elif pos == "bottom":
                        img_left_in = slide_w_in - page_margin - img_w_in
                        img_top_in = max(slide_h_in - page_margin - img_h_in, page_margin)
                        content_left_in = page_margin
                        content_top_in = title_bottom_in
                        content_height_in = max(img_top_in - gutter - content_top_in, 2.0)
                    else:
                        img_left_in = max(slide_w_in - page_margin - img_w_in, page_margin)
                        img_top_in = title_bottom_in

                    slide.shapes.add_picture(image_stream, Inches(img_left_in), Inches(img_top_in), Inches(img_w_in), Inches(img_h_in))
                except Exception:
                    pass

        # Apply content shape positioning
        try:
            content_shape.left = Inches(content_left_in)
            content_shape.top = Inches(content_top_in)
            content_shape.width = Inches(content_width_in)
            content_shape.height = Inches(content_height_in)
        except Exception:
            pass

        # Calculate font size and add content
        approx_chars_per_in = 9.5
        approx_lines_per_in = 1.6
        safe_width = max(content_width_in, 0.1)
        safe_height = max(content_height_in, 0.1)
        est_capacity = int(safe_width * approx_chars_per_in * safe_height * approx_lines_per_in)
        font_size = dynamic_font_size(content_list, max_chars=max(est_capacity, 120), base_size=24, min_size=12)

        try:
            tf = content_shape.text_frame
        except Exception:
            try:
                tf = content_shape.text_frame
            except Exception:
                log.warning("Could not access text frame for content shape")
                continue

        if not tf.paragraphs:
            tf.add_paragraph()
        for idx, line in enumerate(content_list):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = str(line) if line is not None else ""
            run.font.size = font_size
            p.space_after = PptPt(6)

    prs.save(filepath)
    return {"url": _public_url(folder_path, fname), "path": filepath}


def _find_content_placeholder(slide):
    """Find the content placeholder in a slide."""
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.idx == 1:
                return ph
        except Exception:
            pass
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.idx != 0:
                return ph
        except Exception:
            pass
    return None


def search_image(query):
    """Search for an image using configured source."""
    from ..shared.images import search_image as img_search
    return img_search(query)