import re
import os
import ast
import csv
import json
import uuid
import emoji
import math
import time
import base64
import shutil
import datetime
import tarfile
import zipfile
import py7zr
import logging
import requests
from requests import get, post
from requests.auth import HTTPBasicAuth
import threading
import markdown2
import tempfile
from pathlib import Path
from lxml import etree
from PIL import Image
from docx import Document
from docx.shared import Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.shared import qn
from docx.oxml.ns import nsdecls
from docx.oxml import parse_xml
from docx.shared import Pt as DocxPt
from bs4 import BeautifulSoup, NavigableString
from mcp.server.fastmcp import FastMCP, Context
from mcp.server.session import ServerSession
from openpyxl import Workbook, load_workbook
from openpyxl.comments import Comment
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt as PptPt
from pptx.enum.shapes import PP_PLACEHOLDER 
from pptx.parts.image import Image
from pptx.enum.text import MSO_AUTO_SIZE
from io import BytesIO
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem, Image as ReportLabImage, Table as ReportLabTable, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.lib.units import mm

#NonDockerImport
import asyncio
import uvicorn
from typing import Any, List, Iterable, cast
from mcp.server.sse import SseServerTransport
from starlette.requests import Request
from starlette.applications import Starlette
from starlette.routing import Route, Mount 
from starlette.responses import Response, JSONResponse 

SCRIPT_VERSION = "0.8.1"

URL = os.getenv('OWUI_URL')
TOKEN = os.getenv('JWT_SECRET') ## will be deleted in 1.0.0

PERSISTENT_FILES = os.getenv("PERSISTENT_FILES", "false")
FILES_DELAY = int(os.getenv("FILES_DELAY", 60)) 

DEFAULT_PATH_ENV = os.getenv("PYTHONPATH", r"").rstrip("/")
EXPORT_DIR_ENV = os.getenv("FILE_EXPORT_DIR")
EXPORT_DIR = (EXPORT_DIR_ENV or os.path.join(DEFAULT_PATH_ENV, "output")).rstrip("/")
os.makedirs(EXPORT_DIR, exist_ok=True)

BASE_URL_ENV = os.getenv("FILE_EXPORT_BASE_URL")
BASE_URL = (BASE_URL_ENV or "http://localhost:9003/files").rstrip("/")

LOG_LEVEL_ENV = os.getenv("LOG_LEVEL")
LOG_FORMAT_ENV = os.getenv(
    "LOG_FORMAT", "%(asctime)s %(levelname)s %(name)s - %(message)s"
)

DOCS_TEMPLATE_DIR_ENV = os.getenv("DOCS_TEMPLATE_DIR")
DOCS_TEMPLATE_PATH = ((DOCS_TEMPLATE_DIR_ENV or os.path.join(DEFAULT_PATH_ENV, "templates")).rstrip("/"))
os.makedirs(DOCS_TEMPLATE_PATH, exist_ok=True)
PPTX_TEMPLATE = None
DOCX_TEMPLATE = None
XLSX_TEMPLATE = None
PPTX_TEMPLATE_PATH = None
DOCX_TEMPLATE_PATH = None
XLSX_TEMPLATE_PATH = None

TABLE_SEPARATOR_RE = re.compile(r"^\s*\|?(?:\s*:?-+:?\s*\|)+\s*$")

if DOCS_TEMPLATE_PATH and os.path.exists(DOCS_TEMPLATE_PATH):
    logging.debug(f"Template Folder: {DOCS_TEMPLATE_PATH}")
    for root, dirs, files in os.walk(DOCS_TEMPLATE_PATH):
        for file in files:
            fpath = os.path.join(root, file)
            if file.lower().endswith(".pptx") and PPTX_TEMPLATE_PATH is None:
                PPTX_TEMPLATE_PATH = fpath
                logging.debug(f"PPTX template: {PPTX_TEMPLATE_PATH}")
            elif file.lower().endswith(".docx") and DOCX_TEMPLATE_PATH is None:
                DOCX_TEMPLATE_PATH = fpath
            elif file.lower().endswith(".xlsx") and XLSX_TEMPLATE_PATH is None:
                XLSX_TEMPLATE_PATH = fpath
    if PPTX_TEMPLATE_PATH:
        try:
            PPTX_TEMPLATE = Presentation(PPTX_TEMPLATE_PATH)
            logging.debug(f"Using PPTX template: {PPTX_TEMPLATE_PATH}")
        except Exception as e:
            logging.warning(f"PPTX template failed to load : {e}")
            PPTX_TEMPLATE = None
    else:
        logging.debug("No PPTX template found. Creation of a blank document.")
        PPTX_TEMPLATE = None

    if DOCX_TEMPLATE_PATH and os.path.exists(DOCS_TEMPLATE_PATH):
        try:
            DOCX_TEMPLATE = Document(DOCX_TEMPLATE_PATH)
            logging.debug(f"Using DOCX template: {DOCX_TEMPLATE_PATH}")
        except Exception as e:
            logging.warning(f"DOCX template failed to load : {e}")
            DOCX_TEMPLATE = None
    else:
        logging.debug("No DOCX template found. Creation of a blank document.")
        DOCX_TEMPLATE = None
    
    XLSX_TEMPLATE_PATH = os.path.join(DEFAULT_PATH_ENV, "templates", "Default_Template.xlsx")

    if XLSX_TEMPLATE_PATH:
        try:
            XLSX_TEMPLATE = load_workbook(XLSX_TEMPLATE_PATH)
            logging.debug(f"Using XLSX template: {XLSX_TEMPLATE_PATH}")
        except Exception as e:
            logging.warning(f"Failed to load XLSX template: {e}")
            XLSX_TEMPLATE = None
    else:
        logging.debug("No XLSX template found. Creation of a blank document.")
        XLSX_TEMPLATE = None


def search_image(query):
    log.debug(f"Searching for image with query: '{query}'")
    image_source = os.getenv("IMAGE_SOURCE", "unsplash")

    if image_source == "unsplash":
        return search_unsplash(query)
    elif image_source == "local_sd":
        return search_local_sd(query)
    elif image_source == "pexels":
        return search_pexels(query)
    else:
        log.warning(f"Image source unknown : {image_source}")
        return None

def search_local_sd(query: str):
    log.debug(f"Searching for local SD image with query: '{query}'")
    SD_URL = os.getenv("LOCAL_SD_URL")
    SD_USERNAME = os.getenv("LOCAL_SD_USERNAME")
    SD_PASSWORD = os.getenv("LOCAL_SD_PASSWORD")
    DEFAULT_MODEL = os.getenv("LOCAL_SD_DEFAULT_MODEL", "sd_xl_base_1.0.safetensors")
    DEFAULT_STEPS = int(os.getenv("LOCAL_SD_STEPS", 20))
    DEFAULT_WIDTH = int(os.getenv("LOCAL_SD_WIDTH", 512))
    DEFAULT_HEIGHT = int(os.getenv("LOCAL_SD_HEIGHT", 512))
    DEFAULT_CFG_SCALE = float(os.getenv("LOCAL_SD_CFG_SCALE", 1.5))
    DEFAULT_SCHEDULER = os.getenv("LOCAL_SD_SCHEDULER", "Karras")
    DEFAULT_SAMPLE = os.getenv("LOCAL_SD_SAMPLE", "Euler a")

    if not SD_URL:
        log.warning("LOCAL_SD_URL is not defined.")
        return None

    payload = {
        "prompt": query.strip(),
        "steps": DEFAULT_STEPS,
        "width": DEFAULT_WIDTH,
        "height": DEFAULT_HEIGHT,
        "cfg_scale": DEFAULT_CFG_SCALE,
        "sampler_name": DEFAULT_SAMPLE,
        "scheduler": DEFAULT_SCHEDULER,
        "enable_hr": False,
        "hr_upscaler": "Latent",
        "seed": -1,
        "override_settings": {
            "sd_model_checkpoint": DEFAULT_MODEL
        }
    }

    try:
        url = f"{SD_URL}/sdapi/v1/txt2img"
        log.debug(f"Sending request to local SD API at {url}")
        response = requests.post(
            url,
            json=payload,
            headers={"Content-Type": "application/json"},
            auth=HTTPBasicAuth(SD_USERNAME, SD_PASSWORD),
            timeout=30
        )
        response.raise_for_status()
        data = response.json()

        images = data.get("images", [])
        if not images:
            log.warning(f"No image generated for the request : '{query}'")
            return None

        image_b64 = images[0]
        image_data = base64.b64decode(image_b64)

        folder_path = _generate_unique_folder()
        filename = f"{query.replace(' ', '_')}.png"
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)

        with open(filepath, "wb") as f:
            f.write(image_data)

        return _public_url(folder_path, filename)

    except requests.exceptions.Timeout:
        log.error(f"Timeout during generation for : '{query}'")
    except requests.exceptions.RequestException as e:
        log.error(f"Network error : {e}")
    except Exception as e:
        log.error(f"Unexpected error : {e}")

    return None

def search_unsplash(query):
    log.debug(f"Searching Unsplash for query: '{query}'")
    api_key = os.getenv("UNSPLASH_ACCESS_KEY")
    if not api_key:
        log.warning("UNSPLASH_ACCESS_KEY is not set. Cannot search for images.")
        return None
    url = "https://api.unsplash.com/search/photos"
    params = {
        "query": query,
        "per_page": 1,
        "orientation": "landscape"
    }
    headers = {"Authorization": f"Client-ID {api_key}"}
    log.debug(f"Sending request to Unsplash API")
    try:
        response = requests.get(url, params=params, headers=headers)
        log.debug(f"Unsplash API response status: {response.status_code}")
        response.raise_for_status() 
        data = response.json()
        if data.get("results"):
            image_url = data["results"][0]["urls"]["regular"]
            log.debug(f"Found image URL for '{query}': {image_url}")
            return image_url
        else:
            log.debug(f"No results found on Unsplash for query: '{query}'")
    except requests.exceptions.RequestException as e:
        log.error(f"Network error while searching image for '{query}': {e}")
    except json.JSONDecodeError as e:
        log.error(f"Error decoding JSON from Unsplash for '{query}': {e}")
    except Exception as e:
        log.error(f"Unexpected error searching image for '{query}': {e}")
    return None 

def search_pexels(query):
    log.debug(f"Searching Pexels for query: '{query}'")
    api_key = os.getenv("PEXELS_ACCESS_KEY")
    if not api_key:
        log.warning("PEXELS_ACCESS_KEY is not set. Cannot search for images.")
        return None
    url = "https://api.pexels.com/v1/search"
    params = {
        "query": query,
        "per_page": 1,
        "orientation": "landscape"
    }
    headers = {"Authorization": f"{api_key}"}
    log.debug(f"Sending request to Pexels API")
    try:
        response = requests.get(url, params=params, headers=headers)
        log.debug(f"Pexels API response status: {response.status_code}")
        response.raise_for_status() 
        data = response.json()
        if data.get("photos"):
            image_url = data["photos"][0]["src"]["large"]
            log.debug(f"Found image URL for '{query}': {image_url}")
            return image_url
        else:
            log.debug(f"No results found on Pexels for query: '{query}'")
    except requests.exceptions.RequestException as e:
        log.error(f"Network error while searching image for '{query}': {e}")
    except json.JSONDecodeError as e:
        log.error(f"Error decoding JSON from Pexels for '{query}': {e}")
    except Exception as e:
        log.error(f"Unexpected error searching image for '{query}': {e}")
    return None

def _resolve_log_level(val: str | None) -> int:
    if not val:
        return logging.INFO
    v = val.strip()
    if v.isdigit():
        try:
            return int(v)
        except ValueError:
            return logging.INFO
    return getattr(logging, v.upper(), logging.INFO)

logging.basicConfig(
    level=_resolve_log_level(LOG_LEVEL_ENV),
    format=LOG_FORMAT_ENV,
)
log = logging.getLogger("file_export_mcp")
log.setLevel(_resolve_log_level(LOG_LEVEL_ENV))
log.info("Effective LOG_LEVEL -> %s", logging.getLevelName(log.level))

mcp = FastMCP(
    name = "file_export"
)

def dynamic_font_size(content_list, max_chars=400, base_size=28, min_size=12):
    total_chars = sum(len(line) for line in content_list)
    ratio = total_chars / max_chars if max_chars > 0 else 1
    if ratio <= 1:
        return PptPt(base_size)
    else:
        new_size = int(base_size / ratio)
        return PptPt(max(min_size, new_size))

def _public_url(folder_path: str, filename: str) -> str:
    """Build a stable public URL for a generated file."""
    folder = os.path.basename(folder_path).lstrip("/")
    name = filename.lstrip("/")
    return f"{BASE_URL}/{folder}/{name}"

def _generate_unique_folder() -> str:
    folder_name = f"export_{uuid.uuid4().hex[:10]}_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}"
    folder_path = os.path.join(EXPORT_DIR, folder_name)
    os.makedirs(folder_path, exist_ok=True)
    return folder_path

def _generate_filename(folder_path: str, ext: str, filename: str = None) -> tuple[str, str]:
    if not filename:
        filename = f"export_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.{ext}"
    base, ext = os.path.splitext(filename)
    filepath = os.path.join(folder_path, filename)
    counter = 1
    while os.path.exists(filepath):
        filename = f"{base}_{counter}{ext}"
        filepath = os.path.join(folder_path, filename)
        counter += 1
    return filepath, filename

styles = getSampleStyleSheet()
styles.add(ParagraphStyle(
    name="CustomHeading1",
    parent=styles["Heading1"],
    textColor=colors.HexColor("#0A1F44"),
    fontSize=18,
    spaceAfter=16,
    spaceBefore=12,
    alignment=TA_LEFT
))
styles.add(ParagraphStyle(
    name="CustomHeading2",
    parent=styles["Heading2"],
    textColor=colors.HexColor("#1C3F77"),
    fontSize=14,
    spaceAfter=12,
    spaceBefore=10,
    alignment=TA_LEFT
))
styles.add(ParagraphStyle(
    name="CustomHeading3",
    parent=styles["Heading3"],
    textColor=colors.HexColor("#3A6FB0"), 
    fontSize=12,
    spaceAfter=10,
    spaceBefore=8,
    alignment=TA_LEFT
))
styles.add(ParagraphStyle(
    name="CustomNormal",
    parent=styles["Normal"],
    fontSize=11,
    leading=14,
    alignment=TA_LEFT
))
styles.add(ParagraphStyle(
    name="CustomListItem",
    parent=styles["Normal"],
    fontSize=11,
    leading=14,
    alignment=TA_LEFT
))
styles.add(ParagraphStyle(
    name="CustomCode",
    parent=styles["Code"],
    fontSize=10,
    leading=12,
    fontName="Courier",
    backColor=colors.HexColor("#F5F5F5"),
    borderColor=colors.HexColor("#CCCCCC"),
    borderWidth=1,
    leftIndent=10,
    rightIndent=10,
    topPadding=5,
    bottomPadding=5
))
styles.add(ParagraphStyle(
    name="StructuredDocumentTitle",
    parent=styles["CustomHeading1"],
    fontSize=24,
    leading=28,
    alignment=TA_LEFT,
    spaceBefore=0,
    spaceAfter=18
))
styles.add(ParagraphStyle(
    name="StructuredParagraph",
    parent=styles["CustomNormal"],
    leading=16,
    spaceBefore=4,
    spaceAfter=10
))
styles.add(ParagraphStyle(
    name="StructuredListItem",
    parent=styles["CustomListItem"],
    leftIndent=6,
    spaceBefore=2,
    spaceAfter=2
))
styles.add(ParagraphStyle(
    name="StructuredSourcesHeading",
    parent=styles["CustomHeading2"],
    textColor=colors.HexColor("#1F3474"),
    fontSize=13,
    leading=16,
    alignment=TA_LEFT,
    spaceBefore=18,
    spaceAfter=6
))
styles.add(ParagraphStyle(
    name="StructuredSourcesItem",
    parent=styles["CustomNormal"],
    fontSize=10,
    italic=True,
    leading=14,
    leftIndent=4,
    spaceBefore=2,
    spaceAfter=4
))

def render_text_with_emojis(text: str) -> str:
    if not text:
        return ""
    try:
        converted = emoji.emojize(text, language="alias")
        return converted
    except Exception as e:
        log.error(f"Error in emoji conversion: {e}")
        return text

def process_list_items(ul_or_ol_element, is_ordered=False):
    items = []
    bullet_type = '1' if is_ordered else 'bullet'
    for li in ul_or_ol_element.find_all('li', recursive=False):
        li_text_parts = []
        for content in li.contents:
            if isinstance(content, NavigableString):
                li_text_parts.append(str(content))
            elif content.name not in ['ul', 'ol']:
                 li_text_parts.append(content.get_text())
        li_text = ''.join(li_text_parts).strip()
        list_item_paragraph = None
        if li_text:
            rendered_text = render_text_with_emojis(li_text)
            list_item_paragraph = Paragraph(rendered_text, styles["CustomListItem"])
        sub_lists = li.find_all(['ul', 'ol'], recursive=False)
        sub_flowables = []
        if list_item_paragraph:
             sub_flowables.append(list_item_paragraph)
        for sub_list in sub_lists:
            is_sub_ordered = sub_list.name == 'ol'
            nested_items = process_list_items(sub_list, is_sub_ordered)
            if nested_items:
                nested_list_flowable = ListFlowable(
                    cast(Iterable, nested_items),
                    bulletType='1' if is_sub_ordered else 'bullet',
                    leftIndent=10 * mm,
                    bulletIndent=5 * mm,
                    spaceBefore=2,
                    spaceAfter=2
                )
                sub_flowables.append(nested_list_flowable)
        if sub_flowables:
            items.append(ListItem(sub_flowables))
    return items

def render_html_elements(soup):
    log.debug("Starting render_html_elements...")
    story = []
    element_count = 0
    for elem in soup.children:
        element_count += 1
        log.debug(f"Processing element #{element_count}: {type(elem)}, name={getattr(elem, 'name', 'NavigableString')}")
        if isinstance(elem, NavigableString):
            text = str(elem).strip()
            if text:
                log.debug(f"Adding Paragraph from NavigableString: {text[:50]}...")
                story.append(Paragraph(render_text_with_emojis(text), styles["CustomNormal"]))
                story.append(Spacer(1, 6))
        elif hasattr(elem, 'name'):
            tag_name = elem.name
            log.debug(f"Handling tag: <{tag_name}>")
            if tag_name == "h1":
                text = render_text_with_emojis(elem.get_text().strip())
                log.debug(f"Adding H1: {text[:50]}...")
                story.append(Paragraph(text, styles["CustomHeading1"]))
                story.append(Spacer(1, 10))
            elif tag_name == "h2":
                text = render_text_with_emojis(elem.get_text().strip())
                log.debug(f"Adding H2: {text[:50]}...")
                story.append(Paragraph(text, styles["CustomHeading2"]))
                story.append(Spacer(1, 8))
            elif tag_name == "h3":
                text = render_text_with_emojis(elem.get_text().strip())
                log.debug(f"Adding H3: {text[:50]}...")
                story.append(Paragraph(text, styles["CustomHeading3"]))
                story.append(Spacer(1, 6))
            elif tag_name == "p":
                imgs = elem.find_all("img")
                if imgs:
                    for img_tag in imgs:
                        src = img_tag.get("src")
                        alt = img_tag.get("alt", "[Image]")
                        try:
                            if src and src.startswith("http"):
                                log.debug(f"Downloading image from URL: {src}")
                                response = requests.get(src)
                                response.raise_for_status()
                                img_data = BytesIO(response.content)
                                img = Image(img_data, width=200, height=150)
                            else:
                                log.debug(f"Loading local image: {src}")
                                img = Image(src, width=200, height=150)
                            story.append(img)
                            story.append(Spacer(1, 10))
                        except Exception as e:
                            log.error(f"Error loading image {src}: {e}")
                            story.append(Paragraph(f"[Image: {alt}]", styles["CustomNormal"]))
                            story.append(Spacer(1, 6))
                else:
                    text = render_text_with_emojis(elem.get_text().strip())
                    if text:
                        log.debug(f"Adding Paragraph: {text[:50]}...")
                        story.append(Paragraph(text, styles["CustomNormal"]))
                        story.append(Spacer(1, 6))
            elif tag_name in ["ul", "ol"]:
                is_ordered = tag_name == "ol"
                log.debug(f"Processing list (ordered={is_ordered})...")
                items = process_list_items(elem, is_ordered)
                if items:
                    log.debug(f"Adding ListFlowable with {len(items)} items")
                    story.append(ListFlowable(cast(Iterable, items),
                        bulletType='1' if is_ordered else 'bullet',
                        leftIndent=10 * mm,
                        bulletIndent=5 * mm,
                        spaceBefore=6,
                        spaceAfter=10
                    ))
            elif tag_name == "blockquote":
                text = render_text_with_emojis(elem.get_text().strip())
                if text:
                    log.debug(f"Adding Blockquote: {text[:50]}...")
                    story.append(Paragraph(f"{text}", styles["CustomNormal"]))
                    story.append(Spacer(1, 8))
            elif tag_name in ["code", "pre"]:
                text = elem.get_text().strip()
                if text:
                    log.debug(f"Adding Code/Pre block: {text[:50]}...")
                    story.append(Paragraph(text, styles["CustomCode"]))
                    story.append(Spacer(1, 6 if tag_name == "code" else 8))
            elif tag_name == "img":
                src = elem.get("src")
                alt = elem.get("alt", "[Image]")
                log.debug(f"Found <img> tag. src='{src}', alt='{alt}'")
                if src is not None: 
                    try:
                        if src.startswith("image_query:"):

                            query = src.replace("image_query:", "").strip()
                            log.debug(f"Handling image_query: '{query}'")
                            image_url = search_image(query)
                            if image_url:
                                log.debug(f"Downloading image from Unsplash URL: {image_url}")
                                response = requests.get(image_url)
                                log.debug(f"Image download response status: {response.status_code}")
                                response.raise_for_status()
                                img_data = BytesIO(response.content)
                                img = ReportLabImage(img_data, width=200, height=150)
                                log.debug("Adding ReportLab Image object to story (Unsplash)")
                                story.append(img)
                                story.append(Spacer(1, 10))
                            else:
                                log.warning(f"No image found for query: {query}")
                                story.append(Paragraph(f"[Image non trouvee pour: {query}]", styles["CustomNormal"]))
                                story.append(Spacer(1, 6))
                        elif src.startswith("http"):
                            log.debug(f"Downloading image from direct URL: {src}")
                            response = requests.get(src)
                            log.debug(f"Image download response status: {response.status_code}")
                            response.raise_for_status()
                            img_data = BytesIO(response.content)
                            img = ReportLabImage(img_data, width=200, height=150)
                            log.debug("Adding ReportLab Image object to story (Direct URL)")
                            story.append(img)
                            story.append(Spacer(1, 10))
                        else:
                            log.debug(f"Loading local image: {src}")
                            if os.path.exists(src):
                                img = ReportLabImage(src, width=200, height=150)
                                log.debug("Adding ReportLab Image object to story (Local)")
                                story.append(img)
                                story.append(Spacer(1, 10))
                            else:
                               log.error(f"Local image file not found: {src}")
                               story.append(Paragraph(f"[Image locale non trouvee: {src}]", styles["CustomNormal"]))
                               story.append(Spacer(1, 6))
                    except requests.exceptions.RequestException as e:
                        log.error(f"Network error loading image {src}: {e}")
                        story.append(Paragraph(f"[Image (erreur reseau): {alt}]", styles["CustomNormal"]))
                        story.append(Spacer(1, 6))
                    except Exception as e:
                        log.error(f"Error processing image {src}: {e}", exc_info=True) 
                        story.append(Paragraph(f"[Image: {alt}]", styles["CustomNormal"]))
                        story.append(Spacer(1, 6))
                else:
                    log.warning("Image tag found with no 'src' attribute.")
                    story.append(Paragraph(f"[Image: {alt} (source manquante)]", styles["CustomNormal"]))
                    story.append(Spacer(1, 6))
            elif tag_name == "br":
                log.debug("Adding Spacer for <br>")
                story.append(Spacer(1, 6))
            else:
                text = elem.get_text().strip()
                if text:
                    log.debug(f"Adding Paragraph for unknown tag <{tag_name}>: {text[:50]}...")
                    story.append(Paragraph(render_text_with_emojis(text), styles["CustomNormal"]))
                    story.append(Spacer(1, 6))
    log.debug(f"Finished render_html_elements. Story contains {len(story)} elements.")
    return story


def _extract_structured_content(payload: Any) -> list[Any] | None:
    if isinstance(payload, dict) and payload.get("type"):
        return [payload]
    if isinstance(payload, list):
        structured_hint = any(isinstance(item, dict) and item.get("type") for item in payload)
        if structured_hint:
            return payload
    return None


def render_structured_content(structured_blocks: list[Any]) -> list:
    renderer = _StructuredContentRenderer()
    return renderer.build(structured_blocks)


def flatten_structured_blocks(blocks: list[Any] | None, depth: int = 1) -> list[tuple[Any, int]]:
    if not blocks:
        return []
    flattened: list[tuple[Any, int]] = []
    for block in blocks:
        flattened.append((block, depth))
        if isinstance(block, dict):
            child_depth = depth + 1 if (block.get("type") or "").strip().lower() == "section" else depth
            children = block.get("children")
            if isinstance(children, list) and children:
                flattened.extend(flatten_structured_blocks(children, child_depth))
    return flattened


_PARAGRAPH_TYPES = {"paragraph", "text", "body", "description", "summary", "title"}


def _normalize_markup_text(value: Any, inline_only: bool = False) -> tuple[str, str]:
    if value is None:
        return "", ""
    text = str(value)
    html = markdown2.markdown(text, extras=["fenced-code-blocks"])
    soup = BeautifulSoup(html, "html.parser")
    if inline_only:
        # Strip block-level tags before the inline-only pass so that
        # list/paragraph HTML never reaches _apply_formatted_html_to_paragraph.
        # unwrap() keeps the tag's children in place, so text is preserved.
        for tag in soup.find_all(["p", "ul", "ol", "li", "blockquote", "pre", "div", "h1", "h2", "h3", "h4", "h5", "h6"]):
            tag.unwrap()
    allowed = {"strong", "b", "em", "i", "br"}
    for tag in soup.find_all():
        if tag.name not in allowed and tag.name != "body":
            tag.unwrap()
    body = soup.body
    formatted = ""
    if body:
        formatted = "".join(str(child) for child in body.children)
    else:
        formatted = str(soup)
    plain = soup.get_text(" ", strip=True)
    return plain, formatted


def _strip_wrapping_paragraph_tags(html: str) -> str:
    if not html:
        return ""
    trimmed = html.strip()
    stripped = re.sub(r"(?i)^<p>(.*?)</p>$", r"\1", trimmed, flags=re.S)
    return stripped.strip()


def _parse_paragraph_segments(raw_text: str) -> list[dict]:
    segments = []
    if not raw_text:
        return segments
    bullet_re = re.compile(r"^(?:[•\*-])\s+(.*)")
    label_re = re.compile(r"^\*\*(.+?):\*\*\s*(.*)")
    lines = [line.strip() for line in raw_text.splitlines() if line.strip()]
    for line in lines:
        bullet_match = bullet_re.match(line)
        if bullet_match:
            inner = bullet_match.group(1)
            plain, formatted = _normalize_markup_text(inner, inline_only=True)
            formatted = _strip_wrapping_paragraph_tags(formatted)
            if plain:
                segments.append({"type": "bullet", "text": plain, "formatted": formatted})
            continue
        label_match = label_re.match(line)
        if label_match:
            label_text = label_match.group(1).strip()
            label_value = f"{label_text}:"
            label_source = f"**{label_value}**"
            plain_label, formatted_label = _normalize_markup_text(label_source, inline_only=True)
            formatted_label = _strip_wrapping_paragraph_tags(formatted_label)
            rest = label_match.group(2).strip()
            rest_plain, rest_formatted = _normalize_markup_text(rest, inline_only=True)
            rest_formatted = _strip_wrapping_paragraph_tags(rest_formatted)
            if plain_label:
                segments.append({
                    "type": "label",
                    "label": label_value,
                    "formatted_label": formatted_label,
                    "text": rest_plain,
                    "formatted_text": rest_formatted,
                })
            elif rest_plain:
                segments.append({"type": "paragraph", "text": rest_plain, "formatted": rest_formatted})
            continue
        plain, formatted = _normalize_markup_text(line, inline_only=True)
        formatted = _strip_wrapping_paragraph_tags(formatted)
        if plain:
            segments.append({"type": "paragraph", "text": plain, "formatted": formatted})
    return segments


def _collect_nested_child_blocks(item: dict) -> list[Any]:
    nested: list[Any] = []
    for key in ("children", "content", "blocks"):
        value = item.get(key)
        if isinstance(value, dict):
            nested.append(value)
        elif isinstance(value, list):
            nested.extend(value)
    return nested


def _expand_paragraph_block(block: dict) -> list[dict]:
    segments = block.get("segments") or []
    child_blocks = block.get("children")
    if not segments:
        fallback = dict(block)
        fallback.pop("segments", None)
        return [fallback]

    semantic_blocks: list[dict] = []
    bullet_buffer: list[dict] = []

    def flush_bullets() -> None:
        nonlocal bullet_buffer
        if not bullet_buffer:
            return
        semantic_blocks.append({"type": "bullet_list", "items": [dict(item) for item in bullet_buffer]})
        bullet_buffer = []

    for segment in segments:
        seg_type = segment.get("type")
        if seg_type == "bullet":
            text = segment.get("text") or ""
            formatted = segment.get("formatted") or text
            if text:
                bullet_buffer.append({"text": text, "formatted": formatted})
            continue
        flush_bullets()
        if seg_type == "label":
            label_block: dict = {
                "type": "label_paragraph",
                "label": segment.get("label"),
                "formatted_label": segment.get("formatted_label"),
            }
            rest_text = segment.get("text")
            if rest_text:
                label_block["text"] = rest_text
            rest_formatted = segment.get("formatted_text")
            if rest_formatted:
                label_block["formatted"] = rest_formatted
            semantic_blocks.append(label_block)
            continue
        paragraph_text = segment.get("text")
        if paragraph_text:
            paragraph_block: dict = {"type": "paragraph", "text": paragraph_text}
            paragraph_formatted = segment.get("formatted")
            if paragraph_formatted:
                paragraph_block["formatted"] = paragraph_formatted
            semantic_blocks.append(paragraph_block)
    flush_bullets()

    if not semantic_blocks:
        fallback = dict(block)
        fallback.pop("segments", None)
        return [fallback]

    if child_blocks and semantic_blocks:
        semantic_blocks[-1]["children"] = child_blocks

    return semantic_blocks


def _finalize_normalized_block(block: dict) -> list[dict]:
    candidate = dict(block)
    block_type = (candidate.get("type") or "").strip().lower()
    if block_type == "table":                  # protect tables unconditionally first
        candidate.pop("segments", None)
        return [candidate]
    if block_type in _PARAGRAPH_TYPES:
        return _expand_paragraph_block(candidate)
    candidate.pop("segments", None)
    return [candidate]


def normalize_content_for_export(content: Any) -> list[dict]:
    def normalize_item(item: Any) -> list[dict]:
        if item is None:
            return []
        if isinstance(item, str):
            structured = _convert_markdown_to_structured(item)
            normalized_blocks: list[dict] = []
            for block in structured:
                normalized_block = dict(block)
                if normalized_block.get("text"):
                    raw_text = normalized_block["text"]
                    plain, formatted = _normalize_markup_text(raw_text)
                    normalized_block["text"] = plain
                    normalized_block["formatted"] = formatted
                if normalized_block["type"] in _PARAGRAPH_TYPES:
                    normalized_block["segments"] = _parse_paragraph_segments(raw_text)
                if normalized_block.get("type") == "table":
                    table_data = normalized_block.get("data") or []
                    normalized_block["data"] = _normalize_table_rows(table_data)
                normalized_blocks.extend(_finalize_normalized_block(normalized_block))
            return normalized_blocks
        if isinstance(item, dict):
            normalized: dict = {}
            item_type = (item.get("type") or "").strip().lower()

            if not item_type:
                if item.get("items"):
                    item_type = "list"
                elif item.get("cells") or item.get("rows") or item.get("table"):
                    item_type = "table"
                else:
                    item_type = "paragraph"

            normalized["type"] = item_type

            child_sources = _collect_nested_child_blocks(item)
            text_value = _extract_block_text(item)
            if item_type in _PARAGRAPH_TYPES and not text_value and child_sources:
                return normalize_list(child_sources)

            if item_type == "table":
                data = (
                    item.get("data")
                    or item.get("rows")
                    or item.get("cells")
                    or item.get("content")
                )
                normalized["data"] = _normalize_table_rows(data)
            else:
                raw_source = item.get("text") or item.get("title") or item.get("content") or ""
                normalized_plain, normalized_formatted = _normalize_markup_text(raw_source)
                if normalized_plain:
                    normalized["text"] = normalized_plain
                    normalized["formatted"] = normalized_formatted
                    normalized["raw_text"] = raw_source
                    if item_type in _PARAGRAPH_TYPES:
                        normalized["segments"] = _parse_paragraph_segments(raw_source)

            if child_sources:
                nested_children = normalize_list(child_sources)
                if nested_children:
                    normalized["children"] = nested_children

            if item_type == "list" and item.get("items"):
                normalized["items"] = []
                for list_entry in item.get("items") or []:
                    entry_plain, entry_formatted = _normalize_markup_text(list_entry)
                    if entry_plain:
                        normalized["items"].append({"text": entry_plain, "formatted": entry_formatted})
            # if item_type == "table":
            #     table_rows = item.get("data") or item.get("rows") or item.get("cells")
            #     normalized["data"] = _normalize_table_rows(table_rows)
            return _finalize_normalized_block(normalized)
        if isinstance(item, list):
            return normalize_list(item)
        plain, formatted = _normalize_markup_text(item)
        if not plain:
            return []
        fallback_block = {"type": "paragraph", "text": plain, "formatted": formatted}
        fallback_block["segments"] = _parse_paragraph_segments(plain)
        return _finalize_normalized_block(fallback_block)

    def normalize_list(items: list[Any] | None) -> list[dict]:
        if not isinstance(items, list):
            return []
        result: list[dict] = []
        for child in items:
            result.extend(normalize_item(child))
        return result

    return normalize_item(content)


def _normalize_table_rows(data: Any) -> list[list[str]]:
    rows: list[list[str]] = []
    if not isinstance(data, list):
        return rows
    for row in data:
        if isinstance(row, dict):
            cell_values = row.get("cells") or row.get("data") or row.get("row") or []
        elif isinstance(row, list):
            cell_values = row
        else:
            cell_values = [row]
        cells: list[str] = []
        for cell in cell_values:
            plain, _ = _normalize_markup_text(cell)
            cells.append(plain)
        if cells:
            rows.append(cells)
    return rows


class _StructuredContentRenderer:
    def __init__(self):
        self.story: list[Any] = []
        self.section_counters: list[int] = []

    def build(self, blocks: list[Any]) -> list[Any]:
        for block, depth in flatten_structured_blocks(blocks):
            self._render_block(block, depth)
        if not self.story:
            self.story.append(Paragraph("Empty Content", styles["CustomNormal"]))
        return self.story

    def _render_block(self, block: Any, section_depth: int) -> None:
        if isinstance(block, str):
            self._append_paragraph(block)
            return
        if not isinstance(block, dict):
            return

        block_type = (block.get("type") or "").strip().lower()

        if block_type == "title":
            self._render_title(block)
            return

        if block_type == "section":
            self._render_section(block, section_depth)
            return

        if block_type == "table":
            table_data = _normalize_table_rows(
                block.get("data")
                or block.get("rows")
                or block.get("cells")
                or block.get("content")
                or []
            )

            if table_data:
                table = ReportLabTable(table_data, repeatRows=1)

                table.setStyle(TableStyle([
                    ("GRID", (0, 0), (-1, -1), 0.75, colors.HexColor("#9AA0B6")),
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#E8EEF9")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#0A1F44")),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("ALIGN", (0, 0), (-1, 0), "CENTER"),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("LEFTPADDING", (0, 0), (-1, -1), 6),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                    ("TOPPADDING", (0, 0), (-1, -1), 4),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                ]))

                self.story.append(table)
                self.story.append(Spacer(1, 12))

            return

        if block_type == "label_paragraph":
            self._render_label_paragraph(block)
            return

        if block_type in {"paragraph", "text", "body", "description", "summary", "heading", "subheading"}:
            self._append_paragraph(block.get("text"), block.get("formatted"))
            return
        elif block_type in {"sources", "source", "references"}:
            self._render_sources(block)
            return
        elif "list" in block_type or "bullet" in block_type:
            list_flowable = self._build_list_flowable(block, depth=section_depth)
            if list_flowable is not None:
                self.story.append(list_flowable)
                self.story.append(Spacer(1, 10))
            return
        else:
            text = self._extract_text(block)
            self._append_paragraph(text, self._extract_formatted(block))


    def _render_title(self, block: dict) -> None:
        text = self._extract_text(block)
        if text:
            paragraph = Paragraph(render_text_with_emojis(text), styles["StructuredDocumentTitle"])
            self.story.append(paragraph)

    def _render_section(self, block: dict, depth: int) -> None:
        heading_text = self._extract_text(block, keys=("title", "text", "heading"))
        heading_text = heading_text or "Untitled Section"
        numbering = self._increment_section_counter(depth)
        omit_numbering = bool(re.match(r"^\s*\d+[\.)]", heading_text))
        display_text = heading_text if omit_numbering else (f"{numbering} {heading_text}" if numbering else heading_text)
        paragraph = Paragraph(render_text_with_emojis(display_text), self._heading_style(depth))
        self.story.append(paragraph)
        self.story.append(Spacer(1, 16))

    def _render_sources(self, block: dict) -> None:
        heading = block.get("title") or block.get("text") or "Sources"
        self.story.append(Paragraph(render_text_with_emojis(heading), styles["StructuredSourcesHeading"]))
        entries = self._extract_list_entries(block)
        if not entries:
            entries = block.get("children") or []
        for entry in entries:
            text = self._extract_text(entry)
            if text:
                self.story.append(Paragraph(render_text_with_emojis(text), styles["StructuredSourcesItem"]))
        self.story.append(Spacer(1, 6))

    def _build_list_flowable(self, block: dict, depth: int) -> ListFlowable | None:
        entries = self._extract_list_entries(block)
        if not entries:
            return None
        list_items: list[ListItem] = []
        for entry in entries:
            entry_flow: list[Any] = []
            entry_text = self._extract_text(entry)
            if entry_text:
                entry_formatted = self._extract_formatted(entry)
                entry_flow.append(Paragraph(render_text_with_emojis(entry_formatted or entry_text), styles["StructuredListItem"]))
            if isinstance(entry, dict):
                nested = self._extract_list_entries(entry)
                if nested:
                    nested_flowable = self._build_list_flowable({"items": nested, "ordered": entry.get("ordered")}, depth + 1)
                    if nested_flowable:
                        entry_flow.append(nested_flowable)
            if entry_flow:
                list_items.append(ListItem(entry_flow))
        if not list_items:
            return None
        ordered = self._determine_ordering(block)
        return ListFlowable(
            cast(Iterable, list_items),
            bulletType="1" if ordered else "bullet",
            leftIndent=8 * mm * depth,
            bulletIndent=4 * mm,
            spaceBefore=8,
            spaceAfter=8,
        )

    def _determine_ordering(self, block: dict) -> bool:
        ordered = block.get("ordered")
        if isinstance(ordered, str):
            return ordered.lower() in {"true", "1", "yes", "ordered", "ol"}
        if isinstance(ordered, bool):
            return ordered
        block_type = (block.get("type") or "").lower()
        return "ordered" in block_type or block_type in {"ol", "ordered_list"}

    def _extract_list_entries(self, block: dict) -> list[Any] | list:
        entries = block.get("items")
        if isinstance(entries, list) and entries:
            return entries
        entries = block.get("children")
        if isinstance(entries, list) and entries:
            return entries
        entries = block.get("entries")
        if isinstance(entries, list) and entries:
            return entries
        block_type = (block.get("type") or "").lower()
        if block_type in {"bullet", "list_item"}:
            text = block.get("text")
            if text:
                return [text]
        return []

    def _increment_section_counter(self, depth: int) -> str:
        if depth < 1:
            depth = 1
        while len(self.section_counters) < depth:
            self.section_counters.append(0)
        if len(self.section_counters) > depth:
            self.section_counters = self.section_counters[:depth]
        self.section_counters[depth - 1] += 1
        return ".".join(str(num) for num in self.section_counters if num > 0)

    def _heading_style(self, depth: int) -> ParagraphStyle:
        if depth <= 1:
            return styles["CustomHeading1"]
        if depth == 2:
            return styles["CustomHeading2"]
        return styles["CustomHeading3"]

    def _append_paragraph(self, text: str | None, formatted: str | None = None, style: ParagraphStyle | None = None) -> None:
        if not text and not formatted:
            return
        paragraph_style = style or styles["StructuredParagraph"]
        content = formatted or text or ""
        markup = self._prepare_inline_markup(content)
        paragraph = Paragraph(render_text_with_emojis(markup), paragraph_style)
        self.story.append(paragraph)

    def _prepare_inline_markup(self, content: str) -> str:
        if not content:
            return ""
        return (
            content.replace("<strong>", "<b>")
            .replace("</strong>", "</b>")
            .replace("<em>", "<i>")
            .replace("</em>", "</i>")
        )

    def _render_label_paragraph(self, block: dict) -> None:
        label = block.get("label")
        formatted_label = block.get("formatted_label")
        body = block.get("text")
        formatted_body = block.get("formatted")
        plain_parts = [part.strip() for part in (label, body) if part and str(part).strip()]
        formatted_parts = [part for part in (formatted_label, formatted_body) if part]
        combined_text = " ".join(plain_parts) if plain_parts else None
        combined_formatted = " ".join(formatted_parts) if formatted_parts else None
        self._append_paragraph(combined_text, combined_formatted)

    def _extract_text(self, block: Any, keys: tuple[str, ...] = ("text", "title", "content", "description", "label", "name")) -> str:
        if isinstance(block, str):
            return block.strip()
        if not isinstance(block, dict):
            return ""
        for key in keys:
            value = block.get(key)
            if value:
                return str(value).strip()
        return ""

    def _extract_formatted(self, block: Any) -> str | None:
        if isinstance(block, dict):
            return block.get("formatted")
        return None

def _cleanup_files(folder_path: str, delay_minutes: int):
    def delete_files():
        time.sleep(delay_minutes * 60)
        try:
            import shutil
            shutil.rmtree(folder_path) 
            log.debug(f"Folder {folder_path} deleted.")
        except Exception as e:
            logging.error(f"Error deleting files : {e}")
    thread = threading.Thread(target=delete_files)
    thread.start()

def _convert_markdown_to_structured(markdown_content):
    """
    Converts Markdown content into a structured format for Word
    
    Args:
        markdown_content (str): Markdown content
        
    Returns:
        list: List of objects with 'text' and 'type'
    """
    if not markdown_content or not isinstance(markdown_content, str):
        return []
    
    lines = markdown_content.splitlines()
    structured = []
    i = 0

    def parse_table_row(row_line: str) -> list[str]:
        cleaned = row_line.strip().strip("|")
        return [cell.strip() for cell in cleaned.split("|")] if cleaned else []

    while i < len(lines):
        raw_line = lines[i]
        line = raw_line.strip()
        if not line:
            i += 1
            continue

        if line.startswith('|') and i + 1 < len(lines) and TABLE_SEPARATOR_RE.match(lines[i + 1].strip()):
            rows = [parse_table_row(line)]
            i += 2
            while i < len(lines):
                next_line = lines[i].strip()
                if not next_line or '|' not in next_line:
                    break
                rows.append(parse_table_row(next_line))
                i += 1
            structured.append({"type": "table", "data": rows})
            continue

        if line.startswith('# '):
            structured.append({"text": line[2:].strip(), "type": "title"})
        elif line.startswith('## '):
            structured.append({"text": line[3:].strip(), "type": "heading"})
        elif line.startswith('### '):
            structured.append({"text": line[4:].strip(), "type": "subheading"})
        elif line.startswith('#### '):
            structured.append({"text": line[5:].strip(), "type": "subheading"})
        elif line.startswith('- '):
            structured.append({"text": line[2:].strip(), "type": "bullet"})
        elif line.startswith('* '):
            structured.append({"text": line[2:].strip(), "type": "bullet"})
        elif line.startswith('**') and line.endswith('**'):
            structured.append({"text": line[2:-2].strip(), "type": "bold"})
        else:
            structured.append({"text": line, "type": "paragraph"})
        i += 1
    
    return structured


def convert_structured_to_markdown(content: List[dict]) -> str:
    """Render structured content into Markdown before PDF generation."""
    if not isinstance(content, list):
        return ""

    parts: list[str] = []

    def append_paragraph(text: str) -> None:
        normalized = (text or "").strip()
        if normalized:
            parts.append(normalized)

    def handle_item(item: Any) -> None:
        if isinstance(item, str):
            append_paragraph(item)
            return
        if isinstance(item, list):
            for child in item:
                handle_item(child)
            return
        if not isinstance(item, dict):
            return

        item_type = (item.get("type") or "").strip().lower()
        text_value = (item.get("text") or item.get("title") or "").strip()

        if item_type == "title":
            append_paragraph(f"# {text_value}")
        elif item_type == "section":
            if text_value:
                append_paragraph(f"## {text_value}")
        elif item_type in {"paragraph", "text", "source", "sources", "description"}:
            append_paragraph(text_value)
        elif text_value:
            append_paragraph(text_value)

        children = item.get("children")
        if isinstance(children, list):
            for child in children:
                handle_item(child)

    for entry in content:
        handle_item(entry)

    return "\n\n".join(parts)

def _create_excel(data: list[list[str]], filename: str, folder_path: str | None = None, title: str | None = None) -> dict:
    log.debug("Creating Excel file with optional template")
    if folder_path is None:
        folder_path = _generate_unique_folder()
    
    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "xlsx")

    if XLSX_TEMPLATE:
        try:
            log.debug("Loading XLSX template...")
            wb = load_workbook(XLSX_TEMPLATE_PATH) 
            log.debug(f"Template loaded with {len(wb.sheetnames)} sheet(s)")
        except Exception as e:
            log.warning(f"Failed to load XLSX template: {e}")
            wb = Workbook()
    else:
        log.debug("No XLSX template available, creating new workbook")
        wb = Workbook()

    ws = wb.active

    from openpyxl.utils import get_column_letter 

    
    if title:
        ws.title = "".join(c for c in title if c.isalnum() or c in (' ', '-', '_')).strip()[:31]
        title_cell_found = False
        for row in ws.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str) and "title" in cell.value.lower():
                    cell.value = title
                    from openpyxl.styles import Font
                    log.debug(f"Title '{title}' replaced in cell {get_column_letter(cell.column)}{cell.row} containing 'title'")
                    title_cell_found = True
                    break
            if title_cell_found:
                break
    
    start_row, start_col = 1, 1
    if ws.auto_filter and ws.auto_filter.ref:
        try:
            from openpyxl.utils import range_boundaries
            start_col, start_row, _, _ = range_boundaries(ws.auto_filter.ref)
        except: pass

    if not data:
        wb.save(filepath)
        return {"success": True, "filepath": filepath, "filename": filename}

    template_border = ws.cell(start_row, start_col).border
    has_borders = template_border and any([template_border.top.style, template_border.bottom.style, 
                                          template_border.left.style, template_border.right.style])
    
    for r in range(max(len(data) + 10, 50)):
        for c in range(max(len(data[0]) + 5, 20)):
            cell = ws.cell(row=start_row + r, column=start_col + c)
            
            if r < len(data) and c < len(data[0]):
                cell.value = data[r][c]
                if r == 0 and data[r][c]:  
                    from openpyxl.styles import Font
                    cell.font = Font(bold=True)
                if has_borders:  
                    from openpyxl.styles import Border
                    cell.border = Border(top=template_border.top, bottom=template_border.bottom,
                                       left=template_border.left, right=template_border.right)
            else:
                cell.value = None
                if cell.has_style:
                    from openpyxl.styles import Font, PatternFill, Border, Alignment
                    cell.font, cell.fill, cell.border, cell.alignment = Font(), PatternFill(), Border(), Alignment()

    if ws.auto_filter:
        ws.auto_filter.ref = f"{get_column_letter(start_col)}{start_row}:{get_column_letter(start_col + len(data[0]) - 1)}{start_row + len(data) - 1}"
    
    for c in range(len(data[0])):
        max_len = max(len(str(data[r][c])) for r in range(len(data)))
        ws.column_dimensions[get_column_letter(start_col + c)].width = min(max_len + 2, 150)

    wb.save(filepath)

    return {"url": _public_url(folder_path, fname), "path": filepath}
def _create_csv(data: list[list[str]], filename: str, folder_path: str | None = None) -> dict:
    log.debug("Creating CSV file")
    if folder_path is None:
        folder_path = _generate_unique_folder()

    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "csv")

    with open(filepath, "w", newline="", encoding="utf-8") as f:
        if isinstance(data, list):
            csv.writer(f).writerows(data)
        else:
            csv.writer(f).writerow([data])

    return {"url": _public_url(folder_path, fname), "path": filepath}

def _build_markdown_story(source: Any) -> list:
    md_parts: list[str] = []

    def append_paragraph(value: Any) -> None:
        normalized = (value or "").strip()
        if normalized:
            md_parts.append(normalized)

    items = source if isinstance(source, list) else [source]
    for entry in items:
        if isinstance(entry, str):
            append_paragraph(entry)
            continue
        if isinstance(entry, dict):
            t = (entry.get("type") or "").lower()
            text_value = entry.get("text") or entry.get("title") or ""
            if t == "title":
                append_paragraph(f"# {text_value}")
            elif t == "subtitle":
                append_paragraph(f"## {text_value}")
            elif t == "paragraph":
                append_paragraph(entry.get("text", ""))
            elif t == "list":
                append_paragraph("\n".join([f"- {x}" for x in entry.get("items", [])]))
            elif t in ("image", "image_query"):
                query = entry.get("query", "")
                if query:
                    append_paragraph(f"![Image](image_query: {query})")
            else:
                append_paragraph(text_value)
        else:
            append_paragraph(str(entry))

    md_text = "\n\n".join(md_parts) or ""

    def replace_image_query(match: re.Match) -> str:
        query = match.group(1).strip()
        image_url = search_image(query)
        if image_url:
            return f'\n\n<img src="{image_url}" alt="Image: {query}" />\n\n'
        return ""

    md_text = re.sub(r'!\[[^\]]*\]\(\s*image_query:\s*([^)]+)\)', replace_image_query, md_text)
    html = markdown2.markdown(md_text, extras=['fenced-code-blocks','tables','break-on-newline','cuddled-lists'])
    soup = BeautifulSoup(html, "html.parser")
    story = render_html_elements(soup) or [Paragraph("Empty Content", styles["CustomNormal"])]
    return story


def _create_pdf(content: Any, filename: str, folder_path: str | None = None) -> dict:    
    log.debug("Creating PDF file")
    if folder_path is None:
        folder_path = _generate_unique_folder()
    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "pdf")

    normalized_blocks = normalize_content_for_export(content)
    if normalized_blocks:
        story = render_structured_content(normalized_blocks)
    else:
        story = _build_markdown_story(content)

    doc = SimpleDocTemplate(filepath, topMargin=54, bottomMargin=54, leftMargin=54, rightMargin=54)
    try:
        doc.build(story)
    except Exception as e:
        log.error(f"Error building PDF {fname}: {e}", exc_info=True)
        fallback = SimpleDocTemplate(filepath)
        fallback.build([Paragraph("Error in PDF generation", styles["CustomNormal"])])

    return {"url": _public_url(folder_path, fname), "path": filepath}

def _create_presentation(slides_data: list[dict], filename: str, folder_path: str | None = None, title: str | None = None) -> dict:
    if folder_path is None:
        folder_path = _generate_unique_folder()
    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "pptx")
      
    use_template = False
    prs = None
    title_layout = None
    content_layout = None

    if PPTX_TEMPLATE:
        try:
            log.debug("Attempting to load template...")
            src = PPTX_TEMPLATE
            if hasattr(PPTX_TEMPLATE, "slides") and hasattr(PPTX_TEMPLATE, "save"):
                log.debug("Template is a Presentation object, converting to BytesIO")
                buf = BytesIO()
                PPTX_TEMPLATE.save(buf); buf.seek(0)
                src = buf

            tmp = Presentation(src)
            log.debug(f"Template loaded with {len(tmp.slides)} slides")
            if len(tmp.slides) >= 1:
                prs = tmp
                use_template = True

                title_layout = prs.slides[0].slide_layout
                content_layout = prs.slides[1].slide_layout if len(prs.slides) >= 2 else prs.slides[0].slide_layout
                log.debug("Using template layouts")

                for i in range(len(prs.slides) - 1, 0, -1):
                    rId = prs.slides._sldIdLst[i].rId 
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[i]
        except Exception:
            log.error(f"Error loading template: {e}")
            use_template = False
            prs = None

    if not use_template:
        log.debug("No valid template, creating new presentation with default layouts")
        prs = Presentation()
        title_layout = prs.slide_layouts[0]
        content_layout = prs.slide_layouts[1]

    if use_template:
        log.debug("Using template title slide")
        tslide = prs.slides[0]
        if tslide.shapes.title:
            tslide.shapes.title.text = title or ""
            for p in tslide.shapes.title.text_frame.paragraphs:
                for r in p.runs:
                    title_info = next(({'size': PptPt(int(child.attrib.get('sz', 2800))/100), 'bold': child.attrib.get('b', '0') == '1'} for child in title_layout.element.iter() if 'defRPr' in child.tag.split('}')[-1] and 'sz' in child.attrib), {'size': PptPt(28), 'bold': True})

                    r.font.size = title_info['size'] 
                    r.font.bold = title_info['bold']
    else:
        log.debug("Creating new title slide")
        tslide = prs.slides.add_slide(title_layout)
        if tslide.shapes.title:
            tslide.shapes.title.text = title or ""
            for p in tslide.shapes.title.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = PptPt(28); r.font.bold = True

    EMU_PER_IN = 914400
    slide_w_in = prs.slide_width / EMU_PER_IN
    slide_h_in = prs.slide_height / EMU_PER_IN
    log.debug(f"Slide dimensions: {slide_w_in} x {slide_h_in} inches")

    page_margin = 0.5
    gutter = 0.3

    for i, slide_data in enumerate(slides_data):
        log.debug(f"Processing slide {i+1}: {slide_data.get('title', 'Untitled')}")
        if not isinstance(slide_data, dict):
            log.warning(f"Slide data is not a dict, skipping slide {i+1}")
            continue

        slide_title = slide_data.get("title", "Untitled")
        content_list = slide_data.get("content", [])
        if not isinstance(content_list, list):
            content_list = [content_list]
        log.debug(f"Adding slide with title: '{slide_title}'")
        slide = prs.slides.add_slide(content_layout)

        if slide.shapes.title:
            slide.shapes.title.text = slide_title
            for p in slide.shapes.title.text_frame.paragraphs:
                for r in p.runs:
                    title_info = next(({'size': PptPt(int(child.attrib.get('sz', 2800))/100), 'bold': child.attrib.get('b', '0') == '1'} for child in content_layout.element.iter() if 'defRPr' in child.tag.split('}')[-1] and 'sz' in child.attrib), {'size': PptPt(28), 'bold': True})

                    r.font.size = title_info['size'] 
                    r.font.bold = title_info['bold']

        content_shape = None
        try:
            for ph in slide.placeholders:
                try:
                    if ph.placeholder_format.idx == 1:
                        content_shape = ph; break
                except Exception:
                    pass
            if content_shape is None:
                for ph in slide.placeholders:
                    try:
                        if ph.placeholder_format.idx != 0:
                            content_shape = ph; break
                    except Exception:
                        pass
        except Exception:
            log.error(f"Error finding content placeholder: {e}")
            pass

        title_bottom_in = 1.0 
        if slide.shapes.title:
            try:
                title_bottom_emu = slide.shapes.title.top + slide.shapes.title.height
                title_bottom_in = max(title_bottom_emu / EMU_PER_IN, 1.0)
                title_bottom_in += 0.2
            except Exception:
                title_bottom_in = 1.2 

        if content_shape is None:

            content_shape = slide.shapes.add_textbox(Inches(page_margin), Inches(title_bottom_in), Inches(slide_w_in - 2*page_margin), Inches(slide_h_in - title_bottom_in - page_margin))
            log.debug("Creating new textbox for content")
        tf = content_shape.text_frame
        try:
            tf.clear()
        except Exception:
            log.error(f"Error clearing text frame: {e}")
            pass
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        try:
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.05)
            tf.margin_bottom = Inches(0.05)
        except Exception:
            pass

        content_left_in, content_top_in = page_margin, title_bottom_in
        content_width_in = slide_w_in - 2*page_margin
        content_height_in = slide_h_in - (title_bottom_in + page_margin)

        image_query = slide_data.get("image_query")
        if image_query:
            image_url = search_image(image_query)
            if image_url:
                log.debug(f"Searching for image query: '{image_query}'")
                try:
                    log.debug(f"Downloading image from URL: {image_url}")
                    response = requests.get(image_url, timeout=30)
                    response.raise_for_status()
                    image_data = response.content
                    image_stream = BytesIO(image_data)
                    pos = slide_data.get("image_position", "right")
                    size = slide_data.get("image_size", "medium")
                    if size == "small":
                        img_w_in, img_h_in = 2.0, 1.5
                    elif size == "large":
                        img_w_in, img_h_in = 4.0, 3.0
                    else:
                        img_w_in, img_h_in = 3.0, 2.0
                    log.debug(f"Image dimensions: {img_w_in} x {img_h_in} inches")

                    if pos == "left":
                        img_left_in = page_margin
                        img_top_in = title_bottom_in
                        content_left_in = img_left_in + img_w_in + gutter
                        content_top_in = title_bottom_in
                        content_width_in = max(slide_w_in - page_margin - content_left_in, 2.5)
                        content_height_in = slide_h_in - (title_bottom_in + page_margin)
                    elif pos == "right":
                        img_left_in = max(slide_w_in - page_margin - img_w_in, page_margin)
                        img_top_in = title_bottom_in
                        content_left_in = page_margin
                        content_top_in = title_bottom_in
                        content_width_in = max(img_left_in - gutter - content_left_in, 2.5)
                        content_height_in = slide_h_in - (title_bottom_in + page_margin)
                    elif pos == "top":
                        img_left_in = slide_w_in - page_margin - img_w_in
                        img_top_in = title_bottom_in
                        content_left_in = page_margin
                        content_top_in = img_top_in + img_h_in + gutter
                        content_width_in = slide_w_in - 2*page_margin
                        content_height_in = max(slide_h_in - page_margin - content_top_in, 2.0)
                    elif pos == "bottom":
                        img_left_in = slide_w_in - page_margin - img_w_in
                        img_top_in = max(slide_h_in - page_margin - img_h_in, page_margin)
                        content_left_in = page_margin
                        content_top_in = title_bottom_in
                        content_width_in = slide_w_in - 2*page_margin
                        content_height_in = max(img_top_in - gutter - content_top_in, 2.0)
                    else:
                        img_left_in = max(slide_w_in - page_margin - img_w_in, page_margin)
                        img_top_in = title_bottom_in
                        content_left_in = page_margin
                        content_top_in = title_bottom_in
                        content_width_in = max(img_left_in - gutter - content_left_in, 2.5)
                        content_height_in = slide_h_in - (title_bottom_in + page_margin)

                    slide.shapes.add_picture(image_stream, Inches(img_left_in), Inches(img_top_in), Inches(img_w_in), Inches(img_h_in))
                    log.debug(f"Image added at position: left={img_left_in}, top={img_top_in}")
                except Exception:
                    pass

        try:
            content_shape.left = Inches(content_left_in)
            content_shape.top = Inches(content_top_in)
            content_shape.width = Inches(content_width_in)
            content_shape.height = Inches(content_height_in)
        except Exception:
            pass

        approx_chars_per_in = 9.5
        approx_lines_per_in = 1.6
        safe_width = max(content_width_in, 0.1)
        safe_height = max(content_height_in, 0.1)
        est_capacity = int(safe_width * approx_chars_per_in * safe_height * approx_lines_per_in)
        font_size = dynamic_font_size(content_list, max_chars=max(est_capacity, 120), base_size=24, min_size=12)

        try:
            tf = content_shape.text_frame
        except Exception:
            try:
                tf = content_shape.text_frame
            except Exception:
                log.warning("Could not access text frame for content shape")
                continue

        if not tf.paragraphs:
            tf.add_paragraph()
        for idx, line in enumerate(content_list):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = str(line) if line is not None else ""
            run.font.size = font_size
            p.space_after = PptPt(6)

    prs.save(filepath)
    return {"url": _public_url(folder_path, fname), "path": filepath}

def _create_word(content: list[dict] | str, filename: str, folder_path: str | None = None, title: str | None = None) -> dict:
    log.debug("Creating Word document")

    if isinstance(content, str):
        content = _convert_markdown_to_structured(content)
    elif isinstance(content, dict):
        content = [content]
    elif not isinstance(content, list):
        content = []

    if folder_path is None:
        folder_path = _generate_unique_folder()
    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "docx")

    use_template = False
    doc = None

    if DOCX_TEMPLATE:
        try:
            src = DOCX_TEMPLATE
            if hasattr(DOCX_TEMPLATE, "paragraphs") and hasattr(DOCX_TEMPLATE, "save"):
                buf = BytesIO()
                DOCX_TEMPLATE.save(buf)
                buf.seek(0)
                src = buf

            doc = Document(src)
            use_template = True
            log.debug("Using DOCX template")

            for element in doc.element.body:
                if element.tag.endswith('}p') or element.tag.endswith('}tbl'):
                    doc.element.body.remove(element)

        except Exception as e:
            log.warning(f"Failed to load DOCX template: {e}")
            use_template = False
            doc = None

    if not use_template:
        doc = Document()
        log.debug("Creating new Word document without template")

    if title:
        title_paragraph = doc.add_paragraph(title)
        try:
            title_paragraph.style = doc.styles['Title']
        except KeyError:
            try:
                title_paragraph.style = doc.styles['Heading 1']
            except KeyError:
                run = title_paragraph.runs[0] if title_paragraph.runs else title_paragraph.add_run()
                run.font.size = DocxPt(20)
                run.font.bold = True
        title_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        log.debug("Document title added")

    normalized_content = normalize_content_for_export(content)
    if normalized_content:
        _render_structured_docx(doc, normalized_content)
    else:
        _add_docx_paragraph(doc, "")

    doc.save(filepath)
    return {"url": _public_url(folder_path, fname), "path": filepath}


def _render_structured_docx(doc: Document, blocks: list[Any], depth: int = 1) -> None:
    for block in blocks:
        # ── plain string ──────────────────────────────────────────────────────
        if isinstance(block, str):
            _add_docx_paragraph(doc, block)
            continue
        if not isinstance(block, dict):
            continue

        block_type = (block.get("type") or "").strip().lower()
        text = _extract_block_text(block)

        # ── table: first — never allowed to fall through ──────────────────────
        if block_type == "table":
            table_data = _normalize_table_rows(
                block.get("data")
                or block.get("rows")
                or block.get("cells")
                or block.get("content")
                or []
            )

            if table_data:
                cols = max((len(row) for row in table_data), default=0)
                if cols > 0:
                    table = doc.add_table(rows=len(table_data), cols=cols)
                    try:
                        table.style = "Table Grid"
                    except Exception:
                        pass

                    for i, row in enumerate(table_data):
                        for j in range(cols):
                            cell_text = str(row[j]) if j < len(row) else ""
                            cell = table.cell(i, j)
                            cell.text = cell_text

                            for para in cell.paragraphs:
                                para.paragraph_format.space_before = DocxPt(0)
                                para.paragraph_format.space_after = DocxPt(0)

                            if i == 0:
                                for para in cell.paragraphs:
                                    para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                                    for run in para.runs:
                                        run.bold = True

                    doc.add_paragraph()

            continue

        # ── title ─────────────────────────────────────────────────────────────
        if block_type == "title":
            _add_docx_heading(doc, text, level=1, centered=True)
            _render_structured_docx(doc, block.get("children") or [], depth)
            continue

        # ── section ───────────────────────────────────────────────────────────
        if block_type == "section":
            _add_docx_heading(doc, text or block.get("title") or "Section", level=min(depth + 1, 3))
            _render_structured_docx(doc, block.get("children") or [], depth + 1)
            continue

        # ── paragraph-like ────────────────────────────────────────────────────
        if block_type in {"paragraph", "text", "body", "description", "summary"}:
            if text:
                _add_docx_paragraph(doc, text, block.get("formatted"))
            continue

        # ── label paragraph ───────────────────────────────────────────────────
        if block_type == "label_paragraph":
            label_html = block.get("formatted_label") or block.get("label")
            body_html = block.get("formatted") or block.get("text")
            formatted_content = " ".join(part for part in (label_html, body_html) if part)
            plain_content = " ".join(
                part.strip() for part in (block.get("label"), block.get("text"))
                if part and str(part).strip()
            )
            _add_docx_paragraph(doc, plain_content, formatted_content or None)
            continue

        # ── sources / references ──────────────────────────────────────────────
        if block_type in {"sources", "source", "references"}:
            _add_docx_heading(doc, text or "Sources", level=min(depth + 1, 3))
            entries = block.get("children") or block.get("items") or []
            for entry in entries:
                entry_text = _extract_block_text(entry)
                entry_formatted = entry.get("formatted") if isinstance(entry, dict) else None
                if entry_text:
                    para = doc.add_paragraph()
                    _apply_formatted_html_to_paragraph(para, entry_formatted, entry_text)
                    para.paragraph_format.space_before = DocxPt(2)
                    para.paragraph_format.space_after = DocxPt(2)
                    para.paragraph_format.left_indent = DocxPt(12)
                    run = para.runs[0] if para.runs else para.add_run()
                    run.italic = True
            continue

        # ── list / bullet ─────────────────────────────────────────────────────
        if "list" in block_type or "bullet" in block_type:
            entries = block.get("items") or block.get("children") or block.get("entries") or []
            if not entries and block_type in {"bullet", "list_item"}:
                text_value = block.get("text")
                if text_value:
                    entries = [text_value]
            for entry in entries:
                entry_text = _extract_block_text(entry)
                entry_formatted = entry.get("formatted") if isinstance(entry, dict) else None
                if not entry_text:
                    continue
                para = doc.add_paragraph()
                try:
                    para.style = doc.styles["List Bullet"]
                except KeyError:
                    para.style = doc.styles["Normal"]
                para.paragraph_format.left_indent = DocxPt(12 * depth)
                para.paragraph_format.space_before = DocxPt(2)
                para.paragraph_format.space_after = DocxPt(2)
                _apply_formatted_html_to_paragraph(para, entry_formatted, entry_text)
                if isinstance(entry, dict) and entry.get("children"):
                    _render_structured_docx(doc, entry.get("children"), depth + 1)
            continue

        # ── unknown type ──────────────────────────────────────────────────────
        # Reached only for block types not listed above. Safe to render text
        # and recurse — cannot fire for any handled type because every branch
        # above ends with `continue`.
        if text:
            _add_docx_paragraph(doc, text)
        if block.get("children"):
            _render_structured_docx(doc, block.get("children"), depth + 1)
def _add_docx_paragraph(doc: Document, text: str, formatted: str | None = None) -> None:
    if not text and not formatted:
        return
    paragraph = doc.add_paragraph()
    _apply_formatted_html_to_paragraph(paragraph, formatted, text)
    try:
        paragraph.style = doc.styles['Normal']
    except KeyError:
        pass
    paragraph.paragraph_format.space_before = DocxPt(4)
    paragraph.paragraph_format.space_after = DocxPt(6)


def _apply_formatted_html_to_paragraph(paragraph, formatted_html: str | None, fallback_text: str) -> None:
    while paragraph.runs:
        paragraph.runs[0]._element.getparent().remove(paragraph.runs[0]._element)
    content = formatted_html or fallback_text or ""
    if not content:
        return
    soup = BeautifulSoup(content, "html.parser")

    def recurse(node, bold=False, italic=False):
        for child in node.children:
            if isinstance(child, str):
                if child:
                    run = paragraph.add_run(child)
                    run.bold = bold
                    run.italic = italic
            elif child.name in {"strong", "b"}:
                recurse(child, bold=True, italic=italic)
            elif child.name in {"em", "i"}:
                recurse(child, bold=bold, italic=True)
            elif child.name == "br":
                paragraph.add_run().add_break()
            else:
                recurse(child, bold=bold, italic=italic)

    recurse(soup)


def _add_docx_heading(doc: Document, text: str, level: int = 2, centered: bool = False) -> None:
    if not text:
        return
    paragraph = doc.add_paragraph(text)
    style_name = f"Heading {min(max(level, 1), 3)}"
    try:
        paragraph.style = doc.styles[style_name]
    except KeyError:
        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        run.font.size = DocxPt(18 if level <= 2 else 14)
        run.font.bold = True
    if centered:
        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    paragraph.paragraph_format.space_before = DocxPt(6)
    paragraph.paragraph_format.space_after = DocxPt(6)


def _extract_document_title(content: Any, root_title: str | None) -> str | None:
    """
    Implement block-first strategy for document title extraction:
    
    1. If a block with type="title" exists in content:
       - Use ONLY that block as the document title
       - IGNORE the root-level "title" field completely
    
    2. If NO block with type="title" exists:
       - Use the root-level "title" field as the document title
    
    Returns:
        The title text to use, or None if no title should be rendered
    """
    # Normalize content to a list
    content_list: list[Any] = []
    if isinstance(content, dict):
        content_list = [content]
    elif isinstance(content, list):
        content_list = content
    
    # Check for title blocks
    has_title_block = False
    title_text_from_block: str | None = None
    
    for item in content_list:
        if isinstance(item, dict):
            item_type = (item.get("type") or "").strip().lower()
            if item_type == "title":
                has_title_block = True
                # Get text from the title block
                text = item.get("text") or item.get("title") or ""
                if text:
                    title_text_from_block = str(text).strip()
                    break
        elif isinstance(item, str):
            # For string content, check if it's markdown with # title
            if item.strip().startswith("# "):
                has_title_block = True
                title_text_from_block = item.strip()[2:].strip()
                break
    
    # Apply block-first strategy
    if has_title_block:
        # Use only the title block, ignore root_title
        log.debug(f"Title block found, using block title: '{title_text_from_block}'")
        return title_text_from_block
    
    # No title block found, use root-level title
    if root_title:
        log.debug(f"No title block found, using root title: '{root_title}'")
        return root_title.strip()
    
    return None


def _extract_block_text(block: Any) -> str:
    if isinstance(block, str):
        return block.strip()
    if not isinstance(block, dict):
        return ""

    block_type = (block.get("type") or "").strip().lower()
    if block_type == "table":
        return ""

    for key in ("text", "title", "description", "label", "name"):
        value = block.get(key)
        if value:
            return str(value).strip()
    return ""

def _create_raw_file(content: str, filename: str | None, folder_path: str | None = None) -> dict:
    log.debug("Creating raw file")
    if folder_path is None:
        folder_path = _generate_unique_folder()

    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "txt")

    if fname.lower().endswith(".xml") and isinstance(content, str) and not content.strip().startswith("<?xml"):
        content = f'<?xml version="1.0" encoding="UTF-8"?>\n{content}'

    with open(filepath, "w", encoding="utf-8") as f:
        f.write(content or "")

    return {"url": _public_url(folder_path, fname), "path": filepath}

def upload_file(file_path: str, filename: str, file_type: str, token: str) -> dict:
    """
    Upload a file to OpenWebUI server.
    """
    url = f"{URL}/api/v1/files/"
    headers = {
        'Authorization': token,
        'Accept': 'application/json'
    }
    
    with open(file_path, 'rb') as f:
        files = {'file': f}
        response = post(url, headers=headers, files=files)

    if response.status_code != 200:
        return {"error": {"message": f'Error uploading file: {response.status_code}'}}
    else:
        return {
            "file_path_download": f"[Download {filename}.{file_type}](/api/v1/files/{response.json()['id']}/content)"
        }

def download_file(file_id: str, token: str) -> BytesIO:
    """
    Download a file from OpenWebUI server.
    """
    url = f"{URL}/api/v1/files/{file_id}/content"
    headers = {
        'Authorization': token,
        'Accept': 'application/json'
    }
    
    response = get(url, headers=headers)
    
    if response.status_code != 200:
        return {"error": {"message": f'Error downloading the file: {response.status_code}'}}
    else:
        return BytesIO(response._content)
    
def _extract_paragraph_style_info(para):
    """Extrait les informations de style détaillées d'un paragraphe"""
    if not para.runs:
        return {}
    
    first_run = para.runs[0]
    return {
        "font_name": first_run.font.name,
        "font_size": first_run.font.size,
        "bold": first_run.font.bold,
        "italic": first_run.font.italic,
        "underline": first_run.font.underline,
        "color": first_run.font.color.rgb if first_run.font.color else None
    }

def _extract_cell_style_info(cell):
    """Extrait les informations de style d'une cellule"""
    return {
        "style": cell.style.name if hasattr(cell, 'style') else None,
        "text_alignment": cell.paragraphs[0].alignment if cell.paragraphs else None
    }

def _apply_text_to_paragraph(para, new_text):
    """
    Apply new text to a paragraph while preserving formatting.
    """
    original_style = para.style
    original_alignment = para.alignment
    
    original_run_format = None
    if para.runs:
        first_run = para.runs[0]
        original_run_format = {
            "font_name": first_run.font.name,
            "font_size": first_run.font.size,
            "bold": first_run.font.bold,
            "italic": first_run.font.italic,
            "underline": first_run.font.underline,
            "color": first_run.font.color.rgb if first_run.font.color and first_run.font.color.rgb else None
        }
    
    for _ in range(len(para.runs)):
        para._element.remove(para.runs[0]._element)
    
    if isinstance(new_text, list):
        for i, text_item in enumerate(new_text):
            if i > 0:
                para.add_run("\n")
            run = para.add_run(str(text_item))
            if original_run_format:
                _apply_run_formatting(run, original_run_format)
    else:
        run = para.add_run(str(new_text))
        if original_run_format:
            _apply_run_formatting(run, original_run_format)
    
    if original_style:
        try:
            para.style = original_style
        except Exception:
            pass
    if original_alignment is not None:
        try:
            para.alignment = original_alignment
        except Exception:
            pass


def _apply_run_formatting(run, format_dict):
    """
    Apply formatting from a dict to a run.
    """
    try:
        if format_dict.get("font_name"):
            run.font.name = format_dict["font_name"]
    except Exception:
        pass
    
    try:
        if format_dict.get("font_size"):
            run.font.size = format_dict["font_size"]
    except Exception:
        pass
    
    try:
        if format_dict.get("bold") is not None:
            run.font.bold = format_dict["bold"]
    except Exception:
        pass
    
    try:
        if format_dict.get("italic") is not None:
            run.font.italic = format_dict["italic"]
    except Exception:
        pass
    
    try:
        if format_dict.get("underline") is not None:
            run.font.underline = format_dict["underline"]
    except Exception:
        pass
    
    try:
        if format_dict.get("color"):
            from docx.shared import RGBColor
            run.font.color.rgb = format_dict["color"]
    except Exception:
        pass

@mcp.tool(
    name="full_context_document",
    title="Return the structure of a document (docx, xlsx, pptx)",
    description="Return the structure, content, and metadata of a document based on its type (docx, xlsx, pptx). Unified output format with index, type, style, and text."
)

async def full_context_document(
    file_id: str,
    file_name: str,
    ctx: Context[ServerSession, None]
) -> dict:
    """
    Return the structure of a document (docx, xlsx, pptx) based on its file extension.
    The function detects the file type and processes it accordingly.
    Returns:
        dict: A JSON object with the structure of the document.
    """
    try:
        bearer_token = ctx.request_context.request.headers.get("authorization")
        logging.info(f"Recieved authorization header!")
        user_token=bearer_token
    except:
        logging.error(f"Error retrieving authorization header use admin fallback")
        user_token=TOKEN
    try:
        user_file = download_file(file_id,token=user_token)

        if isinstance(user_file, dict) and "error" in user_file:
            return json.dumps(user_file, indent=4, ensure_ascii=False)

        file_extension = os.path.splitext(file_name)[1].lower()
        file_type = file_extension.lstrip('.')

        structure = {
            "file_name": file_name,
            "file_id": file_id,
            "type": file_type,
            "slide_id_order": [],
            "body": [],
        }
        index_counter = 1

        if file_type == "docx":
            doc = Document(user_file)
            
            para_id_counter = 1
            
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                
                style = para.style.name
                style_info = _extract_paragraph_style_info(para)
                element_type = "heading" if style.startswith("Heading") else "paragraph"
                
                para_xml_id = para_id_counter
                
                structure["body"].append({
                    "index": para_id_counter,
                    "para_xml_id": para_xml_id,
                    "id_key": f"pid:{para_xml_id}",
                    "type": element_type,
                    "style": style,
                    "style_info": style_info,
                    "text": text
                })
                para_id_counter += 1
            
            for table_idx, table in enumerate(doc.tables):
                table_xml_id = id(table._element)
                table_info = {
                    "index": para_id_counter,
                    "table_xml_id": table_xml_id,
                    "id_key": f"tid:{table_xml_id}",
                    "type": "table",
                    "style": "Table",
                    "table_id": table_idx,
                    "rows": len(table.rows),
                    "columns": len(table.rows[0].cells) if table.rows else 0,
                    "cells": []
                }
                
                for row_idx, row in enumerate(table.rows):
                    row_data = []
                    for col_idx, cell in enumerate(row.cells):
                        cell_xml_id = id(cell._element)
                        cell_text = cell.text.strip()
                        cell_data = {
                            "row": row_idx,
                            "column": col_idx,
                            "cell_xml_id": cell_xml_id,
                            "id_key": f"tid:{table_xml_id}/cid:{cell_xml_id}",
                            "text": cell_text,
                            "style": cell.style.name if hasattr(cell, 'style') else None
                        }
                        row_data.append(cell_data)
                    table_info["cells"].append(row_data)
                
                structure["body"].append(table_info)
                para_id_counter += 1

        elif file_type == "xlsx":
            wb = load_workbook(user_file, read_only=True, data_only=True)

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                for row_idx, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                    for col_idx, cell in enumerate(row, start=1):
                        if cell is None or str(cell).strip() == "":
                            continue
                        col_letter = sheet.cell(row=row_idx, column=col_idx).column_letter
                        cell_ref = f"{col_letter}{row_idx}"
                        structure["body"].append({
                            "index": cell_ref,
                            "type": "cell",
                            "text": str(cell)
                        })
                        index_counter += 1

        elif file_type == "pptx":
            prs = Presentation(user_file)
            structure["slide_id_order"] = [int(s.slide_id) for s in prs.slides]
            for slide_idx, slide in enumerate(prs.slides):
                title_shape = slide.shapes.title if hasattr(slide.shapes, "title") else None
                title_text = title_shape.text.strip() if (title_shape and getattr(title_shape, "text", "").strip()) else ""

                slide_obj = {
                    "index": slide_idx,
                    "slide_id": int(slide.slide_id),
                    "id_key": f"sid:{int(slide.slide_id)}",
                    "title": title_text,
                    "shapes": []
                }
                

                for shape_idx, shape in enumerate(slide.shapes):
                    key = f"s{slide_idx}/sh{shape_idx}"
                    if hasattr(shape, "image"):
                        shape_id_val = getattr(shape, "shape_id", None) or shape._element.cNvPr.id  
                        slide_obj["shapes"].append({
                            "shape_idx": shape_idx,
                            "shape_id": shape_id_val,
                            "idx_key": key, 
                            "id_key": f"sid:{int(slide.slide_id)}/shid:{int(shape_id_val)}",
                            "kind": "image"
                        })
                        continue

                    if hasattr(shape, "text_frame") and shape.text_frame:
                        kind = "title" if (title_shape is not None and shape is title_shape) else "textbox"

                        paragraphs = []
                        for p in shape.text_frame.paragraphs:
                            text = "".join(run.text for run in p.runs) if p.runs else p.text
                            text = (text or "").strip()
                            if text != "":
                                paragraphs.append(text)

                        shape_id_val = getattr(shape, "shape_id", None) or shape._element.cNvPr.id 
                        slide_obj["shapes"].append({
                            "shape_idx": shape_idx,
                            "shape_id": shape_id_val,
                            "idx_key": key, 
                            "id_key": f"sid:{int(slide.slide_id)}/shid:{int(shape_id_val)}",
                            "kind": kind,
                            "paragraphs": paragraphs
                        })
                        continue


                structure["body"].append(slide_obj)

        else:
            return json.dumps({
                "error": {"message": f"Unsupported file type: {file_type}. Only docx, xlsx, and pptx are supported."}
            }, indent=4, ensure_ascii=False)

        return json.dumps(structure, indent=4, ensure_ascii=False)

    except Exception as e:
        return json.dumps({"error": {"message": str(e)}}, indent=4, ensure_ascii=False)

def add_auto_sized_review_comment(cell, text, author="AI Reviewer"):
    """
    Adds a note to an Excel cell, adjusting the width and height
    so that all the text is visible.
    """
    if not text:
        return

    avg_char_width = 7
    px_per_line = 15
    base_width = 200
    max_width = 500
    min_height = 40

    width = min(max_width, base_width + len(text) * 2)
    chars_per_line = max(1, width // avg_char_width)
    lines = 0
    for paragraph in text.split('\n'):
        lines += math.ceil(len(paragraph) / chars_per_line)
    height = max(min_height, lines * px_per_line)

    comment = Comment(text, author)
    comment.width = width
    comment.height = height
    cell.comment = comment


def _snapshot_runs(p):
    """Return a list of {'text': str, 'font': {...}} for each run in a paragraph."""
    runs = []
    for r in p.runs:
        f = r.font
        font_spec = {
            "name": f.name,
            "size": f.size,
            "bold": f.bold,
            "italic": f.italic,
            "underline": f.underline,
            "color_rgb": getattr(getattr(f.color, "rgb", None), "rgb", None) or getattr(f.color, "rgb", None)
        }
        runs.append({"text": r.text or "", "font": font_spec})
    return runs

def _apply_font(run, font_spec):
    """Apply font specifications to a run."""
    if not font_spec:
        return
    f = run.font
    try:
        if font_spec.get("name") is not None:
            f.name = font_spec["name"]
        if font_spec.get("size") is not None:
            f.size = font_spec["size"]
        if font_spec.get("bold") is not None:
            f.bold = font_spec["bold"]
        if font_spec.get("italic") is not None:
            f.italic = font_spec["italic"]
        if font_spec.get("underline") is not None:
            f.underline = font_spec["underline"]
        rgb = font_spec.get("color_rgb")
        if rgb is not None:
            try:
                f.color.rgb = rgb
            except Exception:
                pass
    except Exception:
        pass

def _set_text_with_runs(shape, new_content):
    """
    Set the text of a shape while preserving the original run-level formatting.
    """

    if not (hasattr(shape, "text_frame") and shape.text_frame):
        return
    tf = shape.text_frame

    if isinstance(new_content, list):
        lines = [str(item) for item in new_content]
    else:
        lines = [str(new_content or "")]

    original_para_styles = []
    original_para_runs = []     

    for p in tf.paragraphs:
        level = int(getattr(p, "level", 0) or 0)
        alignment = getattr(p, "alignment", None)
        original_para_styles.append({"level": level, "alignment": alignment})
        original_para_runs.append(_snapshot_runs(p))

    tf.clear()

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if (i == 0 and tf.paragraphs) else tf.add_paragraph()

        if original_para_styles:
            style = original_para_styles[i] if i < len(original_para_styles) else original_para_styles[-1]
            p.level = style.get("level", 0)
            if style.get("alignment") is not None:
                p.alignment = style["alignment"]

        runs_spec = (
            original_para_runs[i] if i < len(original_para_runs)
            else (original_para_runs[-1] if original_para_runs else [])
        )

        if not runs_spec:
            r = p.add_run()
            r.text = line
            continue

        n = len(runs_spec)
        total = len(line)

        if total == 0:
            for spec in runs_spec:
                r = p.add_run()
                r.text = ""
                _apply_font(r, spec["font"])
        else:
            base, rem = divmod(total, n)
            sizes = [base + (1 if k < rem else 0) for k in range(n)]
            pos = 0
            for k, spec in enumerate(runs_spec):
                seg = line[pos:pos + sizes[k]]
                pos += sizes[k]
                r = p.add_run()
                r.text = seg
                _apply_font(r, spec["font"])

def shape_by_id(slide, shape_id):
    sid = int(shape_id)
    for sh in slide.shapes:
        val = getattr(sh, "shape_id", None) or getattr(getattr(sh, "_element", None), "cNvPr", None)
        val = int(getattr(val, "id", val)) if val is not None else None
        if val == sid:
            return sh
    return None


def ensure_slot_textbox(slide, slot):
    slot = (slot or "").lower()

    def _get(ph_name):
        return getattr(PP_PLACEHOLDER, ph_name, None)

    TITLE = _get("TITLE")
    CENTER_TITLE = _get("CENTER_TITLE")
    SUBTITLE = _get("SUBTITLE")
    BODY = _get("BODY")
    CONTENT = _get("CONTENT")
    OBJECT = _get("OBJECT")

    title_types = {t for t in (TITLE, CENTER_TITLE, SUBTITLE) if t is not None}
    body_types  = {t for t in (BODY, CONTENT, OBJECT) if t is not None}

    def find_placeholder(accepted_types):
        for sh in slide.shapes:
            if not getattr(sh, "is_placeholder", False):
                continue
            pf = getattr(sh, "placeholder_format", None)
            if not pf:
                continue
            try:
                if pf.type in accepted_types:
                    return sh
            except Exception:
                pass
        return None

    if slot == "title":
        ph = find_placeholder(title_types)
        if ph:
            return ph

    if slot == "body":
        ph = find_placeholder(body_types)
        if ph:
            return ph

    from pptx.util import Inches
    if slot == "title":
        return slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    if slot == "body":
        return slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    return slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.5))

def _layout_has(layout, want_title=False, want_body=False):  # ADD
    has_title = has_body = False
    for ph in getattr(layout, "placeholders", []):
        pf = getattr(ph, "placeholder_format", None)
        t = getattr(pf, "type", None) if pf else None
        if t in (getattr(PP_PLACEHOLDER, "TITLE", None),
                 getattr(PP_PLACEHOLDER, "CENTER_TITLE", None),
                 getattr(PP_PLACEHOLDER, "SUBTITLE", None)):
            has_title = True
        if t in (getattr(PP_PLACEHOLDER, "BODY", None),
                 getattr(PP_PLACEHOLDER, "CONTENT", None),
                 getattr(PP_PLACEHOLDER, "OBJECT", None)):
            has_body = True
    return (not want_title or has_title) and (not want_body or has_body)

def _pick_layout_for_slots(prs, anchor_slide, needs_title, needs_body):  # ADD
    if anchor_slide and _layout_has(anchor_slide.slide_layout, needs_title, needs_body):
        return anchor_slide.slide_layout
    for layout in prs.slide_layouts:
        if _layout_has(layout, needs_title, needs_body):
            return layout
    return anchor_slide.slide_layout if anchor_slide else prs.slide_layouts[-1]

def _collect_needs(edit_items):  # ADD
    needs = {}
    for tgt, _ in edit_items:
        if not isinstance(tgt, str):
            continue
        m = re.match(r"^(n\d+):slot:(title|body)$", tgt.strip(), flags=re.I)
        if m:
            ref, slot = m.group(1), m.group(2).lower()
            needs.setdefault(ref, {"title": False, "body": False})
            needs[ref][slot] = True
    return needs

@mcp.tool()
async def edit_document(
    file_id: str,
    file_name: str,
    edits: dict,
    ctx: Context[ServerSession, None]
) -> dict:
    """
    Edits a document (docx, xlsx, pptx) using structured operations.

    Args:
        file_id: Unique identifier for the document.
        file_name: Name of the document file.
        edits: Dictionary with:
            - "ops": List of structural changes.
            - "content_edits": List of content updates.

    ## Supported Formats

    ### PPTX (PowerPoint)
    - ops: 
        - ["insert_after", slide_id, "nK"]
        - ["insert_before", slide_id, "nK"]
        - ["delete_slide", slide_id]
    - content_edits:
        - ["sid:<slide_id>/shid:<shape_id>", text_or_list]
        - ["nK:slot:title", text_or_list]
        - ["nK:slot:body", text_or_list]

    ### DOCX (Word)
    - ops:
        - ["insert_after", para_xml_id, "nK"]
        - ["insert_before", para_xml_id, "nK"]
        - ["delete_paragraph", para_xml_id]
    - content_edits:
        - ["pid:<para_xml_id>", text_or_list]
        - ["tid:<table_xml_id>/cid:<cell_xml_id>", text]
        - ["nK", text_or_list]

    ### XLSX (Excel)
    - ops:
        - ["insert_row", "sheet_name", row_idx]
        - ["delete_row", "sheet_name", row_idx]
        - ["insert_column", "sheet_name", col_idx]
        - ["delete_column", "sheet_name", col_idx]
    - content_edits:
        - ["<ref>", value]

    ## Notes
    - Always call `full_context_document()` first to get IDs.
    - Use cell refs like "A1", "B5".
    - Formatting is preserved.
    - Returns a download link to the edited file.
    """
    temp_folder = f"/app/temp/{uuid.uuid4()}"
    os.makedirs(temp_folder, exist_ok=True)
    try:
        bearer_token = ctx.request_context.request.headers.get("authorization")
        logging.info(f"Recieved authorization header!")
        user_token=bearer_token
    except:
        logging.error(f"Error retrieving authorization header")
        user_token=TOKEN
    try:
        user_file = download_file(file_id, token=user_token)
        if isinstance(user_file, dict) and "error" in user_file:
            return json.dumps(user_file, indent=4, ensure_ascii=False)

        file_extension = os.path.splitext(file_name)[1].lower()
        file_type = file_extension.lstrip('.')

        edited_path = None
        response = None

        if file_type == "docx":
            try:
                doc = Document(user_file)
                
                para_by_xml_id = {}
                table_by_xml_id = {}
                cell_by_xml_id = {}
                
                para_id_counter = 1
                for para in doc.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    para_by_xml_id[para_id_counter] = para
                    para_id_counter += 1
          
                for table in doc.tables:
                    table_xml_id = id(table._element)
                    table_by_xml_id[table_xml_id] = table
                    for row in table.rows:
                        for cell in row.cells:
                            cell_xml_id = id(cell._element)
                            cell_by_xml_id[cell_xml_id] = cell
          
                if isinstance(edits, dict):
                    ops = edits.get("ops", []) or []
                    edit_items = edits.get("content_edits", []) or []
                else:
                    ops = []
                    edit_items = edits
          
                new_refs = {}
                
                for op in ops:
                    if not isinstance(op, (list, tuple)) or not op:
                        continue
                    kind = op[0]
                    
                    if kind == "insert_after" and len(op) >= 3:
                        anchor_xml_id = int(op[1])
                        new_ref = op[2]
                        
                        anchor_para = para_by_xml_id.get(anchor_xml_id)
                        if anchor_para:
                            para_index = doc.paragraphs.index(anchor_para)
          	       	      
                            new_para = doc.add_paragraph()
          	       	      
                            anchor_element = anchor_para._element
                            parent = anchor_element.getparent()
                            parent.insert(parent.index(anchor_element) + 1, new_para._element)
          	       	      
                            new_para.style = anchor_para.style
      	      				
                            new_xml_id = id(new_para._element)
                            new_refs[new_ref] = new_xml_id
                            para_by_xml_id[new_xml_id] = new_para
      	      		
                    elif kind == "insert_before" and len(op) >= 3:
                        anchor_xml_id = int(op[1])
                        new_ref = op[2]
                        
                        anchor_para = para_by_xml_id.get(anchor_xml_id)
                        if anchor_para:
                            new_para = doc.add_paragraph()
          	       	      
                            anchor_element = anchor_para._element
                            parent = anchor_element.getparent()
                            parent.insert(parent.index(anchor_element), new_para._element)
          	       	      
                            new_para.style = anchor_para.style
      	      				
                            new_xml_id = id(new_para._element)
                            new_refs[new_ref] = new_xml_id
                            para_by_xml_id[new_xml_id] = new_para
      	      		
                    elif kind == "delete_paragraph" and len(op) >= 2:
                        para_xml_id = int(op[1])
                        para = para_by_xml_id.get(para_xml_id)
                        if para:
                            parent = para._element.getparent()
                            parent.remove(para._element)
                            para_by_xml_id.pop(para_xml_id, None)
                
                for target, new_text in edit_items:
                    if not isinstance(target, str):
                        continue
      	      		
                    t = target.strip()
                    
                    m = re.match(r"^pid:(\d+)$", t, flags=re.I)
                    if m:
                        para_xml_id = int(m.group(1))
                        para = para_by_xml_id.get(para_xml_id)
                        if para:
                            _apply_text_to_paragraph(para, new_text)
                        continue
      	      		
                    m = re.match(r"^tid:(\d+)/cid:(\d+)$", t, flags=re.I)
                    if m:
                        table_xml_id = int(m.group(1))
                        cell_xml_id = int(m.group(2))
                        cell = cell_by_xml_id.get(cell_xml_id)
                        if cell:
                            for para in cell.paragraphs:
                                for _ in range(len(para.runs)):
                                    para._element.remove(para.runs[0]._element)
          	       	      
                            if cell.paragraphs:
                                first_para = cell.paragraphs[0]
                                first_para.add_run(str(new_text))
                        continue
      	      		
                    m = re.match(r"^n(\d+)$", t, flags=re.I)
                    if m:
                        new_ref = t
                        para_xml_id = new_refs.get(new_ref)
                        if para_xml_id:
                            para = para_by_xml_id.get(para_xml_id)
                            if para:
                                _apply_text_to_paragraph(para, new_text)
                        continue
          
                edited_path = os.path.join(
                    temp_folder, f"{os.path.splitext(file_name)[0]}_edited.docx"
                )
                doc.save(edited_path)
                response = upload_file(
                    file_path=edited_path,
                    filename=f"{os.path.splitext(file_name)[0]}_edited",
                    file_type="docx", 
                    token=user_token
                )
            except Exception as e:
                raise Exception(f"Error during DOCX editing: {e}")

        elif file_type == "xlsx":
            try:
                wb = load_workbook(user_file)
                ws = wb.active

                edit_items = edits.get("content_edits", []) if isinstance(edits, dict) and "content_edits" in edits else edits
          
                for index, new_text in edit_items:
                    try:
                        if isinstance(index, str) and re.match(r"^[A-Z]+[0-9]+$", index.strip().upper()):
                            cell_ref = index.strip().upper()
                        elif isinstance(index, int):
                            cell_ref = f"A{index+1}"
                        else:
                            cell_ref = "A1"

                        cell = ws[cell_ref]
                        cell.value = new_text
                    except Exception:
                        fallback_cell = ws["A1"]
                        fallback_cell.value = new_text

                edited_path = os.path.join(
                    temp_folder, f"{os.path.splitext(file_name)[0]}_edited.xlsx"
                )
                wb.save(edited_path)
                response = upload_file(
                    file_path=edited_path,
                    filename=f"{os.path.splitext(file_name)[0]}_edited",
                    file_type="xlsx", 
                    token=user_token
                )
            except Exception as e:
                raise Exception(f"Error during XLSX editing: {e}")

        elif file_type == "pptx":
            try:
                prs = Presentation(user_file)
                if isinstance(edits, dict):
                    ops = edits.get("ops", []) or []
                    edit_items = edits.get("content_edits", []) or []
                else:
                    ops = []
                    edit_items = edits
                new_ref_needs = _collect_needs(edit_items)
                order = [int(s.slide_id) for s in prs.slides]
                slides_by_id = {int(s.slide_id): s for s in prs.slides}
                new_refs = {}
                
                for op in ops:
                    if not isinstance(op, (list, tuple)) or not op:
                        continue
                    kind = op[0]

                    if kind == "insert_after" and len(op) >= 3:
                        anchor_id = int(op[1])
                        new_ref = op[2]
                        if anchor_id in order:
                            style_donor = slides_by_id.get(order[-1])
                            needs = new_ref_needs.get(new_ref, {"title": False, "body": False})
                            layout = _pick_layout_for_slots(prs, style_donor, needs["title"], needs["body"])
                            new_slide = prs.slides.add_slide(layout)
                            new_sid = int(new_slide.slide_id)
                            sldIdLst = prs.slides._sldIdLst
                            new_sldId = sldIdLst[-1]
                            try:
                                anchor_pos = order.index(anchor_id)
                                sldIdLst.remove(new_sldId)
                                sldIdLst.insert(anchor_pos + 1, new_sldId)
                            except Exception:
                                pass 
                            order.insert(order.index(anchor_id) + 1, new_sid)
                            slides_by_id[new_sid] = new_slide
                            new_refs[new_ref] = new_sid

                    elif kind == "insert_before" and len(op) >= 3:
                        anchor_id = int(op[1])
                        new_ref = op[2]
                        if anchor_id in order:
                            style_donor = slides_by_id.get(order[-1])
                            needs = new_ref_needs.get(new_ref, {"title": False, "body": False})
                            layout = _pick_layout_for_slots(prs, style_donor, needs["title"], needs["body"])
                            new_slide = prs.slides.add_slide(layout)
                            new_sid = int(new_slide.slide_id)

                            sldIdLst = prs.slides._sldIdLst
                            new_sldId = sldIdLst[-1]
                            try:
                                anchor_pos = order.index(anchor_id)
                                sldIdLst.remove(new_sldId)
                                sldIdLst.insert(anchor_pos, new_sldId)
                            except Exception:
                                pass

                            order.insert(order.index(anchor_id), new_sid)
                            slides_by_id[new_sid] = new_slide
                            new_refs[new_ref] = new_sid

                    elif kind == "delete_slide" and len(op) >= 2:
                        sid = int(op[1])
                        if sid in order:
                            i = order.index(sid)
                            sldIdLst = prs.slides._sldIdLst
                            rId = sldIdLst[i].rId
                            prs.part.drop_rel(rId)
                            del sldIdLst[i]
                            order.pop(i)
                            slides_by_id.pop(sid, None)   		 
      		
                for target, new_text in edit_items:
                    if not isinstance(target, str):
                        continue
                    t = target.strip()
                    m = re.match(r"^sid:(\d+)/shid:(\d+)$", t, flags=re.I)
                    if m:
                        slide_id = int(m.group(1))
                        shape_id = int(m.group(2))
                        slide = slides_by_id.get(slide_id)
                        if not slide:
                            continue
                        shape = shape_by_id(slide, shape_id)
                        if not shape:
                            continue
                        _set_text_with_runs(shape, new_text)
                        continue
                    m = re.match(r"^(n\d+):slot:(title|body)$", t, flags=re.I)
                    if m:
                        ref = m.group(1)
                        slot = m.group(2).lower()
                        sid = new_refs.get(ref)
                        if not sid:
                            continue
                        slide = slides_by_id.get(sid)
                        if not slide:
                            continue
                        shape = ensure_slot_textbox(slide, slot)
                        tf = getattr(shape, "text_frame", None)
                        if tf is None:
                            continue
                        if isinstance(new_text, list):
                            tf.clear()
                            tf.text = str(new_text[0]) if new_text else ""
                            for line in new_text[1:]:
                                p = tf.add_paragraph()
                                p.text = str(line)
                                try:
                                    p.level = getattr(tf.paragraphs[0], "level", 0)
                                except Exception:
                                    pass
                        else:
                            tf.text = str(new_text)
                        continue
 

                edited_path = os.path.join(
                    temp_folder, f"{os.path.splitext(file_name)[0]}_edited.pptx"
                )
                prs.save(edited_path)
                response = upload_file(
                    file_path=edited_path,
                    filename=f"{os.path.splitext(file_name)[0]}_edited",
                    file_type="pptx", 
                    token=user_token
                )
            except Exception as e:
                raise Exception(f"Error during PPTX editing: {e}")

        else:
            raise Exception(f"File type not supported: {file_type}")

        shutil.rmtree(temp_folder, ignore_errors=True)

        return response

    except Exception as e:
        shutil.rmtree(temp_folder, ignore_errors=True)
        return json.dumps(
            {"error": {"message": str(e)}},
            indent=4,
            ensure_ascii=False
        )

def _get_pptx_namespaces():
    """Returns XML namespaces for PowerPoint"""
    return {
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p15': 'http://schemas.microsoft.com/office/powerpoint/2012/main',
        'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main'
    }

def _add_native_pptx_comment_zip(pptx_path, slide_num, comment_text, author_id, x=100, y=100):
    """
    Add a native PowerPoint comment by directly manipulating the ZIP file.
        Args:
        pptx_path: Path to the PPTX file
        slide_num: Slide number (1-based)
        comment_text: Comment text
        author_id: Author ID
        x: X position in EMU (not pixels!)
        y: Y position in EMU (not pixels!)
    """
    namespaces = _get_pptx_namespaces()
    
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            zf.extractall(temp_path)
        
        authors_file = temp_path / 'ppt' / 'commentAuthors.xml'
        if authors_file.exists():
            root = etree.parse(str(authors_file)).getroot()
            found = False
            for author in root.findall('.//p:cmAuthor', namespaces):
                if author.get('name') == 'AI Reviewer':
                    author_id = int(author.get('id'))
                    found = True
                    break
            
            if not found:
                existing_ids = [int(a.get('id')) for a in root.findall('.//p:cmAuthor', namespaces)]
                author_id = max(existing_ids) + 1 if existing_ids else 0
                author = etree.SubElement(root, f'{{{namespaces["p"]}}}cmAuthor')
                author.set('id', str(author_id))
                author.set('name', 'AI Reviewer')
                author.set('initials', 'AI')
                author.set('lastIdx', '1')
                author.set('clrIdx', str(author_id % 8))
        else:
            authors_file.parent.mkdir(parents=True, exist_ok=True)
            root = etree.Element(
                f'{{{namespaces["p"]}}}cmAuthorLst',
                nsmap={k: v for k, v in namespaces.items() if k in ['p']}
            )
            author = etree.SubElement(root, f'{{{namespaces["p"]}}}cmAuthor')
            author.set('id', str(author_id))
            author.set('name', 'AI Reviewer')
            author.set('initials', 'AI')
            author.set('lastIdx', '1')
            author.set('clrIdx', '0')
            
            rels_file = temp_path / 'ppt' / '_rels' / 'presentation.xml.rels'
            if rels_file.exists():
                rels_root = etree.parse(str(rels_file)).getroot()
                existing_ids = [int(rel.get('Id')[3:]) for rel in rels_root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship')]
                next_rid = max(existing_ids) + 1 if existing_ids else 1
                
                rel = etree.SubElement(rels_root, '{http://schemas.openxmlformats.org/package/2006/relationships}Relationship')
                rel.set('Id', f'rId{next_rid}')
                rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/commentAuthors')
                rel.set('Target', 'commentAuthors.xml')
                
                with open(rels_file, 'wb') as f:
                    f.write(etree.tostring(rels_root, xml_declaration=True, encoding='UTF-8', pretty_print=True))
        
        with open(authors_file, 'wb') as f:
            f.write(etree.tostring(root, xml_declaration=True, encoding='UTF-8', pretty_print=True))
        
        comments_dir = temp_path / 'ppt' / 'comments'
        comments_dir.mkdir(parents=True, exist_ok=True)
        comment_file = comments_dir / f'comment{slide_num}.xml'
        
        if comment_file.exists():
            comments_root = etree.parse(str(comment_file)).getroot()
        else:
            comments_root = etree.Element(
                f'{{{namespaces["p"]}}}cmLst',
                nsmap={k: v for k, v in namespaces.items() if k in ['p']}
            )
            
            slide_rels_file = temp_path / 'ppt' / 'slides' / '_rels' / f'slide{slide_num}.xml.rels'
            if slide_rels_file.exists():
                slide_rels_root = etree.parse(str(slide_rels_file)).getroot()
            else:
                slide_rels_file.parent.mkdir(parents=True, exist_ok=True)
                slide_rels_root = etree.Element(
                    '{http://schemas.openxmlformats.org/package/2006/relationships}Relationships'
                )
            
            existing_ids = [int(rel.get('Id')[3:]) for rel in slide_rels_root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship')]
            next_rid = max(existing_ids) + 1 if existing_ids else 1
            
            rel = etree.SubElement(slide_rels_root, '{http://schemas.openxmlformats.org/package/2006/relationships}Relationship')
            rel.set('Id', f'rId{next_rid}')
            rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments')
            rel.set('Target', f'../comments/comment{slide_num}.xml')
            
            with open(slide_rels_file, 'wb') as f:
                f.write(etree.tostring(slide_rels_root, xml_declaration=True, encoding='UTF-8', pretty_print=True))
        
        existing_ids = [int(c.get('idx')) for c in comments_root.findall('.//p:cm', namespaces)]
        next_id = max(existing_ids) + 1 if existing_ids else 1
        
        comment = etree.SubElement(comments_root, f'{{{namespaces["p"]}}}cm')
        comment.set('authorId', str(author_id))
        comment.set('dt', datetime.datetime.now().isoformat())
        comment.set('idx', str(next_id))
        
        pos = etree.SubElement(comment, f'{{{namespaces["p"]}}}pos')
        pos.set('x', str(int(x)))
        pos.set('y', str(int(y)))
        
        text_elem = etree.SubElement(comment, f'{{{namespaces["p"]}}}text')
        text_elem.text = comment_text
        
        with open(comment_file, 'wb') as f:
            f.write(etree.tostring(comments_root, xml_declaration=True, encoding='UTF-8', pretty_print=True))
        
        content_types_file = temp_path / '[Content_Types].xml'
        if content_types_file.exists():
            ct_root = etree.parse(str(content_types_file)).getroot()
            ns = {'ct': 'http://schemas.openxmlformats.org/package/2006/content-types'}
            
            has_authors = False
            has_comments = False
            
            for override in ct_root.findall('.//ct:Override', ns):
                if override.get('PartName') == '/ppt/commentAuthors.xml':
                    has_authors = True
                if override.get('PartName') == f'/ppt/comments/comment{slide_num}.xml':
                    has_comments = True
            
            if not has_authors:
                override = etree.SubElement(ct_root, '{http://schemas.openxmlformats.org/package/2006/content-types}Override')
                override.set('PartName', '/ppt/commentAuthors.xml')
                override.set('ContentType', 'application/vnd.openxmlformats-officedocument.presentationml.commentAuthors+xml')
            
            if not has_comments:
                override = etree.SubElement(ct_root, '{http://schemas.openxmlformats.org/package/2006/content-types}Override')
                override.set('PartName', f'/ppt/comments/comment{slide_num}.xml')
                override.set('ContentType', 'application/vnd.openxmlformats-officedocument.presentationml.comments+xml')
            
            with open(content_types_file, 'wb') as f:
                f.write(etree.tostring(ct_root, xml_declaration=True, encoding='UTF-8', pretty_print=True))
        
        with zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for file_path in temp_path.rglob('*'):
                if file_path.is_file():
                    arcname = str(file_path.relative_to(temp_path))
                    zf.write(file_path, arcname)
        
        log.debug(f"Native comment added to slide {slide_num} with idx={next_id}")

@mcp.tool(
    name="review_document",
    title="Review and comment on various document types",
    description="Review an existing document of various types (docx, xlsx, pptx), perform corrections and add comments."
)
async def review_document(
    file_id: str,
    file_name: str,
    review_comments: list[tuple[int | str, str]],
    ctx: Context[ServerSession, None]
) -> dict:
    """
    Generic document review function that works with different document types.
    File type is automatically detected from the file extension.
    Returns a markdown hyperlink for downloading the reviewed document.
    
    For Excel files (.xlsx):
    - The index must be a cell reference (e.g. "A1", "B3", "C10")
    - These correspond to the "index" key returned by the full_context_document() function
    - Never use integer values for Excel cells
    
    For Word files (.docx):
    - The index should be a paragraph ID in the format "pid:<para_xml_id>"
    - These correspond to the "id_key" field returned by the full_context_document() function
    
    For PowerPoint files (.pptx):
    - The index should be a slide ID in the format "sid:<slide_id>"
    - These correspond to the "id_key" field returned by the full_context_document() function
    """
    temp_folder = f"/app/temp/{uuid.uuid4()}"
    os.makedirs(temp_folder, exist_ok=True)
    try:
        bearer_token = ctx.request_context.request.headers.get("authorization")
        logging.info(f"Recieved authorization header!")
        user_token=bearer_token
    except:
        logging.error(f"Error retrieving authorization header")
        user_token=TOKEN    
    try:
        user_file = download_file(file_id, token=user_token)
        if isinstance(user_file, dict) and "error" in user_file:
            return json.dumps(user_file, indent=4, ensure_ascii=False)

        file_extension = os.path.splitext(file_name)[1].lower()
        file_type = file_extension.lstrip('.')

        reviewed_path = None
        response = None

        if file_type == "docx":
            try:
                doc = Document(user_file)
                paragraphs = list(doc.paragraphs)
                para_by_xml_id = {}
                para_id_counter = 1
                
                for para in doc.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    para_by_xml_id[para_id_counter] = para
                    para_id_counter += 1

                for index, comment_text in review_comments:
                    if isinstance(index, int) and 0 <= index < len(paragraphs):
                        para = paragraphs[index]
                        if para.runs:
                            try:
                                doc.add_comment(
                                    runs=[para.runs[0]],
                                    text=comment_text,
                                    author="AI Reviewer",
                                    initials="AI"
                                )
                            except Exception:
                                para.add_run(f"  [AI Comment: {comment_text}]")
                    elif isinstance(index, str) and index.startswith("pid:"):
                        try:
                            para_xml_id = int(index.split(":")[1])
                            para = para_by_xml_id.get(para_xml_id)
                            if para and para.runs:
                                try:
                                    doc.add_comment(
                                        runs=[para.runs[0]],
                                        text=comment_text,
                                        author="AI Reviewer",
                                        initials="AI"
                                    )
                                except Exception:
                                    para.add_run(f"  [AI Comment: {comment_text}]")
                        except Exception:
                            if isinstance(index, int) and 0 <= index < len(paragraphs):
                                para = paragraphs[index]
                                if para.runs:
                                    try:
                                        doc.add_comment(
                                            runs=[para.runs[0]],
                                            text=comment_text,
                                            author="AI Reviewer",
                                            initials="AI"
                                        )
                                    except Exception:
                                        para.add_run(f"  [AI Comment: {comment_text}]")
                reviewed_path = os.path.join(
                    temp_folder, f"{os.path.splitext(file_name)[0]}_reviewed.docx"
                )
                doc.save(reviewed_path)
                response = upload_file(
                    file_path=reviewed_path,
                    filename=f"{os.path.splitext(file_name)[0]}_reviewed",
                    file_type="docx", 
                    token=user_token
                )
            except Exception as e:
                raise Exception(f"Error during DOCX revision: {e}")

        elif file_type == "xlsx":
            try:
                wb = load_workbook(user_file)
                ws = wb.active

                for index, comment_text in review_comments:
                    try:
                        if isinstance(index, str) and re.match(r"^[A-Z]+[0-9]+$", index.strip().upper()):
                            cell_ref = index.strip().upper()
                        elif isinstance(index, int):
                            cell_ref = f"A{index+1}"
                        else:
                            cell_ref = "A1"

                        cell = ws[cell_ref]
                        add_auto_sized_review_comment(cell, comment_text, author="AI Reviewer")

                    except Exception:
                        fallback_cell = ws["A1"]
                        add_auto_sized_review_comment(fallback_cell, comment_text, author="AI Reviewer")

                reviewed_path = os.path.join(
                    temp_folder, f"{os.path.splitext(file_name)[0]}_reviewed.xlsx"
                )
                wb.save(reviewed_path)
                response = upload_file(
                    file_path=reviewed_path,
                    filename=f"{os.path.splitext(file_name)[0]}_reviewed",
                    file_type="xlsx", 
                    token=user_token
                )
            except Exception as e:
                raise Exception(f"Error: {e}")

        elif file_type == "pptx":
            try:
                temp_pptx = os.path.join(temp_folder, "temp_input.pptx")
                with open(temp_pptx, 'wb') as f:
                    f.write(user_file.read())
                
                prs = Presentation(temp_pptx)
                slides_by_id = {int(s.slide_id): s for s in prs.slides}
                
                comments_by_slide = {}
                
                for index, comment_text in review_comments:
                    slide_num = None
                    slide_id = None
                    
                    if isinstance(index, int) and 0 <= index < len(prs.slides):
                        slide_num = index + 1
                        slide_id = list(slides_by_id.keys())[index]
                    elif isinstance(index, str):
                        if index.startswith("sid:") and "/shid:" in index:
                            try:
                                slide_id = int(index.split("/")[0].replace("sid:", ""))
                                if slide_id in slides_by_id:
                                    slide_num = list(slides_by_id.keys()).index(slide_id) + 1
                            except Exception as e:
                                log.warning(f"Failed to parse shape ID: {e}")
                        elif index.startswith("sid:"):
                            try:
                                slide_id = int(index.split(":")[1])
                                if slide_id in slides_by_id:
                                    slide_num = list(slides_by_id.keys()).index(slide_id) + 1
                            except Exception as e:
                                log.warning(f"Failed to parse slide ID: {e}")
                    
                    if slide_num and slide_id:
                        if slide_num not in comments_by_slide:
                            comments_by_slide[slide_num] = []
                        
                        shape_info = ""
                        if "/shid:" in str(index):
                            try:
                                shape_id = int(str(index).split("/shid:")[1])
                                shape_info = f"[Shape {shape_id}] "
                            except:
                                pass
                        
                        comments_by_slide[slide_num].append(f"{shape_info}{comment_text}")
                comment_offset = 0              
                for slide_num, comments in comments_by_slide.items():
                    comment_start_x = 5000
                    comment_start_y = 1000
                    comment_spacing_y = 1500
                    
                    for idx, comment in enumerate(comments):
                        try:
                            y_position = comment_start_y + (idx * comment_spacing_y)
                            
                            _add_native_pptx_comment_zip(
                                pptx_path=temp_pptx,
                                slide_num=slide_num,
                                comment_text=f"• {comment}",
                                author_id=0,
                                x=comment_start_x,
                                y=y_position
                            )
                            log.debug(f"Native PowerPoint comment added to slide {slide_num} at position x={comment_start_x}, y={y_position}")
                        except Exception as e:
                            log.warning(f"Failed to add native comment to slide {slide_num}: {e}", exc_info=True)
                            prs_fallback = Presentation(temp_pptx)
                            slide = prs_fallback.slides[slide_num - 1]
                            left = top = Inches(0.2)
                            width = Inches(4)
                            height = Inches(1)
                            textbox = slide.shapes.add_textbox(left, top, width, height)
                            text_frame = textbox.text_frame
                            p = text_frame.add_paragraph()
                            p.text = f"AI Reviewer: {comment}"
                            p.font.size = PptPt(10)
                            prs_fallback.save(temp_pptx)

                reviewed_path = os.path.join(
                    temp_folder, f"{os.path.splitext(file_name)[0]}_reviewed.pptx"
                )
                shutil.copy(temp_pptx, reviewed_path)
                
                response = upload_file(
                    file_path=reviewed_path,
                    filename=f"{os.path.splitext(file_name)[0]}_reviewed",
                    file_type="pptx", 
                    token=user_token
                )
            except Exception as e:
                raise Exception(f"Error when revising PPTX: {e}")

        else:
            raise Exception(f"File type not supported : {file_type}")

        shutil.rmtree(temp_folder, ignore_errors=True)

        return response

    except Exception as e:
        shutil.rmtree(temp_folder, ignore_errors=True)
        return json.dumps(
            {"error": {"message": str(e)}},
            indent=4,
            ensure_ascii=False
        )

@mcp.tool()
async def create_file(data: dict, persistent: bool = PERSISTENT_FILES) -> dict:
    """ "{"data": {"format":"pdf","filename":"report.pdf","content":[{"type":"title","text":"..."},{"type":"paragraph","text":"..."}],"title":"..."}}
"{"data": {"format":"docx","filename":"doc.docx","content":[{"type":"title","text":"..."},{"type":"list","items":[...]}],"title":"..."}}"
"{"data": {"format":"pptx","filename":"slides.pptx","slides_data":[{"title":"...","content":[...],"image_query":"...","image_position":"left|right|top|bottom","image_size":"small|medium|large"}],"title":"..."}}"
"{"data": {"format":"xlsx","filename":"data.xlsx","content":[["Header1","Header2"],["Val1","Val2"]],"title":"..."}}"
"{"data": {"format":"csv","filename":"data.csv","content":[[...]]}}"
"{"data": {"format":"txt|xml|py|etc","filename":"file.ext","content":"string"}}" """
    log.debug("Creating file via tool")
    folder_path = _generate_unique_folder()
    format_type = (data.get("format") or "").lower()
    filename = data.get("filename")
    content = data.get("content")
    title = data.get("title")

    if format_type == "pdf":
        result = _create_pdf(content if content is not None else "", filename, folder_path=folder_path)
    elif format_type == "pptx":
        result = _create_presentation(data.get("slides_data", []), filename, folder_path=folder_path, title=title)
    elif format_type == "docx":
        result = _create_word(content if content is not None else [], filename, folder_path=folder_path, title=title)
    elif format_type == "xlsx":
        result = _create_excel(content if content is not None else [], filename, folder_path=folder_path, title=title)
    elif format_type == "csv":
        result = _create_csv(content if content is not None else [], filename, folder_path=folder_path)
    else:
        use_filename = filename or f"export.{format_type or 'txt'}"
        result = _create_raw_file(content if content is not None else "", use_filename, folder_path=folder_path)

    if not persistent:
        _cleanup_files(folder_path, FILES_DELAY)

    return {"url": result["url"]}

@mcp.tool()
async def generate_and_archive(files_data: list[dict], archive_format: str = "zip", archive_name: str = None, persistent: bool = PERSISTENT_FILES) -> dict:
    """files_data=[{"format":"pdf","filename":"report.pdf","content":[{"type":"title","text":"..."},{"type":"paragraph","text":"..."}],"title":"..."},{"format":"docx","filename":"doc.docx","content":[{"type":"title","text":"..."},{"type":"list","items":[...]}],"title":"..."},{"format":"pptx","filename":"slides.pptx","slides_data":[{"title":"...","content":[...],"image_query":"...","image_position":"left|right|top|bottom","image_size":"small|medium|large"}],"title":"..."},{"format":"xlsx","filename":"data.xlsx","content":[["Header1","Header2"],["Val1","Val2"]],"title":"..."},{"format":"csv","filename":"data.csv","content":[[...]]},{"format":"txt|xml|py|etc","filename":"file.ext","content":"string"}]"""
    log.debug("Generating archive via tool")
    folder_path = _generate_unique_folder()
    generated_paths: list[str] = []

    for file_info in files_data or []:
        fmt = (file_info.get("format") or "").lower()
        fname = file_info.get("filename")
        content = file_info.get("content")
        title = file_info.get("title")

        try:
            if fmt == "pdf":
                res = _create_pdf(content if content is not None else "", fname, folder_path=folder_path)
            elif fmt == "pptx":
                res = _create_presentation(file_info.get("slides_data", []), fname, folder_path=folder_path, title=title)
            elif fmt == "docx":
                res = _create_word(content if content is not None else [], fname, folder_path=folder_path, title=title)
            elif fmt == "xlsx":
                res = _create_excel(content if content is not None else [], fname, folder_path=folder_path, title=title)
            elif fmt == "csv":
                res = _create_csv(content if content is not None else [], fname, folder_path=folder_path)
            else:
                use_fname = fname or f"export.{fmt or 'txt'}"
                res = _create_raw_file(content if content is not None else "", use_fname, folder_path=folder_path)
        except Exception as e:
            log.error(f"Error generating file {fname or '<no name>'}: {e}", exc_info=True)
            raise

        generated_paths.append(res["path"])

    timestamp = datetime.datetime.now().strftime('%Y%m%d_%H%M%S')
    archive_basename = f"{archive_name or 'archive'}_{timestamp}"
    archive_filename = f"{archive_basename}.zip" if archive_format.lower() not in ("7z", "tar.gz") else f"{archive_basename}.{archive_format}"
    archive_path = os.path.join(folder_path, archive_filename)

    if archive_format.lower() == "7z":
        with py7zr.SevenZipFile(archive_path, mode='w') as archive:
            for p in generated_paths:
                archive.write(p, os.path.relpath(p, folder_path))
    elif archive_format.lower() == "tar.gz":
        with tarfile.open(archive_path, "w:gz") as tar:
            for p in generated_paths:
                tar.add(p, arcname=os.path.relpath(p, folder_path))
    else:
        with zipfile.ZipFile(archive_path, 'w') as zipf:
            for p in generated_paths:
                zipf.write(p, os.path.relpath(p, folder_path))

    if not persistent:
        _cleanup_files(folder_path, FILES_DELAY)

    return {"url": _public_url(folder_path, archive_filename)}

from sse_starlette.sse import EventSourceResponse

class SimpleRequestContext:
    def __init__(self, request):
        self.request = request

class SimpleCtx:
    def __init__(self, request):
        self.request_context = SimpleRequestContext(request)

async def handle_sse(request: Request) -> Response:
    """Handle SSE transport for MCP - supports both GET and POST"""
    
    if request.method == "POST":
        try:
            message = await request.json()
            log.debug(f"Received POST message: {message}")
            
            response = {
                "jsonrpc": "2.0",
                "id": message.get("id"),
                "result": None
            }
            
            method = message.get("method")
            
            if method == "initialize":
                response["result"] = {
                    "protocolVersion": "2024-11-05",
                    "capabilities": {
                        "tools": {}
                    },
                    "serverInfo": {
                        "name": "file_export_mcp",
                        "version": SCRIPT_VERSION
                    }
                }
            elif method == "tools/list":
                response["result"] = {
                    "tools": [
                        {
                            "name": "create_file",
                            "description": "Create files in various formats (pdf, docx, pptx, xlsx, csv, txt, xml, py, etc.). Supports rich content including titles, paragraphs, lists, tables, images via queries, and more.",
                            "inputSchema": {
                                "type": "object",
                                "properties": {
                                    "data": {
                                        "type": "object",
                                        "description": "File data configuration",
                                        "properties": {
                                            "format": {
                                                "type": "string",
                                                "enum": ["pdf", "docx", "pptx", "xlsx", "csv", "txt", "xml", "py", "json", "md"],
                                                "description": "Output file format"
                                            },
                                            "filename": {
                                                "type": "string",
                                                "description": "Name of the file to create (optional, will be auto-generated if not provided)"
                                            },
                                            "title": {
                                                "type": "string",
                                                "description": "Document title (for docx, pptx, xlsx, pdf)"
                                            },
                                            "content": {
                                                "description": "Content varies by format. For pdf/docx: array of objects with type/text. For xlsx/csv: 2D array. For pptx: use slides_data instead. For txt/xml/py: string",
                                                "oneOf": [
                                                    {"type": "array"},
                                                    {"type": "string"}
                                                ]
                                            },
                                            "slides_data": {
                                                "type": "array",
                                                "description": "For pptx format only: array of slide objects",
                                                "items": {
                                                    "type": "object",
                                                    "properties": {
                                                        "title": {"type": "string"},
                                                        "content": {
                                                            "type": "array",
                                                            "items": {"type": "string"}
                                                        },
                                                        "image_query": {
                                                            "type": "string",
                                                            "description": "Search query for image (Unsplash, Pexels, or local SD)"
                                                        },
                                                        "image_position": {
                                                            "type": "string",
                                                            "enum": ["left", "right", "top", "bottom"],
                                                            "description": "Position of the image on the slide"
                                                        },
                                                        "image_size": {
                                                            "type": "string",
                                                            "enum": ["small", "medium", "large"],
                                                            "description": "Size of the image"
                                                        }
                                                    }
                                                }
                                            }
                                        },
                                        "required": ["format"]
                                    },
                                    "persistent": {
                                        "type": "boolean",
                                        "description": "Whether to keep files permanently (default: false, files deleted after delay)"
                                    }
                                },
                                "required": ["data"]
                            }
                        },
                        {
                            "name": "generate_and_archive",
                            "description": "Generate multiple files at once and create an archive (zip, 7z, or tar.gz). Perfect for creating document packages with multiple formats.",
                            "inputSchema": {
                                "type": "object",
                                "properties": {
                                    "files_data": {"type": "array", "description": "Array of file data objects"},
                                    "archive_format": {"type": "string", "enum": ["zip", "7z", "tar.gz"]},
                                    "archive_name": {"type": "string"},
                                    "persistent": {"type": "boolean"}
                                },
                                "required": ["files_data"]
                            }
                        },
                        {
                            "name": "full_context_document",
                            "description": "Return the structure, content, and metadata of a document",
                            "inputSchema": {
                                "type": "object",
                                "properties": {
                                    "file_id": {"type": "string", "description": "The file ID from OpenWebUI"},
                                    "file_name": {"type": "string", "description": "The name of the file"}
                                },
                                "required": ["file_id", "file_name"]
                            }
                        },
                        {
                            "name": "review_document",
                            "description": "Review and comment on various document types (docx, xlsx, pptx)",
                            "inputSchema": {
                                "type": "object",
                                "properties": {
                                    "file_id": {"type": "string", "description": "The file ID from OpenWebUI"},
                                    "file_name": {"type": "string", "description": "The name of the file"},
                                    "review_comments": {
                                        "type": "array",
                                        "description": "Array of file configurations (same structure as create_file data)",
                                        "items": {
                                            "type": "object",
                                            "properties": {
                                                "format": {"type": "string"},
                                                "filename": {"type": "string"},
                                                "title": {"type": "string"},
                                                "content": {
                                                    "oneOf": [
                                                        {"type": "array"},
                                                        {"type": "string"}
                                                    ]
                                                },
                                                "slides_data": {"type": "array"}
                                            },
                                            "required": ["format"]
                                        }
                                    },
                                    "archive_format": {
                                        "type": "string",
                                        "enum": ["zip", "7z", "tar.gz"],
                                        "description": "Archive format (default: zip)"
                                    },
                                    "archive_name": {
                                        "type": "string",
                                        "description": "Name of the archive file (without extension, timestamp will be added)"
                                    },
                                    "persistent": {
                                        "type": "boolean",
                                        "description": "Whether to keep the archive permanently"
                                    }
                                },
                                "required": ["files_data"]
                            }
                        },
                        {
                            "name": "full_context_document",
                            "description": "Extract and return the complete structure, content, and metadata of a document (docx, xlsx, pptx). Returns a JSON structure with indexed elements (paragraphs, headings, tables, cells, slides, images) that can be referenced for editing or review.",
                            "inputSchema": {
                                "type": "object",
                                "properties": {
                                    "file_id": {
                                        "type": "string",
                                        "description": "The file ID from OpenWebUI file upload"
                                    },
                                    "file_name": {
                                        "type": "string",
                                        "description": "The name of the file with extension (e.g., 'report.docx', 'data.xlsx', 'presentation.pptx')"
                                    }
                                },
                                "required": ["file_id", "file_name"]
                            }
                        },
                        {
                            "name": "edit_document",
                            "description": "Edit an existing document (docx, xlsx, pptx) using structured operations. Supports inserting/deleting elements and updating content. ALWAYS call full_context_document() first to get proper IDs and references. Preserves formatting and returns a download link for the edited file.",
                            "inputSchema": {
                                "type": "object",
                                "properties": {
                                    "file_id": {
                                        "type": "string",
                                        "description": "The file ID from OpenWebUI"
                                    },
                                    "file_name": {
                                        "type": "string",
                                        "description": "The name of the file with extension"
                                    },
                                    "edits": {
                                        "type": "object",
                                        "description": "Edit operations and content changes",
                                        "properties": {
                                            "ops": {
                                                "type": "array",
                                                "description": "Structural operations (insert/delete). For PPTX: ['insert_after', slide_id, 'nK'], ['insert_before', slide_id, 'nK'], ['delete_slide', slide_id]. For DOCX: ['insert_after', para_xml_id, 'nK'], ['insert_before', para_xml_id, 'nK'], ['delete_paragraph', para_xml_id]. For XLSX: ['insert_row', 'sheet_name', row_idx], ['delete_row', 'sheet_name', row_idx], ['insert_column', 'sheet_name', col_idx], ['delete_column', 'sheet_name', col_idx]",
                                                "items": {
                                                    "type": "array",
                                                    "items": {
                                                        "oneOf": [
                                                            {"type": "string"},
                                                            {"type": "integer"}
                                                        ]
                                                    }
                                                }
                                            },
                                            "content_edits": {
                                                "type": "array",
                                                "description": "Content updates as [target, new_text] pairs. For PPTX: ['sid:<slide_id>/shid:<shape_id>', text], ['nK:slot:title', text], ['nK:slot:body', text]. For DOCX: ['pid:<para_xml_id>', text], ['tid:<table_xml_id>/cid:<cell_xml_id>', text], ['nK', text]. For XLSX: ['A1', value], ['B5', value]",
                                                "items": {
                                                    "type": "array",
                                                    "minItems": 2,
                                                    "maxItems": 2,
                                                    "items": [
                                                        {
                                                            "type": "string",
                                                            "description": "Target reference (element ID or cell ref)"
                                                        },
                                                        {
                                                            "description": "New content (text string or array of strings for lists)",
                                                            "oneOf": [
                                                                {"type": "string"},
                                                                {"type": "array", "items": {"type": "string"}},
                                                                {"type": "number"},
                                                                {"type": "boolean"}
                                                            ]
                                                        }
                                                    ]
                                                }
                                            }
                                        }
                                    }
                                },
                                "required": ["file_id", "file_name", "edits"]
                            }
                        },
                        {
                            "name": "review_document",
                            "description": "Review and add comments/corrections to an existing document (docx, xlsx, pptx). Returns a download link for the reviewed document with comments added. For Excel files, the index MUST be a cell reference (e.g., 'A1', 'B5', 'C10') as returned by full_context_document. For Word/PowerPoint, use integer indices.",
                            "inputSchema": {
                                "type": "object",
                                "properties": {
                                    "file_id": {
                                        "type": "string",
                                        "description": "The file ID from OpenWebUI"
                                    },
                                    "file_name": {
                                        "type": "string",
                                        "description": "The name of the file with extension"
                                    },
                                    "review_comments": {
                                        "type": "array",
                                        "description": "Array of [index, comment_text] tuples. For Excel: index must be a cell reference string like 'A1', 'B3'. For Word: integer paragraph index. For PowerPoint: integer slide index.",
                                        "items": {
                                            "type": "array",
                                            "minItems": 2,
                                            "maxItems": 2,
                                            "items": [
                                                {
                                                    "description": "Index/reference: For Excel use cell reference (e.g., 'A1'), for Word/PowerPoint use integer",
                                                    "oneOf": [
                                                        {"type": "string"},
                                                        {"type": "integer"}
                                                    ]
                                                },
                                                {
                                                    "type": "string",
                                                    "description": "Comment or correction text"
                                                }
                                            ]
                                        }
                                    }
                                },
                                "required": ["file_id", "file_name", "review_comments"]
                            }
                        }
                    ]
                }
            elif method == "tools/call":
                params = message.get("params", {})
                tool_name = params.get("name")
                arguments = params.get("arguments", {}) or {}
                ctx = SimpleCtx(request)

                try:
                    if tool_name == "create_file":
                        result = await create_file(**arguments)
                        response["result"] = {
                            "content": [
                                {"type": "text", "text": f"File created successfully: {result.get('url')}"}
                            ]
                        }

                    elif tool_name == "generate_and_archive":
                        result = await generate_and_archive(**arguments)
                        response["result"] = {
                            "content": [
                                {"type": "text", "text": f"Archive created successfully: {result.get('url')}"}
                            ]
                        }

                    elif tool_name == "full_context_document":
                        arguments.setdefault("ctx", ctx)
                        result = await full_context_document(**arguments)
                        response["result"] = {
                            "content": [
                                {"type": "text", "text": result}
                            ]
                        }

                    elif tool_name == "edit_document":
                        arguments.setdefault("ctx", ctx)
                        result = await edit_document(**arguments)
                        response["result"] = {
                            "content": [
                                {"type": "text", "text": json.dumps(result, indent=2, ensure_ascii=False)}
                            ]
                        }

                    elif tool_name == "review_document":
                        arguments.setdefault("ctx", ctx)
                        result = await review_document(**arguments)
                        response["result"] = {
                            "content": [
                                {"type": "text", "text": json.dumps(result, indent=2, ensure_ascii=False)}
                            ]
                        }

                    else:
                        response["error"] = {
                            "code": -32601,
                            "message": f"Tool not found: {tool_name}"
                        }
                except Exception as e:
                    log.error(f"Error executing tool {tool_name}: {e}", exc_info=True)
                    response["error"] = {
                        "code": -32000,
                        "message": str(e)
                    }
            else:
                response["error"] = {
                    "code": -32601,
                    "message": f"Method not found: {method}"
                }
            
            return JSONResponse(response)
            
        except Exception as e:
            log.error(f"Error handling POST request: {e}", exc_info=True)
            return JSONResponse(
                {"jsonrpc": "2.0", "error": {"code": -32000, "message": str(e)}},
                status_code=500
            )
    
    else:
        async def event_generator():
            """Generator for SSE events"""
            try:
                yield {
                    "event": "endpoint",
                    "data": json.dumps({
                        "endpoint": "/messages"
                    })
                }
                
                import asyncio
                while True:
                    await asyncio.sleep(15)
                    yield {
                        "event": "ping",
                        "data": ""
                    }
                    
            except asyncio.CancelledError:
                log.info("SSE connection closed by client")
                raise
            except Exception as e:
                log.error(f"SSE Error: {e}", exc_info=True)
                yield {
                    "event": "error",
                    "data": json.dumps({"error": str(e)})
                }
        
        return EventSourceResponse(event_generator())

async def handle_messages(request: Request) -> Response:
    """Handle POST requests to /messages endpoint"""
    try:
        data = await request.json()
        return JSONResponse({"jsonrpc": "2.0", "result": data})
    except Exception as e:
        log.error(f"Message handling error: {e}", exc_info=True)
        return JSONResponse(
            {"jsonrpc": "2.0", "error": {"code": -32000, "message": str(e)}},
            status_code=500
        )

async def health_check(request: Request) -> Response:
    """Health check endpoint"""
    return JSONResponse({"status": "healthy", "server": "file_export_mcp"})

app = Starlette(
    debug=True,
    routes=[
        Route("/sse", endpoint=handle_sse, methods=["GET", "POST"]),
        Route("/messages", endpoint=handle_messages, methods=["POST"]),
        Route("/health", endpoint=health_check, methods=["GET"]),
    ]
)

if __name__ == "__main__":

    mode = (os.getenv("MODE", "SSE"))
 
    if mode == "sse":
        port = int(os.getenv("MCP_HTTP_PORT", "9004"))
        host = os.getenv("MCP_HTTP_HOST", "0.0.0.0")
            
        log.info(f"Starting file_export_mcp version {SCRIPT_VERSION}")
        log.info(f"Starting file_export_mcp in SSE mode on http://{host}:{port}")
        log.info(f"SSE endpoint: http://{host}:{port}/sse")
        log.info(f"Messages endpoint: http://{host}:{port}/messages")
            
        uvicorn.run(
            app,
            host=host,
            port=port,
            access_log=False,
            log_level="info",
            use_colors=False
        )
    elif mode == "http":
        port = int(os.getenv("MCP_HTTP_PORT", "9004"))
        host = os.getenv("MCP_HTTP_HOST", "0.0.0.0")
        
        log.info(f"Starting file_export_mcp version {SCRIPT_VERSION}")
        log.info(f"Starting file_export_mcp in http mode on http://{host}:{port}")
        log.info(f"HTTP endpoint: http://{host}:{port}/mcp")

        mcp.run(
            transport="streamable-http"
        )
    else:
        log.info("Starting file_export_mcp in stdio mode version {SCRIPT_VERSION}")
        mcp.run()
